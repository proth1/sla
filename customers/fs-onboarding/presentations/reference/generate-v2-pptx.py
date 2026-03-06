#!/usr/bin/env python3
"""Generate v2-roadmap PowerPoint from HTML content."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ASSETS = "/Users/proth/repos/sla/customers/fs-onboarding/presentations/reference/assets"

# Brand colors
NAVY = RGBColor(0x0D, 0x21, 0x37)
DEEP_BLUE = RGBColor(0x1A, 0x3A, 0x6B)
BLUE = RGBColor(0x25, 0x63, 0xEB)
TEAL = RGBColor(0x00, 0xB4, 0xD8)
GOLD = RGBColor(0xD9, 0x77, 0x06)
EMERALD = RGBColor(0x05, 0x96, 0x69)
ROSE = RGBColor(0xDC, 0x26, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x64, 0x74, 0x8B)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_gradient_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = NAVY
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = DEEP_BLUE
    fill.gradient_stops[1].position = 1.0


def add_footer(slide, dark=False):
    left_color = WHITE if dark else LIGHT_GRAY
    tf = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(6), Inches(0.4)).text_frame
    tf.paragraphs[0].font.size = Pt(8)
    tf.paragraphs[0].font.color.rgb = left_color
    tf.paragraphs[0].text = "Software Onboarding Transformation -- v2"
    tf2 = slide.shapes.add_textbox(Inches(10), Inches(7.0), Inches(2.8), Inches(0.4)).text_frame
    tf2.paragraphs[0].font.size = Pt(8)
    tf2.paragraphs[0].font.color.rgb = left_color
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
    tf2.paragraphs[0].text = "Confidential"


def add_title_text(slide, text, x, y, w, h, size=32, color=WHITE, bold=True, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def add_body_text(slide, text, x, y, w, h, size=14, color=DARK_TEXT, bold=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_after = Pt(6)
    p.line_spacing = Pt(size * 1.5)
    return tf


def add_metric_box(slide, x, y, value, label, value_color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.8), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = value
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = value_color
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(9)
    run2.font.color.rgb = LIGHT_GRAY


def add_metric_box_dark(slide, x, y, value, label, value_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.5), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x44)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = value
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = value_color
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)


def add_card(slide, x, y, w, h, title, body, accent_color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 0.05), Inches(0.06), Inches(h - 0.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    # Title
    tf = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.12), Inches(w - 0.35), Inches(0.35))
    tf.text_frame.word_wrap = True
    p = tf.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = DEEP_BLUE
    # Body
    tf2 = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.45), Inches(w - 0.35), Inches(h - 0.55))
    tf2.text_frame.word_wrap = True
    p2 = tf2.text_frame.paragraphs[0]
    p2.text = body
    p2.font.size = Pt(9)
    p2.font.color.rgb = LIGHT_GRAY
    p2.line_spacing = Pt(13)


def add_pain_point(slide, x, y, w, severity, text):
    color_map = {"Critical": ROSE, "High": GOLD, "Medium": BLUE}
    c = color_map.get(severity, BLUE)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.05), Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = c
    bar.line.fill.background()
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.05), Inches(y), Inches(w - 0.05), Inches(0.35))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xFE, 0xF2, 0xF2) if severity == "Critical" else RGBColor(0xFF, 0xFB, 0xEB) if severity == "High" else RGBColor(0xEF, 0xF6, 0xFF)
    bg.line.fill.background()
    tf = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.02), Inches(w - 0.25), Inches(0.32))
    tf.text_frame.word_wrap = True
    p = tf.text_frame.paragraphs[0]
    run1 = p.add_run()
    run1.text = f"{severity}: "
    run1.font.size = Pt(9)
    run1.font.bold = True
    run1.font.color.rgb = c
    run2 = p.add_run()
    run2.text = text
    run2.font.size = Pt(9)
    run2.font.color.rgb = DARK_TEXT


def add_quote(slide, x, y, w, text, attribution):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.05), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.05), Inches(y), Inches(w - 0.05), Inches(0.55))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    tf = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.04), Inches(w - 0.35), Inches(0.3))
    tf.text_frame.word_wrap = True
    p = tf.text_frame.paragraphs[0]
    p.text = f'"{text}"'
    p.font.size = Pt(9)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x33, 0x41, 0x55)
    tf2 = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.33), Inches(w - 0.35), Inches(0.2))
    p2 = tf2.text_frame.paragraphs[0]
    p2.text = attribution
    p2.font.size = Pt(8)
    p2.font.color.rgb = LIGHT_GRAY


def add_roadmap_box(slide, x, y, w, h, phase_title, subtitle, items, bg_color, border_color, title_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)
    tf = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.1), Inches(w - 0.3), Inches(0.3))
    tf.text_frame.paragraphs[0].text = phase_title
    tf.text_frame.paragraphs[0].font.size = Pt(11)
    tf.text_frame.paragraphs[0].font.bold = True
    tf.text_frame.paragraphs[0].font.color.rgb = title_color
    if subtitle:
        tf2 = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.35), Inches(w - 0.3), Inches(0.2))
        tf2.text_frame.paragraphs[0].text = subtitle
        tf2.text_frame.paragraphs[0].font.size = Pt(8)
        tf2.text_frame.paragraphs[0].font.color.rgb = LIGHT_GRAY
        tf2.text_frame.paragraphs[0].font.bold = True
    items_y = y + (0.55 if subtitle else 0.4)
    tf3 = slide.shapes.add_textbox(Inches(x + 0.15), Inches(items_y), Inches(w - 0.3), Inches(h - items_y + y - 0.1))
    tf3.text_frame.word_wrap = True
    for i, item in enumerate(items):
        p = tf3.text_frame.paragraphs[0] if i == 0 else tf3.text_frame.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(8)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(2)
        p.level = 0


def add_diagram(slide, filename, x, y, max_w, max_h):
    """Add a PNG diagram scaled to fit within max_w x max_h while preserving aspect ratio."""
    from PIL import Image
    path = os.path.join(ASSETS, filename)
    if not os.path.exists(path):
        add_body_text(slide, f"[Diagram: {filename}]", x, y, max_w, max_h, size=10, color=LIGHT_GRAY)
        return
    im = Image.open(path)
    ratio = im.width / im.height
    # Fit within bounds
    w = max_w
    h = w / ratio
    if h > max_h:
        h = max_h
        w = h * ratio
    # Add white container background
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(max_w), Inches(max_h))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.color.rgb = BORDER
    bg.line.width = Pt(1)
    # Center image in container
    img_x = x + (max_w - w) / 2
    img_y = y + (max_h - h) / 2
    slide.shapes.add_picture(path, Inches(img_x), Inches(img_y), Inches(w), Inches(h))


def add_roi_card(slide, x, y, w, h, title, value, desc):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xEC, 0xFD, 0xF5)
    shape.line.color.rgb = EMERALD
    shape.line.width = Pt(2)
    tf = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.12), Inches(w - 0.4), Inches(0.25))
    tf.text_frame.paragraphs[0].text = title
    tf.text_frame.paragraphs[0].font.size = Pt(11)
    tf.text_frame.paragraphs[0].font.bold = True
    tf.text_frame.paragraphs[0].font.color.rgb = EMERALD
    tf2 = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.38), Inches(w - 0.4), Inches(0.35))
    tf2.text_frame.paragraphs[0].text = value
    tf2.text_frame.paragraphs[0].font.size = Pt(24)
    tf2.text_frame.paragraphs[0].font.bold = True
    tf2.text_frame.paragraphs[0].font.color.rgb = EMERALD
    tf3 = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.75), Inches(w - 0.4), Inches(h - 0.85))
    tf3.text_frame.word_wrap = True
    tf3.text_frame.paragraphs[0].text = desc
    tf3.text_frame.paragraphs[0].font.size = Pt(9)
    tf3.text_frame.paragraphs[0].font.color.rgb = LIGHT_GRAY


def add_table(slide, x, y, w, rows, col_widths=None):
    num_rows = len(rows)
    num_cols = len(rows[0])
    table_shape = slide.shapes.add_table(num_rows, num_cols, Inches(x), Inches(y), Inches(w), Inches(0.3 * num_rows))
    table = table_shape.table
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                if i == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = DARK_TEXT
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DEEP_BLUE
            elif i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)
    if col_widths:
        for j, cw in enumerate(col_widths):
            table.columns[j].width = Inches(cw)
    return table


def add_section_slide(title, badge_text, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide)
    if badge_text:
        tf = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(4), Inches(0.35))
        tf.text_frame.paragraphs[0].text = badge_text
        tf.text_frame.paragraphs[0].font.size = Pt(10)
        tf.text_frame.paragraphs[0].font.bold = True
        tf.text_frame.paragraphs[0].font.color.rgb = TEAL
    add_title_text(slide, title, 0.8, 2.6, 10, 0.8, size=36, color=WHITE)
    add_body_text(slide, subtitle, 0.8, 3.5, 10, 0.6, size=14, color=RGBColor(0xCC, 0xCC, 0xDD))
    add_footer(slide, dark=True)
    return slide


def add_topic_triple(topic_num, badge_text, title, subtitle,
                     current_title, pain_points, quotes, extra_content_fn,
                     rec_title, recommendations, roadmap):
    """Create the 3-slide pattern: section, current state, recommendations."""
    # Section slide
    add_section_slide(f"{topic_num}. {title}", badge_text, subtitle)

    # Current state slide
    cs = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(cs, RGBColor(0xF8, 0xFA, 0xFC))
    add_title_text(cs, current_title, 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
    y_pos = 1.0
    for sev, txt in pain_points:
        add_pain_point(cs, 0.6, y_pos, 5.8, sev, txt)
        y_pos += 0.4
    q_y = y_pos + 0.1
    for qtxt, qattr in quotes:
        add_quote(cs, 0.6, q_y, 5.8, qtxt, qattr)
        q_y += 0.65
    if extra_content_fn:
        extra_content_fn(cs)
    add_footer(cs)

    # Recommendations slide
    rec = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(rec, RGBColor(0xF8, 0xFA, 0xFC))
    add_title_text(rec, rec_title, 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
    col = 0
    row = 0
    for i, (rtitle, rbody, rcolor) in enumerate(recommendations):
        cx = 0.6 + col * 6.15
        cy = 1.0 + row * 1.65
        add_card(rec, cx, cy, 5.9, 1.5, rtitle, rbody, rcolor)
        col += 1
        if col >= 2:
            col = 0
            row += 1
    # Roadmap row
    rm_y = 4.5
    rm_w = 2.9
    for i, (rm_title, rm_sub, rm_items, rm_bg, rm_border, rm_tc) in enumerate(roadmap):
        add_roadmap_box(rec, 0.6 + i * 3.05, rm_y, rm_w, 2.3, rm_title, rm_sub, rm_items, rm_bg, rm_border, rm_tc)
    add_footer(rec)


# ============================================================
# SLIDE 0: Title
# ============================================================
slide0 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide0)
add_title_text(slide0, "Software Onboarding Transformation", 0.8, 1.5, 11, 0.8, size=40, color=WHITE)
add_title_text(slide0, "Topic Deep Dive & Implementation Roadmap", 0.8, 2.3, 11, 0.6, size=22, color=WHITE, bold=True)
add_body_text(slide0, "11 governance topics. Current state assessment, detailed recommendations, and 30/60/90/120-day roadmaps for each. Grounded in 9 stakeholder sessions, 30+ interviews, and 16 identified process gaps.", 0.8, 3.1, 10, 0.8, size=13, color=RGBColor(0xCC, 0xCC, 0xDD))
add_metric_box_dark(slide0, 0.8, 4.3, "6-9 mo", "Current Cycle", RGBColor(0xFC, 0xA5, 0xA5))
add_metric_box_dark(slide0, 3.6, 4.3, "19 days", "Target Cycle", RGBColor(0x6E, 0xE7, 0xB7))
add_metric_box_dark(slide0, 6.4, 4.3, "11", "Topic Areas", WHITE)
add_metric_box_dark(slide0, 9.2, 4.3, "16", "Gaps Identified", WHITE)
add_footer(slide0, dark=True)

# ============================================================
# SLIDE 1: Executive Summary
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide1, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(slide1, "Executive Summary", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_body_text(slide1, "Nine discovery sessions with 30+ stakeholders revealed an end-to-end process taking 6-9 months, routing through 18 sequential committees, and accepting requests from 5+ disconnected channels.", 0.6, 0.9, 11.5, 0.5, size=12, color=DARK_TEXT)
add_roi_card(slide1, 0.6, 1.6, 3.7, 1.3, "Intake Consolidation", "70% faster", "Single gateway replaces 5+ channels. Completeness gate eliminates 40%+ rework.")
add_roi_card(slide1, 4.6, 1.6, 3.7, 1.3, "Parallel Evaluation", "18 -> 5", "Sequential committees replaced with 5 concurrent streams.")
add_roi_card(slide1, 8.6, 1.6, 3.7, 1.3, "3-Pathway Routing", "~30 day", "Vendor Affinity pathway for low-risk advisor tools.")
# Metrics row
metrics = [("6-9 mo", "End-to-End Duration"), ("18", "Sequential Committees"), ("5+", "Intake Channels"),
           ("60+", "AI Queue Items"), ("2", "Legal Partners"), ("~80", "REA Questions")]
for i, (val, lbl) in enumerate(metrics):
    add_metric_box(slide1, 0.6 + i * 2.05, 3.2, val, lbl, ROSE)
add_footer(slide1)

# ============================================================
# SLIDE 2: E2E Workflow Pain Points
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(slide2, "End-to-End Workflow -- Pain Point Overlay", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_body_text(slide2, "The 5-phase lifecycle with red callouts marking the biggest pain points at each stage, validated by stakeholder interviews.", 0.6, 0.9, 11, 0.4, size=12, color=DARK_TEXT)
add_diagram(slide2, "top-level.png", 0.6, 1.4, 12.0, 3.2)
# Phase pipeline
phases = [("SP1: Intake", "5+ channels, ~80 Q form, no pre-screen"),
          ("SP2: Planning", "No prioritization, binary Buy/Build only"),
          ("SP3: Eval & DD", "18 sequential committees, no security baseline"),
          ("SP4: Contracting", "2-3 mo contracting, 2 legal partners, no rework loop"),
          ("SP5: Go-Live", "No conditional approval, no ownership tracking")]
for i, (phase, desc) in enumerate(phases):
    x = 0.6 + i * 2.45
    add_card(slide2, x, 4.8, 2.3, 1.4, phase, desc, ROSE)
# Pipeline visual
for i, (phase, _) in enumerate(phases):
    x = 0.6 + i * 2.45
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(3.3), Inches(2.3), Inches(0.5)) if False else None
add_footer(slide2)

# ============================================================
# SLIDE 3: High-Level Roadmap
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(slide3, "Implementation Roadmap -- 30/60/90/120 Days", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_body_text(slide3, "Phased approach: planning and consolidation first, then progressive automation aligned to stakeholder urgency and gap priorities.", 0.6, 0.85, 11, 0.4, size=12, color=DARK_TEXT)
roadmap_data = [
    ("Days 1-30: Plan & Consolidate", "Zero automation -- define standards",
     ["Map and deduplicate intake fields across 5+ channels", "Define prioritization scoring formula",
      "Define security baseline controls (3-tier)", "Map 18 committees to 5 parallel streams",
      "Define Vendor Affinity pathway rules", "Establish deal-killer no-go list", "Executive alignment and sign-off"],
     RGBColor(0xEF, 0xF6, 0xFF), BLUE, BLUE),
    ("Days 31-60: Quick Wins", "First measurable cycle time reduction",
     ["Deploy unified intake form (GAP-1)", "Implement completeness quality gate (GAP-9)",
      "Deploy deal-killer pre-screen (GAP-16)", "Add status notifications at phase boundaries (GAP-7)",
      "Implement prioritization scoring (GAP-2)", "Add mandatory ownership assignment (GAP-14)"],
     RGBColor(0xF0, 0xFD, 0xFA), TEAL, RGBColor(0x08, 0x91, 0xB2)),
    ("Days 61-90: Structural Changes", "Pathway routing and parallel evaluation",
     ["Deploy 3-pathway routing (GAP-11)", "Implement progressive forms (GAP-3)",
      "Deploy tiered security assessment (GAP-12)", "Restructure SP3: 5 parallel streams",
      "Add finance rework loop (GAP-4)", "Implement conditional approval (GAP-13)"],
     RGBColor(0xFF, 0xFB, 0xEB), GOLD, GOLD),
    ("Days 91-120: Optimize", "Advanced automation and measurement",
     ["Deploy workload visibility dashboard (GAP-10)", "Implement AI fast-track pathway (GAP-6)",
      "Deploy exception routing (GAP-8)", "Implement pre-onboarding idea funnel (GAP-15)",
      "Measure before/after metrics", "Executive reporting on cycle time reduction"],
     RGBColor(0xF0, 0xFD, 0xF4), EMERALD, EMERALD),
]
for i, (title, sub, items, bg_c, border_c, title_c) in enumerate(roadmap_data):
    add_roadmap_box(slide3, 0.5 + i * 3.1, 1.4, 2.95, 5.3, title, sub, items, bg_c, border_c, title_c)
add_footer(slide3)

# ============================================================
# SLIDE 4: Day 30 Detail
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(slide4, "Day 30 -- Process Improvements & Consolidations", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_body_text(slide4, "The first 30 days focus entirely on planning artifacts. No automation deployed -- but these decisions unlock everything that follows.", 0.6, 0.85, 11, 0.4, size=12, color=DARK_TEXT)
day30_cards = [
    ("Intake Field Inventory", "Map every question across ServiceNow, REA (~80 Q), AI forms, Power Apps, email. Deduplicate. Define minimum viable fields for unified intake.", BLUE),
    ("Prioritization Formula", "Define quantitative scoring: business impact, strategic alignment, urgency, risk tier, capacity. Draft OB-DMN-5 decision table.", BLUE),
    ("Security Baseline Definition", "Define 3-tier control hierarchy: Baseline (automated), Elevated (automated + manual), Major (full assessment + pen test).", BLUE),
    ("Committee Consolidation Plan", "Map all 18 committees to 5 parallel evaluation streams. Identify which reviews can run simultaneously. Draft communication plan.", GOLD),
    ("Vendor Affinity Rules", "Define 3rd pathway criteria: no org cost, advisor-direct purchase, vendor partnership program. Update OB-DMN-2 routing table.", GOLD),
    ("Executive Alignment", "Present discovery findings and roadmap to executive sponsors. Confirm committee consolidation. Assign owners to every Day 60 deliverable.", GOLD),
]
for i, (t, b, c) in enumerate(day30_cards):
    col = i % 3
    row = i // 3
    add_card(slide4, 0.6 + col * 4.05, 1.5 + row * 2.5, 3.85, 2.3, t, b, c)
add_footer(slide4)

# ============================================================
# TOPICS 1-11 (3 slides each)
# ============================================================

# TOPIC 1: INTAKE
add_section_slide("1. Intake", "SP1: Request & Triage", "Initial request capture, portfolio review, and solution discovery. The front door to the entire onboarding lifecycle.")

s6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s6, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(s6, "Intake -- Current State", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_diagram(s6, "sp1-request-triage.png", 0.6, 0.9, 12.0, 2.5)
add_table(s6, 0.6, 3.6, 5.8, [
    ["Channel", "System", "Issue"],
    ["TBC Intake", "ServiceNow", "Inconsistent compliance"],
    ["REA Form", "~80 questions", "Frequently completed incorrectly"],
    ["AI Use Case", "Separate form", "2 forms at different stages"],
    ["Rapid Risk", "Power Apps", "Used as onboarding workaround"],
    ["Email/Chat", "Unstructured", "No tracking"],
])
add_pain_point(s6, 6.8, 3.6, 5.6, "Critical", "Business partners submit same information to multiple teams")
add_pain_point(s6, 6.8, 4.05, 5.6, "Critical", "Multiple groups run independent programs, seeking fastest path")
add_pain_point(s6, 6.8, 4.5, 5.6, "High", "REA form frequently filled out incorrectly -- even experienced users struggle")
add_quote(s6, 6.8, 5.1, 5.6, "Multiple business cases submitted on a single form.", "-- Legal Director")
add_footer(s6)

s7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s7, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(s7, "Intake -- Recommendations", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_card(s7, 0.6, 1.0, 5.9, 1.4, "Unified Intake Gateway (GAP-1)", "Single entry point absorbs all 5+ channels. Dynamic routing by request type: Standard, AI, Exception, Vendor Partnership. Eliminates duplicate submissions.", EMERALD)
add_card(s7, 6.7, 1.0, 5.9, 1.4, "Completeness Quality Gate (GAP-9)", "Automated screening before human review. AI-assisted pre-screening validates minimum viable fields, problem statement quality, and risk flags.", EMERALD)
add_card(s7, 0.6, 2.55, 5.9, 1.4, "Deal-Killer Pre-Screen (GAP-16)", "DMN-driven check against managed no-go list. Certain AI models are complete no-gos -- reject immediately with explanation rather than consuming weeks.", GOLD)
add_card(s7, 6.7, 2.55, 5.9, 1.4, "Existing Portfolio Check", "Before formal intake, check if existing solution satisfies the need. Route to 'Leverage Existing' -- avoid the entire onboarding process for duplicates.", BLUE)
rm_colors = [
    (RGBColor(0xEF, 0xF6, 0xFF), BLUE, BLUE),
    (RGBColor(0xF0, 0xFD, 0xFA), TEAL, RGBColor(0x08, 0x91, 0xB2)),
    (RGBColor(0xFF, 0xFB, 0xEB), GOLD, GOLD),
    (RGBColor(0xF0, 0xFD, 0xF4), EMERALD, EMERALD),
]
intake_rm = [
    ("Day 30", None, ["Map all intake fields", "Deduplicate questions", "Define minimum viable fields", "Draft deal-killer no-go list"]),
    ("Day 60", None, ["Deploy unified intake form", "Activate completeness gate", "Deploy deal-killer pre-screen", "Redirect all channels"]),
    ("Day 90", None, ["Progressive form strategy", "AI-assisted pre-screening", "Form analytics: completion rates"]),
    ("Day 120", None, ["Exception routing live", "Pre-onboarding idea funnel", "Measure: form completion %, rejection rate"]),
]
for i, ((t, sub, items), (bg, brd, tc)) in enumerate(zip(intake_rm, rm_colors)):
    add_roadmap_box(s7, 0.6 + i * 3.1, 4.2, 2.9, 2.7, t, sub, items, bg, brd, tc)
add_footer(s7)

# TOPIC 2: PRIORITIZATION
add_section_slide("2. Prioritization", "SP2: Planning & Routing", "Backlog ranking, pathway selection, and resource allocation. The missing link between intake and evaluation.")

s9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s9, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(s9, "Prioritization -- Current State", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_diagram(s9, "sp2-planning-routing.png", 0.6, 0.9, 12.0, 1.8)
add_pain_point(s9, 0.6, 2.9, 5.8, "Critical", "No formal prioritization standard or formula")
add_pain_point(s9, 0.6, 3.35, 5.8, "Critical", "No force-ranking mechanism across the enterprise")
add_pain_point(s9, 0.6, 3.8, 5.8, "High", "Each requestor views their request as most important")
add_pain_point(s9, 0.6, 4.25, 5.8, "High", "Auto-approval for low-risk items not supported -- all follow same heavyweight process")
add_quote(s9, 6.8, 2.9, 5.6, "Teams horse trade internally before involving Legal.", "-- Process Manager")
add_quote(s9, 6.8, 3.6, 5.6, "Competing initiatives across the organization with no arbitration.", "-- Procurement Lead")
add_quote(s9, 6.8, 4.3, 5.6, "Prioritization can be partially driven by scoring intake form responses.", "-- Governance Specialist")
add_footer(s9)

s10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s10, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(s10, "Prioritization -- Recommendations", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_card(s10, 0.6, 1.0, 5.9, 1.5, "DMN-Driven Scoring (GAP-2)", "New OB-DMN-5: inputs = business impact, strategic alignment, urgency, risk tier, capacity. Output: P1/P2/P3 priority score. Replaces 'horse trading' with objective ranking.", EMERALD)
add_card(s10, 6.7, 1.0, 5.9, 1.5, "3-Pathway Routing (GAP-11)", "Replace binary Buy/Build gateway with tri-state routing: Buy, Build, Vendor Affinity. OB-DMN-2 expanded. Right-sized governance per pathway.", TEAL)
add_card(s10, 0.6, 2.65, 5.9, 1.5, "Low-Complexity Bypass", "SP2 gateway: requests that don't need full evaluation skip directly to pathway routing. Auto-approval for low-risk, low-complexity items.", BLUE)
add_card(s10, 6.7, 2.65, 5.9, 1.5, "Transparent Queue Position", "Requestors see their position in the ranked backlog. Status notifications at every transition. Data-driven capacity decisions.", GOLD)
prio_rm = [
    ("Day 30", None, ["Define scoring formula", "Draft OB-DMN-5 table", "Get governance sign-off", "Define VA criteria"]),
    ("Day 60", None, ["Deploy scoring in SP2", "Publish queue position", "Status notifications"]),
    ("Day 90", None, ["3-pathway routing live", "OB-DMN-2 updated with VA", "Low-complexity bypass active"]),
    ("Day 120", None, ["Measure: queue wait time", "Refine scoring weights", "Executive reporting"]),
]
for i, ((t, sub, items), (bg, brd, tc)) in enumerate(zip(prio_rm, rm_colors)):
    add_roadmap_box(s10, 0.6 + i * 3.1, 4.4, 2.9, 2.5, t, sub, items, bg, brd, tc)
add_footer(s10)

# TOPIC 3: FUNDING
add_section_slide("3. Funding / Finance", "SP3-SP4: Evaluation & Contracting", "Financial analysis, budget authorization, and TCO/ROI assessment. Establishes fiscal justification before commitments.")

s12 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s12, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(s12, "Funding / Finance -- Current State", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_pain_point(s12, 0.6, 1.0, 5.8, "Critical", "Formal denial sends request back to beginning -- even for minor corrections")
add_pain_point(s12, 0.6, 1.45, 5.8, "High", "Cannot reroute coding matrix issues to FP&A -- must email outside system")
add_pain_point(s12, 0.6, 1.9, 5.8, "High", "Cannot make minor cosmetic corrections independently")
add_pain_point(s12, 0.6, 2.35, 5.8, "Medium", "Vendor partnership products require no funding but validation still required")
add_quote(s12, 0.6, 2.9, 5.8, "Formal denial sends request back to beginning of process.", "-- Finance Representative")
add_table(s12, 6.8, 1.0, 5.5, [
    ["Role", "RACI"],
    ["Finance", "R/A -- performs TCO, ROI, NPV"],
    ["Business", "C -- provides business case and financial inputs"],
    ["Governance", "C -- authorizes budget"],
    ["Procurement", "C -- procurement routing"],
    ["Oversight", "I -- informed of financial decisions"],
])
add_footer(s12)

s13 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s13, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(s13, "Funding / Finance -- Recommendations", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_card(s13, 0.6, 1.0, 5.9, 1.5, "Finance Rework Loop (GAP-4)", "Add correction loop in SP4: minor issues loop back to correction task without formal denial. Only substantive issues restart the process.", EMERALD)
add_card(s13, 6.7, 1.0, 5.9, 1.5, "Vendor Affinity Bypass", "VA products (no org cost) skip funding validation entirely. OB-DMN-2 VA output bypasses Task_FinancialAnalysis in SP3 and negotiation in SP4.", TEAL)
add_card(s13, 0.6, 2.65, 5.9, 1.5, "FP&A Integration", "Direct routing to FP&A for coding matrix issues within the process -- no external email required. Automated routing based on correction type.", GOLD)
add_card(s13, 6.7, 2.65, 5.9, 1.5, "Progressive Financial Data", "SP1: budget estimate only. SP2: strategic alignment. SP3: full TCO/ROI model. No financial question appears in more than one form.", BLUE)
fund_rm = [
    ("Day 30", None, ["Map financial questions", "Define correction vs. denial", "Identify FP&A routing rules"]),
    ("Day 60", None, ["Progressive fields in form", "FP&A routing path defined"]),
    ("Day 90", None, ["Finance rework loop deployed", "VA bypass active", "Automated FP&A routing"]),
    ("Day 120", None, ["Measure: restart rate", "Average correction cycle time", "VA bypass metrics"]),
]
for i, ((t, sub, items), (bg, brd, tc)) in enumerate(zip(fund_rm, rm_colors)):
    add_roadmap_box(s13, 0.6 + i * 3.1, 4.4, 2.9, 2.5, t, sub, items, bg, brd, tc)
add_footer(s13)

# TOPIC 4: SOURCING
add_section_slide("4. Sourcing", "SP3: Evaluation & Due Diligence", "Vendor landscape assessment, vendor selection, and risk evaluation. Includes procurement execution and due diligence.")

s15 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s15, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(s15, "Sourcing -- Current State", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_diagram(s15, "sp3-evaluation-dd.png", 0.6, 0.9, 12.0, 3.0)
add_pain_point(s15, 0.6, 4.1, 5.8, "Critical", "Business partners bypass sourcing, arriving with pre-selected vendors")
add_pain_point(s15, 0.6, 4.55, 5.8, "High", "Sourcing team makes quasi-legal decisions due to legal bottleneck")
add_pain_point(s15, 0.6, 5.0, 5.8, "High", "Solutions underutilized -- multi-million dollar tools used for 1/3 of capabilities")
add_pain_point(s15, 0.6, 5.45, 5.8, "Critical", "18 sequential committee reviews -- meetings 5+ weeks apart")
add_quote(s15, 6.8, 4.1, 5.6, "Same use case presented to multiple committees sequentially.", "-- Business Representative")
add_quote(s15, 6.8, 4.8, 5.6, "Architecture teams assessing financials. Financial teams assessing architecture.", "-- Technical Lead")
add_footer(s15)

s16 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s16, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(s16, "Sourcing -- Recommendations", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_card(s16, 0.6, 1.0, 5.9, 1.5, "5 Parallel Evaluation Streams", "Replace 18 sequential committees with 5 concurrent streams: Tech Architecture, Security, Risk/Compliance/Legal, Financial/TCO, Vendor Landscape.", EMERALD)
add_card(s16, 6.7, 1.0, 5.9, 1.5, "Vendor Affinity Reduced Scope (GAP-11)", "VA pathway: skip financial analysis and vendor landscape assessment. Compliance review and security assessment only -- right-sized due diligence.", TEAL)
add_card(s16, 0.6, 2.65, 5.9, 1.5, "Vendor DD via Message Flow", "Structured cross-pool communication: DD request sent to vendor pool, vendor responds with evidence. 5-day SLA timer on response.", GOLD)
add_card(s16, 6.7, 2.65, 5.9, 1.5, "Post-Purchase Utilization Tracking", "Post-onboarding includes utilization review tasks. Prevents multi-million dollar tools being used for 1/3 of capabilities. Annual review.", BLUE)
src_rm = [
    ("Day 30", None, ["Map 18 committees to 5 streams", "Define concurrent review rules", "Get committee chair buy-in"]),
    ("Day 60", None, ["Pilot parallel evaluation", "Vendor DD template standardized"]),
    ("Day 90", None, ["5 parallel streams operational", "VA reduced scope deployed", "5-day SLA timers active"]),
    ("Day 120", None, ["Measure: eval cycle time", "Post-purchase utilization", "Vendor response SLA"]),
]
for i, ((t, sub, items), (bg, brd, tc)) in enumerate(zip(src_rm, rm_colors)):
    add_roadmap_box(s16, 0.6 + i * 3.1, 4.4, 2.9, 2.5, t, sub, items, bg, brd, tc)
add_footer(s16)

# TOPIC 5: CYBERSECURITY
add_section_slide("5. Cybersecurity", "SP3: Evaluation & Due Diligence", 'Security assessment, vulnerability management, and security architecture review. The "secure by design" gap.')

s18 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s18, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(s18, "Cybersecurity -- Current State", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_pain_point(s18, 0.6, 1.0, 5.8, "Critical", 'No defined "secure by design" standard at the organization')
add_pain_point(s18, 0.6, 1.45, 5.8, "Critical", "Teams don't know minimum security controls; enforcement 'fairly loose'")
add_pain_point(s18, 0.6, 1.9, 5.8, "High", "Only a fraction of systems covered by identity management")
add_pain_point(s18, 0.6, 2.35, 5.8, "High", "End-of-year compliance failures common; PII in uncontrolled systems")
add_pain_point(s18, 0.6, 2.8, 5.8, "High", "Security tools also contain AI -- circular dependency for approvals")
add_quote(s18, 6.8, 1.0, 5.6, "AI technologies still in formation; no public standards for securing new tech.", "-- AI Security Lead")
add_quote(s18, 6.8, 1.7, 5.6, "MCP technology didn't exist 18 months ago -- new tech requires new baselines.", "-- AI Security Lead")
add_footer(s18)

s19 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s19, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(s19, "Cybersecurity -- Recommendations", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_card(s19, 0.6, 1.0, 3.85, 2.8, "Tiered Security Baseline (GAP-12)", "DMN-driven assessment level based on risk tier + data classification + AI component:\n\nBaseline: Automated checks only\nElevated: Automated + manual review\nMajor: Full assessment + pen test", EMERALD)
add_card(s19, 4.7, 1.0, 3.85, 2.8, "Time-Bound Risk Acceptance (GAP-13)", '"Company wants product by Q2 but security controls not ready until Q3." Conditional approval with monitoring timer. Auto-escalation if conditions not met by expiry.', TEAL)
add_card(s19, 8.8, 1.0, 3.85, 2.8, "Automated Baseline Checks", "Baseline tier runs as service task in Automation lane -- no human review needed. Only Elevated/Major route to manual review. Reduces security team workload.", GOLD)
cyber_rm = [
    ("Day 30", None, ["Define 3-tier control hierarchy", "Map requirements per tier", "Security Lead sign-off"]),
    ("Day 60", None, ["Document baseline checklist", "Draft DMN table for routing"]),
    ("Day 90", None, ["Tiered assessment deployed", "Automated baseline active", "Conditional approval in SP5"]),
    ("Day 120", None, ["Measure: review cycle time", "Baseline automation rate", "Conditional monitoring"]),
]
for i, ((t, sub, items), (bg, brd, tc)) in enumerate(zip(cyber_rm, rm_colors)):
    add_roadmap_box(s19, 0.6 + i * 3.1, 4.2, 2.9, 2.7, t, sub, items, bg, brd, tc)
add_footer(s19)

# TOPIC 6: EA
add_section_slide("6. Enterprise Architecture (Buy Only)", "SP3-SP4: Evaluation & Build", "Architecture review, integration planning, and technology standards compliance. Ensures software fits the enterprise landscape.")

s21 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s21, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(s21, "Enterprise Architecture -- Current State", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_diagram(s21, "sp4-contracting-build.png", 0.6, 0.9, 12.0, 2.0)
add_pain_point(s21, 0.6, 3.1, 5.8, "High", "Architecture review group recently reduced in staffing")
add_pain_point(s21, 0.6, 3.55, 5.8, "High", "Architecture teams assessing financials (outside their expertise)")
add_pain_point(s21, 0.6, 4.0, 5.8, "High", "Integration costs beyond initial purchase not captured in Buy/Build binary")
add_quote(s21, 0.6, 4.6, 5.8, "Some 'buy' scenarios also require 'build' -- AI tool generating emails needs compliance capture integration.", "-- Product Director")
add_table(s21, 6.8, 3.1, 5.5, [
    ["Task", "Sub-Process"],
    ["Technical Architecture Review", "SP3"],
    ["Perform Proof of Concept", "SP4"],
    ["Technical Risk Evaluation", "SP4"],
    ["Vendor Technical Demo", "Vendor Pool"],
    ["Vendor Deployment Support", "Vendor Pool"],
])
add_footer(s21)

s22 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s22, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(s22, "Enterprise Architecture -- Recommendations", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_card(s22, 0.6, 1.0, 5.9, 1.5, "Integration Governance for Buy Path", "When Buy pathway identifies integration work, trigger SP4 Task_DefineBuildReqs. Hybrid buy+build pattern captured explicitly.", EMERALD)
add_card(s22, 6.7, 1.0, 5.9, 1.5, "PDLC for Build Path", "Build pathway enters nested PDLC: Architecture Review, Development, Testing & QA, Integration Testing. Retry loop on test failures.", TEAL)
add_diagram(s22, "pdlc.png", 6.7, 2.65, 5.9, 1.0)
add_card(s22, 0.6, 2.65, 5.9, 1.5, "Focused Review Scope", "Architecture teams review architecture only -- no financial assessment. Financial analysis runs in parallel, not through architecture review.", GOLD)
add_card(s22, 6.7, 3.8, 5.9, 1.0, "Standards Compliance Checklist", "Automated technology standards check before PoC. DMN-driven: integration patterns, data flow, security architecture fit.", BLUE)
ea_rm = [
    ("Day 30", None, ["Define Buy path checklist", "Document integration triggers", "Clarify review scope"]),
    ("Day 60", None, ["Architecture decoupled from financial", "Standards checklist documented"]),
    ("Day 90", None, ["Integration governance active", "PDLC for Build path", "Parallel review enforced"]),
    ("Day 120", None, ["Measure: review cycle time", "Integration cost capture rate", "PoC success rate"]),
]
for i, ((t, sub, items), (bg, brd, tc)) in enumerate(zip(ea_rm, rm_colors)):
    add_roadmap_box(s22, 0.6 + i * 3.1, 4.4, 2.9, 2.5, t, sub, items, bg, brd, tc)
add_footer(s22)

# TOPIC 7: COMPLIANCE
add_section_slide("7. Compliance", "Cross-Cutting: SP1-SP5", "Regulatory compliance review, legal assessment, and data protection obligations across the full lifecycle.")

s24 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s24, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(s24, "Compliance -- Current State", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_table(s24, 0.6, 1.0, 5.5, [
    ["Regulation", "Scope"],
    ["OCC 2023-17", "Third-party risk management"],
    ["GDPR / CCPA", "Data processing and privacy"],
    ["SOX 302/404", "Internal controls, financial reporting"],
    ["DORA", "Digital operational resilience"],
    ["EU AI Act", "AI system risk classification"],
    ["NIST CSF 2.0", "Cybersecurity framework"],
])
add_body_text(s24, "Compliance review appears in all 5 phases plus the Vendor Pool -- highest breadth of any governance topic.", 0.6, 3.4, 5.5, 0.4, size=10, color=DARK_TEXT)
add_pain_point(s24, 6.8, 1.0, 5.6, "Critical", "No consistent standard for what requires escalation")
add_pain_point(s24, 6.8, 1.45, 5.6, "High", "End-of-year compliance failures common; PII in uncontrolled systems")
add_pain_point(s24, 6.8, 1.9, 5.6, "High", "Different teams have varying risk acceptance thresholds")
add_quote(s24, 6.8, 2.5, 5.6, "Information trickles through various processes instead of arriving complete upfront.", "-- Legal Director")
add_footer(s24)

s25 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s25, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(s25, "Compliance -- Recommendations", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_card(s25, 0.6, 1.0, 3.85, 2.8, "Phase-Boundary Quality Gates", "Every phase transition validates compliance prerequisites before allowing the request to advance. Prevents incomplete information from cascading.", EMERALD)
add_card(s25, 4.7, 1.0, 3.85, 2.8, "DMN-Driven Escalation", "Risk tier determines escalation: Minimal = single approver, Limited = committee quorum, High = full advisory board. OB-DMN-3 routes automatically.", TEAL)
add_card(s25, 8.8, 1.0, 3.85, 2.8, "Regulatory Traceability", "Every task carries regulatory annotations via Camunda properties. Audit trail links decisions to OCC, GDPR, SOX, DORA requirements.", GOLD)
comp_rm = [
    ("Day 30", None, ["Map regulations to phases", "Define escalation criteria", "Document quality gates"]),
    ("Day 60", None, ["Phase-boundary gates active", "Compliance notification"]),
    ("Day 90", None, ["OB-DMN-3 routing deployed", "3-tier review operational", "Regulatory annotations"]),
    ("Day 120", None, ["Compliance audit trail", "Measure: failure rate", "Traceability dashboard"]),
]
for i, ((t, sub, items), (bg, brd, tc)) in enumerate(zip(comp_rm, rm_colors)):
    add_roadmap_box(s25, 0.6 + i * 3.1, 4.2, 2.9, 2.7, t, sub, items, bg, brd, tc)
add_footer(s25)

# TOPIC 8: AI GOVERNANCE
add_section_slide("8. AI Governance", "SP3: Evaluation & Due Diligence", "AI and ML risk classification, model validation, and bias assessment. The highest-urgency domain with 60+ items in queue.")

s27 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s27, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(s27, "AI Governance -- Current State", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_pain_point(s27, 0.6, 1.0, 5.8, "Critical", "Overly restrictive AI addendum causes extended vendor negotiations")
add_pain_point(s27, 0.6, 1.45, 5.8, "Critical", "AI tools perceived as taking too long -- stakeholders want AI-first solutions")
add_pain_point(s27, 0.6, 1.9, 5.8, "High", "60+ items in AI governance queue")
add_pain_point(s27, 0.6, 2.35, 5.8, "High", "AI process involves 3 separate committees")
add_pain_point(s27, 0.6, 2.8, 5.8, "High", "Multiple tools submitted for same function -- no AI strategy alignment")
add_quote(s27, 6.8, 1.0, 5.6, "AI should be treated as the first use case with planned expansion afterward.", "-- Product Director")
add_quote(s27, 6.8, 1.7, 5.6, "'Crawl, walk, run' messaging rejected -- clear target state needed.", "-- Product Director")
add_quote(s27, 6.8, 2.4, 5.6, "Need to communicate non-starter models/vendors early to prevent wasted effort.", "-- Business Representative")
add_footer(s27)

s28 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s28, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(s28, "AI Governance -- Recommendations", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_card(s28, 0.6, 1.0, 5.9, 1.5, "AI Fast-Track Pathway (GAP-6)", "AI tools run AI Gov Review + Security Assessment only. Pre-defined AI risk posture reduces decision latency. Target: 2-week end-to-end vs. 6-9 months.", EMERALD)
add_card(s28, 6.7, 1.0, 5.9, 1.5, "Deal-Killer Pre-Screen for AI (GAP-16)", "Known no-go AI models/vendors rejected at intake. Enterprise Risk Management red/green light decision matrix in SP1. Prevents wasted review cycles.", TEAL)
add_card(s28, 0.6, 2.65, 5.9, 1.5, "Streamlined AI Addendum", "Working with outside firm to redefine AI stipulations. Predefined scenarios get streamlined approval; complex cases get full review. Reduces vendor pushback.", GOLD)
add_card(s28, 6.7, 2.65, 5.9, 1.5, "3-Committee Consolidation", "AI Risk Working Group, AI Cyber Review, AI Risk Review, AI Governance Committee consolidated into single AI Governance Review task in SP3.", BLUE)
ai_rm = [
    ("Day 30", None, ["Define AI risk posture", "Build AI no-go list", "Define fast-track criteria", "Map 4 AI committees to 1"]),
    ("Day 60", None, ["AI deal-killer active", "Committee consolidation communicated"]),
    ("Day 90", None, ["AI fast-track piloted", "Streamlined addendum drafted"]),
    ("Day 120", None, ["AI fast-track operational", "Measure: queue depth", "Target: 2 weeks"]),
]
for i, ((t, sub, items), (bg, brd, tc)) in enumerate(zip(ai_rm, rm_colors)):
    add_roadmap_box(s28, 0.6 + i * 3.1, 4.4, 2.9, 2.5, t, sub, items, bg, brd, tc)
add_footer(s28)

# TOPIC 9: PRIVACY
add_section_slide("9. Privacy (Legal)", "Cross-Cutting: SP1-SP4", "Data protection impact assessment, data classification, and cross-border transfer compliance. GDPR/CCPA obligations throughout the lifecycle.")

s30 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s30, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(s30, "Privacy -- Current State", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_pain_point(s30, 0.6, 1.0, 5.8, "High", "Different teams have varying risk acceptance thresholds for data handling")
add_pain_point(s30, 0.6, 1.45, 5.8, "High", "PII found in uncontrolled systems -- end-of-year compliance failures")
add_pain_point(s30, 0.6, 1.9, 5.8, "High", "AI, offshoring, and subcontracting requirements create rework when discovered late")
add_pain_point(s30, 0.6, 2.35, 5.8, "Medium", "AI-related requests add complexity with evolving privacy guidance")
add_table(s30, 6.8, 1.0, 5.5, [
    ["Task", "Phase", "Privacy Relevance"],
    ["Gather Requirements Docs", "SP1", "Data classification"],
    ["Preliminary Analysis", "SP2", "Data flow identification"],
    ["Risk, Compliance & Legal", "SP3", "DPIA, cross-border review"],
    ["Refine Requirements", "SP4", "DPA drafting"],
    ["Perform PoC", "SP4", "Data handling validation"],
])
add_footer(s30)

s31 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s31, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(s31, "Privacy -- Recommendations", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_card(s31, 0.6, 1.0, 3.85, 2.8, "Early Data Classification", "SP1 intake captures data classification at submission: PII, PHI, financial data, public. Classification drives automated routing of privacy review requirements.", EMERALD)
add_card(s31, 4.7, 1.0, 3.85, 2.8, "Progressive Privacy Disclosure", "SP1: basic data type. SP2: data flow direction. SP3: full DPIA when warranted. SP4: DPA terms. No privacy question repeated across phases.", TEAL)
add_card(s31, 8.8, 1.0, 3.85, 2.8, "AI Privacy Triggers", "When AI component detected: automatic cross-border data assessment, training data provenance check, and model card review. Integrated into AI Gov branch.", GOLD)
priv_rm = [
    ("Day 30", None, ["Define data classification", "Map privacy questions per phase", "Document DPIA triggers"]),
    ("Day 60", None, ["Data classification in intake", "Automated privacy routing"]),
    ("Day 90", None, ["Progressive privacy fields", "AI privacy triggers active", "DPIA automation"]),
    ("Day 120", None, ["Measure: rework rate", "DPIA completion time", "Cross-border compliance"]),
]
for i, ((t, sub, items), (bg, brd, tc)) in enumerate(zip(priv_rm, rm_colors)):
    add_roadmap_box(s31, 0.6 + i * 3.1, 4.2, 2.9, 2.7, t, sub, items, bg, brd, tc)
add_footer(s31)

# TOPIC 10: COMMERCIAL COUNSEL
add_section_slide("10. Commercial Counsel", "SP4: Contracting & Build", "Contract negotiation, SOW/MSA drafting, redlining, and execution. The 2-3 month bottleneck with only 2 legal partners.")

s33 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s33, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(s33, "Commercial Counsel -- Current State", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_diagram(s33, "sp4-contracting-build.png", 0.6, 0.9, 12.0, 2.0)
add_pain_point(s33, 0.6, 3.1, 5.8, "Critical", "Contracting alone takes 2-3 months")
add_pain_point(s33, 0.6, 3.55, 5.8, "Critical", "Only 2 legal partners handling all vendor contracts")
add_pain_point(s33, 0.6, 4.0, 5.8, "Critical", "CLM system requested for 5-6 years without funding")
add_pain_point(s33, 0.6, 4.45, 5.8, "High", "Sourcing team makes quasi-legal decisions due to legal bottleneck")
add_table(s33, 6.8, 3.1, 5.5, [
    ["Task", "Owner"],
    ["Negotiate Contract Terms", "Contracting Lane (R/A)"],
    ["Finalize and Execute Contract", "Contracting Lane (R/A)"],
    ["Vendor Contract Review", "Vendor Pool"],
    ["Vendor Contract Execution", "Vendor Pool"],
])
add_quote(s33, 6.8, 4.9, 5.5, "Contract lifecycle management system requested for 5-6 years without funding.", "-- Vendor Risk Lead")
add_footer(s33)

s34 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s34, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(s34, "Commercial Counsel -- Recommendations", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_card(s34, 0.6, 1.0, 5.9, 1.5, "Finance Rework Loop (GAP-4)", "Minor corrections loop back without formal denial. Only substantive contractual issues require process restart. Saves weeks per correction cycle.", EMERALD)
add_card(s34, 6.7, 1.0, 5.9, 1.5, "Vendor Affinity Streamlined Contracting", "VA pathway: compliance review + contract execution only. Skip negotiation (pre-approved terms). 5-day contracting SLA vs. 2-3 months.", TEAL)
add_card(s34, 0.6, 2.65, 5.9, 1.5, "Contract Execution via Message Flow", "Structured cross-pool communication: draft to vendor, vendor redlines, signed contract returned. 5-day SLA timer with escalation.", GOLD)
add_card(s34, 6.7, 2.65, 5.9, 1.5, "Template-Based Contracting", "Pre-approved MSA/SLA templates for common vendor types. Standard terms for VA products. Custom negotiation reserved for High risk only.", BLUE)
legal_rm = [
    ("Day 30", None, ["Document contract templates", "Define correction criteria", "Define VA contracting scope"]),
    ("Day 60", None, ["Template library established", "Vendor message flow standardized"]),
    ("Day 90", None, ["Rework loop active in SP4", "VA contracting deployed", "5-day SLA timers"]),
    ("Day 120", None, ["Measure: contracting cycle time", "Redlining iterations", "Template usage rate"]),
]
for i, ((t, sub, items), (bg, brd, tc)) in enumerate(zip(legal_rm, rm_colors)):
    add_roadmap_box(s34, 0.6 + i * 3.1, 4.4, 2.9, 2.5, t, sub, items, bg, brd, tc)
add_footer(s34)

# TOPIC 11: TPRM
add_section_slide("11. Third-Party Risk Management", "Cross-Cutting: Full Lifecycle", "Vendor lifecycle management, ongoing monitoring, and concentration risk. The thread that runs through every phase -- highest task count of any governance topic.")

s36 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s36, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(s36, "TPRM -- Current State", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_pain_point(s36, 0.6, 1.0, 5.8, "Critical", "No authoritative source for app/technology ownership")
add_pain_point(s36, 0.6, 1.45, 5.8, "Critical", "Incomplete CMDB -- unclear who owns support, requirements, lifecycle")
add_pain_point(s36, 0.6, 1.9, 5.8, "High", "Tech owner identification less rigorous than business owner identification")
add_pain_point(s36, 0.6, 2.35, 5.8, "High", "Single vendor spanning 10+ business units, each requiring separate assessment")
add_pain_point(s36, 0.6, 2.8, 5.8, "High", "1,200+ institutions requiring documentation and compliance tracking")
add_table(s36, 6.8, 1.0, 5.5, [
    ["Role", "State"],
    ["Business Owner", "Reasonably well-identified"],
    ["Vendor Owner", "Identified by risk team (team of 6)"],
    ["Product Owner", "Weakly tracked; often differs from requester"],
])
add_body_text(s36, "TPRM touches 13 tasks across all 5 phases + Vendor Pool -- highest task count of any governance topic. Governance lane is R/A.", 6.8, 2.8, 5.5, 0.5, size=10, color=DARK_TEXT)
add_footer(s36)

s37 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s37, RGBColor(0xF8, 0xFA, 0xFC))
add_title_text(s37, "TPRM -- Recommendations", 0.6, 0.3, 10, 0.5, size=24, color=DEEP_BLUE)
add_card(s37, 0.6, 1.0, 5.9, 1.5, "Mandatory Ownership Assignment (GAP-14)", "SP5 Task_AssignOwnership: required Business, Technical, Support Owner fields before close. Feeds CMDB. Annual ownership validation.", EMERALD)
add_card(s37, 6.7, 1.0, 5.9, 1.5, "Vendor-Level Aggregation", "When vendor spans 10+ BUs, aggregate vendor-level compliance artifacts. Separate request-level evaluation from vendor-level due diligence.", TEAL)
add_card(s37, 0.6, 2.65, 5.9, 1.5, "Risk Tier Drives Monitoring Cadence", "OB-DMN-1 risk tier determines post-onboarding monitoring: High = quarterly, Limited = semi-annual, Minimal = annual. Automated SLA timers.", GOLD)
add_card(s37, 6.7, 2.65, 5.9, 1.5, "Status Visibility for All Stakeholders (GAP-7)", "Real-time status at every phase transition. Requestors, risk team, and oversight all see current position. Portfolio view.", BLUE)
add_diagram(s37, "sp5-uat-golive.png", 0.6, 4.3, 12.0, 1.0)
tprm_rm = [
    ("Day 30", None, ["Define ownership fields", "Document aggregation rules", "Map monitoring to risk tier"]),
    ("Day 60", None, ["Ownership assignment in SP5", "Status notifications deployed"]),
    ("Day 90", None, ["Vendor aggregation deployed", "Risk-based monitoring cadence"]),
    ("Day 120", None, ["Workload dashboard", "Measure: ownership completeness", "Vendor DD reuse rate"]),
]
for i, ((t, sub, items), (bg, brd, tc)) in enumerate(zip(tprm_rm, rm_colors)):
    add_roadmap_box(s37, 0.6 + i * 3.1, 5.5, 2.9, 1.6, t, sub, items, bg, brd, tc)
add_footer(s37)

# ============================================================
# SLIDE 38: Next Steps
# ============================================================
slide38 = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide38)
add_title_text(slide38, "Next Steps", 0.8, 0.8, 10, 0.6, size=36, color=WHITE)
add_body_text(slide38, "Discovery complete. 11 topic areas assessed. 16 gaps identified and prioritized. The 30/60/90/120-day roadmap begins with planning and consolidation -- zero automation in the first 30 days.", 0.8, 1.5, 10, 0.6, size=13, color=RGBColor(0xCC, 0xCC, 0xDD))

next_steps = [
    ("1. Executive Alignment (Week 1)", "Review discovery findings with executive sponsors. Confirm committee consolidation approach. Assign owners to each Day 30 deliverable.", BLUE),
    ("2. Standards Definition (Weeks 2-3)", "Intake field inventory, prioritization scoring formula, security baseline controls, Vendor Affinity pathway rules, deal-killer no-go list.", TEAL),
    ("3. Quick Wins Target (Day 60)", "Unified intake, completeness gate, deal-killer pre-screen, status notifications, prioritization scoring, ownership assignment.", GOLD),
    ("4. Full Operation (Day 120)", "3-pathway routing, 5 parallel evaluation streams, AI fast-track, workload dashboard, before/after metrics. Target: 90%+ improvement.", EMERALD),
]
for i, (title, desc, color) in enumerate(next_steps):
    y = 2.4 + i * 1.15
    shape = slide38.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y), Inches(11.5), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x44)
    shape.line.fill.background()
    bar = slide38.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y + 0.05), Inches(0.06), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tf = slide38.shapes.add_textbox(Inches(1.1), Inches(y + 0.1), Inches(10.5), Inches(0.3))
    tf.text_frame.paragraphs[0].text = title
    tf.text_frame.paragraphs[0].font.size = Pt(13)
    tf.text_frame.paragraphs[0].font.bold = True
    tf.text_frame.paragraphs[0].font.color.rgb = WHITE
    tf2 = slide38.shapes.add_textbox(Inches(1.1), Inches(y + 0.45), Inches(10.5), Inches(0.45))
    tf2.text_frame.word_wrap = True
    tf2.text_frame.paragraphs[0].text = desc
    tf2.text_frame.paragraphs[0].font.size = Pt(10)
    tf2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
add_footer(slide38, dark=True)

# Save
output_path = "/Users/proth/repos/sla/customers/fs-onboarding/presentations/reference/v2-roadmap.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
