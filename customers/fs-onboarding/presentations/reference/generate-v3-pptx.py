#!/usr/bin/env python3
"""Generate KPMG-branded PowerPoint v3: Discovery Deep Dive and Implementation Roadmap.

Sources:
  - 14 stakeholder sessions (Feb-Mar 2026)
  - 5 deep-dive sessions (Mar 5-6, 2026)
  - 24 gaps across 11 governance domains
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── KPMG Brand Colors ──────────────────────────────────────────────────────────
KPMG_BLUE = RGBColor(0x00, 0x33, 0x8D)
KPMG_MEDIUM_BLUE = RGBColor(0x00, 0x5E, 0xB8)
KPMG_LIGHT_BLUE = RGBColor(0x00, 0x91, 0xDA)
KPMG_DARK_NAVY = RGBColor(0x00, 0x1D, 0x48)
KPMG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
KPMG_LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
KPMG_DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
KPMG_GOLD = RGBColor(0xD4, 0xA8, 0x43)
KPMG_EMERALD = RGBColor(0x00, 0xA8, 0x6B)
KPMG_ROSE = RGBColor(0xC6, 0x28, 0x28)
KPMG_AMBER = RGBColor(0xE6, 0x7E, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


# ── Helper Functions ────────────────────────────────────────────────────────────

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_gradient_bg(slide, c1, c2):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[1].color.rgb = c2
    fill.gradient_angle = 135


def add_text_box(slide, left, top, width, height, text, font_size=12,
                 bold=False, color=KPMG_DARK_GRAY, font_name='Open Sans',
                 alignment=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_footer(slide, text="KPMG Confidential", dark=False):
    color = RGBColor(0x99, 0x99, 0x99) if not dark else RGBColor(0x99, 0xAA, 0xCC)
    add_text_box(slide, Inches(0.5), Inches(7.0), Inches(4), Inches(0.4),
                 text, font_size=8, color=color)
    add_text_box(slide, Inches(9), Inches(7.0), Inches(4), Inches(0.4),
                 "March 2026", font_size=8, color=color, alignment=PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, title, body, accent_color=KPMG_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = KPMG_WHITE
    shape.line.fill.background()
    shape.shadow.inherit = False
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    add_text_box(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.3), Inches(0.35),
                 title, font_size=11, bold=True, color=KPMG_BLUE)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.45), width - Inches(0.3), height - Inches(0.55),
                 body, font_size=9, color=RGBColor(0x66, 0x66, 0x66))


def add_metric(slide, left, top, value, label, value_color=KPMG_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = KPMG_WHITE
    shape.line.fill.background()
    add_text_box(slide, left, top + Inches(0.1), Inches(2), Inches(0.5),
                 value, font_size=24, bold=True, color=value_color,
                 font_name='Open Sans Condensed', alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(0.65), Inches(2), Inches(0.3),
                 label.upper(), font_size=7, color=RGBColor(0x66, 0x66, 0x66),
                 alignment=PP_ALIGN.CENTER)


def add_section_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide, KPMG_MEDIUM_BLUE, KPMG_BLUE)
    add_text_box(slide, Inches(1), Inches(2.5), Inches(10), Inches(1.2),
                 title, font_size=40, bold=True, color=KPMG_WHITE, font_name='Open Sans Condensed')
    if subtitle:
        add_text_box(slide, Inches(1), Inches(3.8), Inches(9), Inches(0.8),
                     subtitle, font_size=14, color=RGBColor(0xCC, 0xDD, 0xFF))
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.55), Inches(3), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = KPMG_GOLD
    line.line.fill.background()
    add_footer(slide, dark=True)
    return slide


def add_current_state_slide(title, findings, quote=None, quote_source=None, raci=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, KPMG_WHITE)
    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
                 title, font_size=24, bold=True, color=KPMG_BLUE, font_name='Open Sans Condensed')

    cols = 2
    rows = (len(findings) + 1) // 2
    for i, (ftitle, fbody, is_positive) in enumerate(findings):
        col = i % cols
        row = i // cols
        left = Inches(0.6 + col * 6.2)
        top = Inches(1.2 + row * 1.15)
        accent = KPMG_EMERALD if is_positive else KPMG_ROSE
        add_card(slide, left, top, Inches(5.9), Inches(1.0), ftitle, fbody, accent)

    if quote:
        y = Inches(1.2 + rows * 1.15 + 0.1)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), y, Inches(0.06), Inches(0.7))
        bar.fill.solid()
        bar.fill.fore_color.rgb = KPMG_GOLD
        bar.line.fill.background()
        add_text_box(slide, Inches(0.8), y, Inches(11), Inches(0.5),
                     f'"{quote}"', font_size=10, color=KPMG_DARK_GRAY)
        if quote_source:
            add_text_box(slide, Inches(0.8), y + Inches(0.45), Inches(11), Inches(0.25),
                         f"-- {quote_source}", font_size=8, color=RGBColor(0x66, 0x66, 0x66))

    if raci:
        add_text_box(slide, Inches(0.6), Inches(6.7), Inches(12), Inches(0.25),
                     raci, font_size=8, color=RGBColor(0x66, 0x66, 0x66))
    add_footer(slide)
    return slide


def add_recommendations_slide(title, recs, roadmap):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, KPMG_WHITE)
    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
                 title, font_size=24, bold=True, color=KPMG_BLUE, font_name='Open Sans Condensed')

    for i, (rtitle, rbody, impact) in enumerate(recs[:4]):
        col = i % 2
        row = i // 2
        left = Inches(0.6 + col * 6.2)
        top = Inches(1.1 + row * 1.3)
        add_card(slide, left, top, Inches(5.9), Inches(1.15), rtitle,
                 f"{rbody}\n{impact}", KPMG_EMERALD)

    y_road = Inches(3.9)
    colors = [KPMG_EMERALD, KPMG_LIGHT_BLUE, KPMG_MEDIUM_BLUE, KPMG_BLUE]
    labels = ["30 Days", "60 Days", "90 Days", "120 Days"]
    for i, (items, label) in enumerate(zip(roadmap, labels)):
        left = Inches(0.6 + i * 3.1)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y_road, Inches(2.9), Inches(0.05))
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors[i]
        bar.line.fill.background()
        add_text_box(slide, left, y_road + Inches(0.08), Inches(2.9), Inches(0.3),
                     label, font_size=10, bold=True, color=KPMG_BLUE)
        text = "\n".join(f"  {item}" for item in items)
        add_text_box(slide, left, y_road + Inches(0.4), Inches(2.9), Inches(2.8),
                     text, font_size=8, color=RGBColor(0x66, 0x66, 0x66))

    add_footer(slide)
    return slide


def add_bpmn_desc_slide(title, sp_id, description):
    """Add a BPMN process description slide (text-only, no image)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, KPMG_LIGHT_GRAY)
    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
                 f"Process Model: {title}", font_size=24, bold=True,
                 color=KPMG_BLUE, font_name='Open Sans Condensed')
    # SP ID badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(0.6), Inches(1.2), Inches(2.5), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = KPMG_BLUE
    badge.line.fill.background()
    add_text_box(slide, Inches(0.6), Inches(1.22), Inches(2.5), Inches(0.36),
                 sp_id, font_size=10, bold=True, color=KPMG_WHITE, alignment=PP_ALIGN.CENTER)
    # Description
    add_text_box(slide, Inches(0.6), Inches(1.9), Inches(12), Inches(4.5),
                 description, font_size=12, color=KPMG_DARK_GRAY)
    add_footer(slide)
    return slide


def add_table_slide(title, headers, rows, subtitle=None):
    """Add a slide with a text-based table layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, KPMG_WHITE)
    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
                 title, font_size=24, bold=True, color=KPMG_BLUE, font_name='Open Sans Condensed')
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.4),
                     subtitle, font_size=11, color=KPMG_DARK_GRAY)

    y_start = Inches(1.6) if subtitle else Inches(1.3)
    n_cols = len(headers)
    col_w = 12.0 / n_cols

    # Header row
    for j, h in enumerate(headers):
        left = Inches(0.6 + j * col_w)
        hbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y_start, Inches(col_w - 0.1), Inches(0.4))
        hbar.fill.solid()
        hbar.fill.fore_color.rgb = KPMG_BLUE
        hbar.line.fill.background()
        add_text_box(slide, left + Inches(0.1), y_start + Inches(0.05), Inches(col_w - 0.3), Inches(0.3),
                     h, font_size=9, bold=True, color=KPMG_WHITE)

    # Data rows
    for r_idx, row in enumerate(rows):
        row_y = y_start + Inches(0.45 + r_idx * 0.55)
        bg_color = KPMG_WHITE if r_idx % 2 == 0 else KPMG_LIGHT_GRAY
        for j, cell in enumerate(row):
            left = Inches(0.6 + j * col_w)
            cell_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, row_y,
                                              Inches(col_w - 0.1), Inches(0.5))
            cell_bg.fill.solid()
            cell_bg.fill.fore_color.rgb = bg_color
            cell_bg.line.fill.background()
            add_text_box(slide, left + Inches(0.1), row_y + Inches(0.05),
                         Inches(col_w - 0.3), Inches(0.4),
                         cell, font_size=8, color=KPMG_DARK_GRAY)

    add_footer(slide)
    return slide


# ╔═══════════════════════════════════════════════════════════════════════════════╗
# ║  OPENING SLIDES (1-5)                                                       ║
# ╚═══════════════════════════════════════════════════════════════════════════════╝

# ── Slide 1: Title ──────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, KPMG_BLUE, KPMG_MEDIUM_BLUE)
add_text_box(slide, Inches(1), Inches(1.5), Inches(10), Inches(1.5),
             "Software Onboarding\nProcess Transformation", font_size=44, bold=True,
             color=KPMG_WHITE, font_name='Open Sans Condensed')
add_text_box(slide, Inches(1), Inches(3.3), Inches(10), Inches(0.6),
             "v3 Discovery Deep Dive and Implementation Roadmap",
             font_size=18, bold=True, color=RGBColor(0xCC, 0xDD, 0xFF))
add_text_box(slide, Inches(1), Inches(4.2), Inches(9), Inches(1),
             "Based on 14 stakeholder interviews across Architecture, Product, Security, Risk Management, "
             "Vendor Management, Finance, Legal, and Compliance teams. 24 gaps identified across "
             "11 governance domains with actionable 30/60/90/120-day implementation plans.",
             font_size=11, color=RGBColor(0xBB, 0xCC, 0xEE))
badges = ["14 Stakeholder Sessions", "11 Governance Domains", "24 Gap Findings", "120-Day Roadmap"]
for i, b in enumerate(badges):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1 + i * 2.6), Inches(5.5), Inches(2.3), Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x00, 0x2A, 0x70)
    shape.line.color.rgb = RGBColor(0x33, 0x55, 0xAA)
    shape.line.width = Pt(1)
    add_text_box(slide, Inches(1 + i * 2.6), Inches(5.53), Inches(2.3), Inches(0.35),
                 b, font_size=9, color=RGBColor(0xCC, 0xDD, 0xFF), alignment=PP_ALIGN.CENTER)
add_footer(slide, dark=True)

# ── Slide 2: Executive Summary ──────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, KPMG_WHITE)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
             "Executive Summary", font_size=28, bold=True, color=KPMG_BLUE, font_name='Open Sans Condensed')
add_text_box(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.7),
             "The software onboarding process spans 6 to 9 months end-to-end, driven by 18 sequential "
             "committees, 5+ disconnected intake channels, and critical resource bottlenecks in Security and Legal. "
             "The START initiative created centralized awareness but did not integrate underlying team processes.",
             font_size=11, color=KPMG_DARK_GRAY)

metrics = [("6-9 mo", "Current E2E Cycle", KPMG_ROSE), ("18", "Committees", KPMG_ROSE),
           ("335/yr", "Assessments", KPMG_BLUE), ("24", "Gaps Identified", KPMG_AMBER)]
for i, (val, lab, clr) in enumerate(metrics):
    add_metric(slide, Inches(0.6 + i * 2.5), Inches(1.8), val, lab, clr)

add_text_box(slide, Inches(0.6), Inches(3.1), Inches(6), Inches(0.3),
             "Highest-ROI Investment Areas", font_size=14, bold=True, color=KPMG_DARK_NAVY,
             font_name='Open Sans Condensed')
investments = [
    ("Parallel Evaluation", "Replace 18 sequential committees with 5 parallel evaluation streams. Bounded by slowest SLA, not sequential sum."),
    ("Unified Intake + Deal-Killer", "Consolidate 5+ channels. Block non-starters at day 1. Route 3 request types differently."),
    ("Contract Automation", "Automate contract review for the 2-person team handling 30+ contracts/month."),
]
for i, (it, ib) in enumerate(investments):
    add_card(slide, Inches(0.6 + i * 4.1), Inches(3.45), Inches(3.9), Inches(1.1), it, ib, KPMG_EMERALD)

# Quote
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.8), Inches(0.06), Inches(0.7))
bar.fill.solid()
bar.fill.fore_color.rgb = KPMG_GOLD
bar.line.fill.background()
add_text_box(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(0.5),
             '"Somebody needs to be empowered at the firm to say I own this"',
             font_size=10, color=KPMG_DARK_GRAY)
add_text_box(slide, Inches(0.8), Inches(5.25), Inches(11), Inches(0.25),
             "-- TPRM Lead", font_size=8, color=RGBColor(0x66, 0x66, 0x66))
add_footer(slide)

# ── Slide 3: E2E Workflow Pain Points ────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, KPMG_WHITE)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
             "End-to-End Workflow: Current Pain Points", font_size=24, bold=True,
             color=KPMG_BLUE, font_name='Open Sans Condensed')

phases = [
    ("SP1: Intake", "3 request types\nnot routed", KPMG_ROSE),
    ("SP2: Prioritize", "Whoever screams\nloudest", KPMG_ROSE),
    ("SP3: Evaluate", "Sequential DART\n3x AI vendor contact", KPMG_ROSE),
    ("SP4: Contract", "2 ppl / 30+ contracts\n1.5yr security", KPMG_ROSE),
    ("SP5: Go-Live", "No conditional\napproval", KPMG_AMBER),
]
for i, (name, issue, color) in enumerate(phases):
    left = Inches(0.4 + i * 2.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.2), Inches(2.3), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_text_box(slide, left, Inches(1.22), Inches(2.3), Inches(0.4),
                 name, font_size=10, bold=True, color=KPMG_WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(1.6), Inches(2.3), Inches(0.5),
                 issue, font_size=8, color=RGBColor(0xFF, 0xDD, 0xDD), alignment=PP_ALIGN.CENTER)

# Bottleneck cards
add_text_box(slide, Inches(0.6), Inches(2.5), Inches(6), Inches(0.3),
             "Critical Bottlenecks", font_size=14, bold=True, color=KPMG_DARK_NAVY,
             font_name='Open Sans Condensed')
bottlenecks = [
    ("Contract Negotiation", "2 people / 30+ contracts monthly. Up to 1.5 years for security exhibits."),
    ("Security Review", '"Biggest bottleneck... reason our SLA takes 2 weeks." Cannot scale for AI.'),
    ("Sequential DART", "Requester contacts 5-6 teams independently. No visibility into sequence."),
    ("AI Governance", "60+ items in queue. 3 committees. 3x vendor contact for AI reviews."),
]
for i, (bt, bb) in enumerate(bottlenecks):
    col = i % 2
    row = i // 2
    add_card(slide, Inches(0.6 + col * 6.2), Inches(2.9 + row * 1.15),
             Inches(5.9), Inches(1.0), bt, bb, KPMG_ROSE)

# Strengths
add_text_box(slide, Inches(0.6), Inches(5.3), Inches(6), Inches(0.3),
             "What Works Well", font_size=14, bold=True, color=KPMG_DARK_NAVY,
             font_name='Open Sans Condensed')
strengths = [
    ("Architecture Governance", "Dedicated Facilitator pre-screens artifacts, manages follow-ups. Best-practice model."),
    ("TPRM Improvement", "Reduced DD from 144 to 75 days. 335 assessments/year with 8-person team."),
    ("AI Tool Adoption", "Cursor AI for diagrams, custom AI agent for pattern conformance scanning."),
]
for i, (st, sb) in enumerate(strengths):
    add_card(slide, Inches(0.6 + i * 4.1), Inches(5.6), Inches(3.9), Inches(0.95), st, sb, KPMG_EMERALD)

add_footer(slide)

# ── Slide 4: System Landscape ────────────────────────────────────────────────────
add_table_slide(
    "System Landscape and Integration Gaps",
    ["System", "Function", "Integration Status"],
    [
        ["ServiceNow (START)", "Process intake", "Entry point, no downstream integration"],
        ["Ariba", "Registration, NDAs, contracts", "API to Oracle only"],
        ["OneTrust", "Risk assessments, tracking, control gaps", "Standalone, manual PDF exports"],
        ["Oracle", "Accounts payable", "API from Ariba"],
        ["Confluence + Catfox", "Architecture approvals", "Plugin-based, no workflow integration"],
        ["JIRA", "Architecture ticket tracking", "Per-team, no cross-functional view"],
        ["Camunda 8 (target)", "E2E process orchestration", "Not yet deployed"],
    ],
    subtitle='"We have all these systems but none of them talk to each other" -- TPRM Lead'
)

# ── Slide 5: By the Numbers ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, KPMG_WHITE)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
             "By the Numbers: Quantified Pain", font_size=28, bold=True,
             color=KPMG_BLUE, font_name='Open Sans Condensed')

num_metrics = [
    ("335/yr", "Assessments", KPMG_BLUE),
    ("28-29d", "RAE Actual\n(14d target)", KPMG_ROSE),
    ("75d", "DD Internal Review", KPMG_ROSE),
    ("2 ppl", "30+ Contracts/mo", KPMG_ROSE),
]
for i, (val, lab, clr) in enumerate(num_metrics):
    add_metric(slide, Inches(0.6 + i * 3.1), Inches(1.2), val, lab, clr)

num_metrics_2 = [
    ("1.5yr", "Security Exhibit\nNegotiation", KPMG_ROSE),
    ("2-3/10", "Business Council\nQuorum", KPMG_AMBER),
    ("60-90d", "Competitor\nBenchmark", KPMG_LIGHT_BLUE),
    ("18", "Committees\n(target: 5 streams)", KPMG_AMBER),
]
for i, (val, lab, clr) in enumerate(num_metrics_2):
    add_metric(slide, Inches(0.6 + i * 3.1), Inches(2.8), val, lab, clr)

# Bottom quote
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.4), Inches(0.06), Inches(0.7))
bar.fill.solid()
bar.fill.fore_color.rgb = KPMG_GOLD
bar.line.fill.background()
add_text_box(slide, Inches(0.8), Inches(4.4), Inches(11), Inches(0.5),
             '"That team desperately needs automation... The firm needs them to be automated"',
             font_size=10, color=KPMG_DARK_GRAY)
add_text_box(slide, Inches(0.8), Inches(4.85), Inches(11), Inches(0.25),
             "-- TPRM Lead (on Legal/Contracts)", font_size=8, color=RGBColor(0x66, 0x66, 0x66))
add_footer(slide)


# ╔═══════════════════════════════════════════════════════════════════════════════╗
# ║  CROSS-CUTTING SLIDES (6-9)                                                 ║
# ╚═══════════════════════════════════════════════════════════════════════════════╝

# ── Slide 6: Section Divider ─────────────────────────────────────────────────────
add_section_slide("Cross-Cutting Process Design",
                  "Concierge model, simultaneous engagement, and organizational design patterns "
                  "that apply across all 11 governance domains.")

# ── Slide 7: Concierge / Quarterback Model ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, KPMG_WHITE)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
             "Process Quarterback / Concierge Model", font_size=24, bold=True,
             color=KPMG_BLUE, font_name='Open Sans Condensed')

add_text_box(slide, Inches(0.6), Inches(1.1), Inches(5.8), Inches(0.3),
             "Architecture Today (Best Practice)", font_size=14, bold=True,
             color=KPMG_DARK_NAVY, font_name='Open Sans Condensed')
arch_today = [
    ("Pre-Screen Artifacts", "Remove incomplete items from ARB/SDRB agenda before meeting."),
    ("Manage Follow-Ups", "Capture action items, track completion, run meetings efficiently."),
    ("JIRA Integration", "Ticket tracking and domain-based auto-assignment to 4 EA leaders."),
]
for i, (at, ab) in enumerate(arch_today):
    add_card(slide, Inches(0.6), Inches(1.45 + i * 0.95), Inches(5.8), Inches(0.85), at, ab, KPMG_EMERALD)

add_text_box(slide, Inches(6.8), Inches(1.1), Inches(5.8), Inches(0.3),
             "Proposed E2E Extension", font_size=14, bold=True,
             color=KPMG_DARK_NAVY, font_name='Open Sans Condensed')
proposed = [
    ("Single Point of Contact", "Eliminates 5-6 team self-navigation for requesters."),
    ("Automated DART Formation", "Concierge triggers all review streams simultaneously."),
    ("Quality Gates at Boundaries", "Reject incomplete before wasting SME time."),
]
for i, (pt, pb) in enumerate(proposed):
    add_card(slide, Inches(6.8), Inches(1.45 + i * 0.95), Inches(5.8), Inches(0.85), pt, pb, KPMG_LIGHT_BLUE)

# Quotes
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.5), Inches(0.06), Inches(1.2))
bar.fill.solid()
bar.fill.fore_color.rgb = KPMG_GOLD
bar.line.fill.background()
add_text_box(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.5),
             '"If any artifact is not in good order, remove it from the agenda"',
             font_size=10, color=KPMG_DARK_GRAY)
add_text_box(slide, Inches(0.8), Inches(4.9), Inches(11), Inches(0.25),
             "-- Architecture Lead", font_size=8, color=RGBColor(0x66, 0x66, 0x66))
add_text_box(slide, Inches(0.8), Inches(5.25), Inches(11), Inches(0.5),
             '"Exactly like what we\'re talking about for the quarterback from a broader end-to-end"',
             font_size=10, color=KPMG_DARK_GRAY)
add_text_box(slide, Inches(0.8), Inches(5.65), Inches(11), Inches(0.25),
             "-- Consulting Team", font_size=8, color=RGBColor(0x66, 0x66, 0x66))
add_footer(slide)

# ── Slide 8: Simultaneous Engagement ────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, KPMG_WHITE)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
             "Simultaneous Engagement Model", font_size=24, bold=True,
             color=KPMG_BLUE, font_name='Open Sans Condensed')

# Current state - sequential
add_text_box(slide, Inches(0.6), Inches(1.1), Inches(5.8), Inches(0.3),
             "Current: Sequential (weeks to months)", font_size=14, bold=True,
             color=KPMG_ROSE, font_name='Open Sans Condensed')
seq_steps = [
    "TBC approval",
    "Requester contacts Architecture... wait...",
    "Requester contacts Security... wait...",
    "Requester contacts Compliance... wait...",
    "Requester contacts AI Gov... wait...",
    "Timeline: dependent on requester initiative",
]
for i, s in enumerate(seq_steps):
    add_text_box(slide, Inches(0.8), Inches(1.5 + i * 0.32), Inches(5.5), Inches(0.3),
                 f"  {s}", font_size=9, color=KPMG_DARK_GRAY)

# Proposed state - parallel
add_text_box(slide, Inches(6.8), Inches(1.1), Inches(5.8), Inches(0.3),
             "Proposed: Parallel (bounded by slowest SLA)", font_size=14, bold=True,
             color=KPMG_EMERALD, font_name='Open Sans Condensed')
par_streams = [
    ("Architecture", "HLD review, domain assignment"),
    ("Security", "Tiered assessment by risk category"),
    ("Compliance", "Regulatory review + OneTrust evidence"),
    ("AI Governance", "Consolidated single-stream review"),
    ("Legal", "Contract terms + NDA execution"),
]
for i, (stream, desc) in enumerate(par_streams):
    add_card(slide, Inches(6.8), Inches(1.5 + i * 0.7), Inches(5.8), Inches(0.6),
             stream, desc, KPMG_EMERALD)

# Quote
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.2), Inches(0.06), Inches(0.7))
bar.fill.solid()
bar.fill.fore_color.rgb = KPMG_GOLD
bar.line.fill.background()
add_text_box(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(0.5),
             '"I disagree with that because that makes us a bottleneck... '
             'there should be simultaneous engagement of all the major players that have a vote"',
             font_size=10, color=KPMG_DARK_GRAY)
add_text_box(slide, Inches(0.8), Inches(5.65), Inches(11), Inches(0.25),
             "-- Architecture Lead", font_size=8, color=RGBColor(0x66, 0x66, 0x66))
add_footer(slide)

# ── Slide 9: 3 Request Types + Pod Model ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, KPMG_WHITE)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
             "3 Request Types and Distributed Pod Model", font_size=24, bold=True,
             color=KPMG_BLUE, font_name='Open Sans Condensed')

# Request types table
add_text_box(slide, Inches(0.6), Inches(1.1), Inches(6), Inches(0.3),
             "3 Request Types (New in v3)", font_size=14, bold=True,
             color=KPMG_DARK_NAVY, font_name='Open Sans Condensed')
req_types = [
    ("Defined Need", "Business owner knows requirements, vendor selected. Standard 5-phase path. Most common."),
    ("Forced Update", "Existing vendor, product changes (on-prem to SaaS, EOL, new AI). Re-evaluation at SP3."),
    ("Speculative / Exploratory", "Advisory support, no sponsorship. Idea funnel (pre-SP1), NOT standard process."),
]
for i, (rt, rd) in enumerate(req_types):
    accent = [KPMG_BLUE, KPMG_AMBER, KPMG_LIGHT_BLUE][i]
    add_card(slide, Inches(0.6), Inches(1.5 + i * 0.95), Inches(5.8), Inches(0.85), rt, rd, accent)

# Pod model
add_text_box(slide, Inches(6.8), Inches(1.1), Inches(5.8), Inches(0.3),
             "Distributed Pod Model", font_size=14, bold=True,
             color=KPMG_DARK_NAVY, font_name='Open Sans Condensed')
pods = [
    ("Cybersecurity Pod", "Own prioritization, meeting frequency, review speed"),
    ("Architecture Pod", "Technical review cadence, domain assignment"),
    ("Legal/Contracts Pod", "Template usage, negotiation approach"),
    ("AI Governance Pod", "AI risk posture, review depth"),
    ("TPRM Pod", "Assessment methodology, vendor scoring"),
]
for i, (pn, pd) in enumerate(pods):
    add_card(slide, Inches(6.8), Inches(1.5 + i * 0.73), Inches(5.8), Inches(0.63),
             pn, f"{pd}  |  Central team ensures consistency", KPMG_MEDIUM_BLUE)

# Quote
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.3), Inches(0.06), Inches(0.7))
bar.fill.solid()
bar.fill.fore_color.rgb = KPMG_GOLD
bar.line.fill.background()
add_text_box(slide, Inches(0.8), Inches(5.3), Inches(11), Inches(0.5),
             '"The process never solved the step of requester who\'s never been through onboarding"',
             font_size=10, color=KPMG_DARK_GRAY)
add_text_box(slide, Inches(0.8), Inches(5.75), Inches(11), Inches(0.25),
             "-- Vendor Management Lead", font_size=8, color=RGBColor(0x66, 0x66, 0x66))
add_footer(slide)


# ╔═══════════════════════════════════════════════════════════════════════════════╗
# ║  ROADMAP SLIDE (10)                                                         ║
# ╚═══════════════════════════════════════════════════════════════════════════════╝

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, KPMG_WHITE)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
             "Implementation Roadmap: 30 / 60 / 90 / 120 Days", font_size=24, bold=True,
             color=KPMG_BLUE, font_name='Open Sans Condensed')

roadmap_data = [
    ("Days 1-30: Foundation", KPMG_EMERALD, [
        "Map all intake fields, deduplicate",
        "Define 3 request types + routing rules",
        "NDA timing decision (Sec + Legal)",
        "Map 18 committees to 5 streams",
        "Draft concierge role description",
        "Simultaneous engagement rules",
        "Deal-killer no-go list",
        "Prioritization scoring w/ capacity",
        "Baseline RAE (28d) and DD (75d) metrics",
        "Map OneTrust API integration points"]),
    ("Days 31-60: Quick Wins", KPMG_LIGHT_BLUE, [
        "Unified intake form (3 request types)",
        "Completeness quality gate at intake",
        "Deal-killer pre-screen active",
        "NDA gate between SP1 and SP3",
        "Status notifications at boundaries",
        "Async voting for Business Council",
        "Single AI questionnaire deployed",
        "OneTrust-Camunda integration design",
        "Contract template standardization",
        "Pilot simultaneous engagement"]),
    ("Days 61-90: Automation", KPMG_MEDIUM_BLUE, [
        "3-pathway routing: Buy/Build/Enable",
        "5 parallel evaluation streams live",
        "Concierge orchestration active",
        "Tiered security assessment in SP3",
        "AI fast-track piloted (2-week target)",
        "OneTrust assessment automation in SP3",
        "Automated baseline security checks",
        "AI-assisted pre-screening pilot",
        "Finance rework loop in SP4",
        "Capacity-aware routing"]),
    ("Days 91-120: Scale", KPMG_BLUE, [
        "Exception routing (rapid RA)",
        "Pre-onboarding idea funnel",
        "Workload visibility dashboard",
        "Contract review automation",
        "Full deviation reporting",
        "Distributed pod model pilot",
        "Progressive forms + stage awareness",
        "Annual ownership validation",
        "Executive KPI reporting",
        "Before/after metrics documented"]),
]

for i, (title, color, items) in enumerate(roadmap_data):
    left = Inches(0.5 + i * 3.15)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.2), Inches(2.95), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text_box(slide, left, Inches(1.3), Inches(2.95), Inches(0.35),
                 title, font_size=11, bold=True, color=KPMG_BLUE)
    text = "\n".join(f"  {item}" for item in items)
    add_text_box(slide, left, Inches(1.7), Inches(2.95), Inches(4.5),
                 text, font_size=8, color=RGBColor(0x66, 0x66, 0x66))

# ROI metrics
roi = [("60-70%", "Cycle Time Reduction"), ("18 to 5", "Sequential to Parallel"),
       ("3 Pathways", "Buy/Build/Enable"), ("Day 1", "Non-Starter ID")]
for i, (v, l) in enumerate(roi):
    left = Inches(0.6 + i * 3.15)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(6.2), Inches(2.9), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = KPMG_BLUE
    shape.line.fill.background()
    add_text_box(slide, left, Inches(6.22), Inches(2.9), Inches(0.4),
                 v, font_size=20, bold=True, color=KPMG_GOLD,
                 font_name='Open Sans Condensed', alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(6.6), Inches(2.9), Inches(0.3),
                 l, font_size=9, color=RGBColor(0xCC, 0xDD, 0xFF), alignment=PP_ALIGN.CENTER)

add_footer(slide)


# ╔═══════════════════════════════════════════════════════════════════════════════╗
# ║  11 GOVERNANCE TOPIC DEEP DIVES (Slides 11-54, 4 per topic)                ║
# ╚═══════════════════════════════════════════════════════════════════════════════╝

# ── Topic 1: Intake (slides 11-14) ──────────────────────────────────────────────
add_section_slide("Intake",
                  "Initial request capture, portfolio review, and solution discovery.")

add_bpmn_desc_slide("SP1: Request and Triage", "SP1_RequestTriage",
    "Request capture through unified intake gateway, existing solution review, "
    "bypass routing for leveraged software, initial triage with 2-day SLA timer, "
    "and DMN-driven pathway routing (OB-DMN-2). v3 adds 3 request type routing: "
    "Defined Need (standard), Forced Updates (re-evaluation at SP3), "
    "Speculative/Exploratory (idea funnel pre-SP1).")

add_current_state_slide(
    "Intake: Current State Findings",
    [
        ("5+ Disconnected Channels",
         "ServiceNow, REA (~80Q), AI forms, Power Apps, email, feedback platform, advisor tool approvals.",
         False),
        ("3 Request Types Treated Identically",
         "Defined Need, Forced Updates, Speculative all follow same heavyweight process.",
         False),
        ("Tool Decay",
         "Financial form shows 2024, locked, password unknown, owner left org.",
         False),
        ("No Pre-Screen",
         "Obvious no-gos reach SME teams. No deal-killer check.",
         False),
    ],
    quote="You submit a formal intake request so they have trackability",
    quote_source="Product Manager"
)

add_recommendations_slide(
    "Intake: Recommendations",
    [
        ("Unified Intake Gateway (GAP-1)",
         "Single entry absorbs all channels. Dynamic routing by request type.",
         "HIGH IMPACT"),
        ("NDA Gate (GAP-17)",
         "Mandatory NDA before detailed evaluation. Security+Legal aligned.",
         "NEW in v3"),
        ("Deal-Killer Pre-Screen (GAP-16)",
         "DMN-driven no-go list. Reject known non-starters immediately.",
         "QUICK WIN"),
        ("3 Request Type Routing",
         "OB-DMN-2 expanded with requestType. Speculative to idea funnel, Forced Updates skip intake.",
         "NEW in v3"),
    ],
    [
        ["Map fields", "NDA timing decision", "3 request types"],
        ["Unified form", "NDA gate", "Completeness gate"],
        ["Progressive forms", "Informal channel redirect"],
        ["Exception routing", "Idea funnel", "Metrics"],
    ]
)

# ── Topic 2: Prioritization (slides 15-18) ──────────────────────────────────────
add_section_slide("Prioritization",
                  "Scoring, queue management, and governance decision cadence.")

add_bpmn_desc_slide("SP2: Planning and Routing", "SP2_PlanningRouting",
    "Preliminary assessment, backlog management, and DMN-driven pathway routing (OB-DMN-2). "
    "v3 adds capacity-aware scoring via OB-DMN-5 with business impact, strategic alignment, "
    "urgency, risk tier, and capacity impact inputs producing P1/P2/P3 output.")

add_current_state_slide(
    "Prioritization: Current State Findings",
    [
        ("No Formal Standard",
         "Teams 'horse trade.' Each requestor views their request as most important.",
         False),
        ("Whoever Screams Loudest",
         "EVP support pushes reviews down. No effective prioritization possible.",
         False),
        ("Business Council Quorum",
         "Monthly meetings, 2-3 of 8-10 attend. Evolved to email voting.",
         False),
        ("Capacity Disconnected",
         "PI planning allocates capacity. Exception requests force pivots.",
         False),
    ],
    quote="If the general pace of onboarding is accelerated, then we could theoretically handle these one-off cases with more capacity",
    quote_source="Product Manager"
)

add_recommendations_slide(
    "Prioritization: Recommendations",
    [
        ("DMN-Driven Scoring (GAP-2)",
         "OB-DMN-5: business impact, alignment, urgency, risk tier, capacity. P1/P2/P3 output.",
         "HIGH IMPACT"),
        ("Async Governance (GAP-24)",
         "Replace monthly meetings with 48-hour async SLA. Email voting formalized.",
         "NEW in v3"),
        ("3-Pathway Routing (GAP-11)",
         "Buy, Build, Enable (Vendor Affinity). Right-sized governance per pathway.",
         "HIGH IMPACT"),
        ("Capacity-Aware Routing",
         "Score includes capacity impact. Route around constrained teams.",
         "NEW in v3"),
    ],
    [
        ["Capacity impact scoring", "Business Council quorum fix", "Async decision rules"],
        ["Async voting deployed", "Queue position visible"],
        ["Capacity-aware routing", "Acceleration metrics baseline"],
        ["Queue wait time metrics", "Capacity utilization", "Decision cycle time"],
    ]
)

# ── Topic 3: Funding (slides 19-22) ─────────────────────────────────────────────
add_section_slide("Funding",
                  "Financial analysis, budget approval, and coding matrix workflows.")

add_bpmn_desc_slide("SP2/SP3: Financial Touchpoints", "SP2_SP3_Finance",
    "Financial analysis spans 4 phases: SP1 (estimate), SP2 (justification), "
    "SP3 (full TCO/ROI), SP4 (final approval). Progressive financial data "
    "collection reduces requester burden. v3 adds Enable pathway bypass for "
    "Vendor Affinity products that require no org funding.")

add_current_state_slide(
    "Funding: Current State Findings",
    [
        ("Formal Denial Restarts Process",
         "Minor corrections require full restart. Cannot reroute to FP&A.",
         False),
        ("Enable Pathway Friction",
         "Vendor Affinity tools require no org funding but validation still required.",
         False),
        ("Outdated Tools",
         "Finance form locked, password unknown, completed offline.",
         False),
        ("4-Phase Financial Touchpoints",
         "Financial analysis in SP1, SP2, SP3, SP4. Progressive collection needed.",
         False),
    ],
    quote=None
)

add_recommendations_slide(
    "Funding: Recommendations",
    [
        ("Finance Rework Loop (GAP-4)",
         "Minor corrections loop without denial. Pattern matches PDLC retry.",
         "MEDIUM"),
        ("Enable Pathway Bypass",
         "Vendor Affinity products skip funding validation entirely.",
         "NEW in v3"),
        ("FP&A Integration",
         "Direct routing for coding matrix issues. No external email.",
         "MEDIUM"),
        ("Progressive Financial Data",
         "SP1: estimate. SP2: justification. SP3: full TCO/ROI.",
         "LOW EFFORT"),
    ],
    [
        ["Map financial fields across phases", "Enable bypass criteria"],
        ["Finance rework loop deployed", "FP&A direct routing"],
        ["Enable pathway financial bypass active", "Progressive forms"],
        ["Measure: rework rate", "Enable pathway adoption"],
    ]
)

# ── Topic 4: Sourcing (slides 23-26) ────────────────────────────────────────────
add_section_slide("Sourcing / Evaluation",
                  "DART formation, committee coordination, and vendor due diligence.")

add_bpmn_desc_slide("SP3: Evaluation and Due Diligence", "SP3_EvalDD",
    "5 parallel evaluation streams: Technical Architecture, Security Assessment, "
    "Risk/Compliance/Legal, Financial Analysis, Vendor Landscape Assessment. "
    "Vendor due diligence questionnaire (830 questions with skip logic). "
    "v3 replaces sequential DART formation with simultaneous engagement "
    "triggered by concierge. 5-day SLA timer on vendor response.")

add_current_state_slide(
    "Sourcing: Current State Findings",
    [
        ("Sequential DART Formation",
         "Requester independently contacts 5-6 teams. No visibility into sequence.",
         False),
        ("ARB + TBC + AI Gov Sequential",
         "Same case presented to multiple committees. Meetings weeks apart.",
         False),
        ("Security as Primary Bottleneck",
         "Security is our biggest bottleneck. 2-week SLA driver. Understaffed.",
         False),
        ("Governance Facilitator Excellence",
         "Architecture has dedicated facilitator: pre-screens, removes incomplete items, runs ARB/SDRB.",
         True),
    ],
    quote="I disagree with that because that makes us a bottleneck... simultaneous engagement of all major players",
    quote_source="Architecture Lead"
)

add_recommendations_slide(
    "Sourcing: Recommendations",
    [
        ("Simultaneous Engagement (GAP-19)",
         "Parallel engagement of all review streams. Concierge triggers automatically.",
         "HIGH IMPACT"),
        ("Governance Facilitator Model",
         "Scale Architecture's pre-screening model to E2E. Reject incomplete artifacts.",
         "NEW in v3"),
        ("Domain Auto-Assignment",
         "4 EA leaders route by domain. Auto-assign by bandwidth via OB-DMN-2.",
         "NEW in v3"),
        ("5 Parallel Evaluation Streams",
         "Technical, Security, AI, Compliance, Vendor Landscape simultaneously.",
         "HIGH IMPACT"),
    ],
    [
        ["ARB/SDRB/TBC scope boundaries", "Simultaneous engagement rules", "Concierge role"],
        ["Pilot simultaneous engagement", "Catfox integration", "Domain auto-assignment"],
        ["5 streams operational", "Artifact quality gate", "Concierge orchestration"],
        ["DART formation time (target: automated)", "Evaluation cycle time"],
    ]
)

# ── Topic 5: Cybersecurity (slides 27-30) ────────────────────────────────────────
add_section_slide("Cybersecurity",
                  "Security assessment, tiered review, and AI review consolidation.")

add_bpmn_desc_slide("SP3: Security Assessment Stream", "SP3_Security",
    "Tiered security assessment driven by DMN risk-based categories: "
    "local install, platform, hybrid, pre-approved, updates, module additions. "
    "Baseline tier (automated checks), Elevated tier (auto + manual), "
    "Major tier (full assessment + penetration testing). "
    "v3 consolidates 3x AI vendor contact into single stream.")

add_current_state_slide(
    "Cybersecurity: Current State Findings",
    [
        ("Biggest Bottleneck",
         "Security is the reason SLA takes 2 weeks. Understaffed team.",
         False),
        ("Cannot Scale for AI",
         '"I need three of me right now" -- specifically for AI reviews.',
         False),
        ("Zero Automation",
         "All routing manual. No workflow integration. Individual knowledge dependency.",
         False),
        ("3x Vendor Contact for AI",
         "Tech risk, cybersecurity, TPRM each contact vendors independently.",
         False),
    ],
    quote="Security is our biggest bottleneck... Enterprise architecture would get SLA way down without security constraints",
    quote_source="Architecture Lead"
)

add_recommendations_slide(
    "Cybersecurity: Recommendations",
    [
        ("Tiered Security Assessment (GAP-12)",
         "DMN-driven: Baseline (automated), Elevated (auto+manual), Major (full+pentest).",
         "HIGH IMPACT"),
        ("Risk-Based Categories",
         "6 types: local install, platform, hybrid, pre-approved, updates, module additions.",
         "NEW in v3"),
        ("AI Review Consolidation",
         "Single stream replaces 3x vendor contact. Eliminate redundancy.",
         "NEW in v3"),
        ("NDA Timing Resolution",
         "Security+Legal align on early NDA. Executive decision at Day 30.",
         "NEW in v3"),
    ],
    [
        ["6 risk-based assessment categories", "NDA timing resolution", "Map security capacity"],
        ["Automated baseline checks", "Tiered assessment DMN", "NDA gate active"],
        ["Security: Elevated/Major only", "AI vendor contact consolidated"],
        ["Security review cycle time by tier", "Baseline automation rate", "AI queue depth"],
    ]
)

# ── Topic 6: Enterprise Architecture (slides 31-34) ─────────────────────────────
add_section_slide("Enterprise Architecture",
                  "Architecture review boards, artifact frameworks, and domain governance.")

add_bpmn_desc_slide("SP3: Architecture Review Stream", "SP3_Architecture",
    "Architecture review within SP3 parallel evaluation. ARB (2-week SLA for complex), "
    "SDRB (same-day for prepared teams). Artifact framework: HLD pre-contract (all), "
    "SD/SYDD post-contract (local), PSS post-contract (enterprise/AI). "
    "v3 adds AI-assisted pre-screening and preparation quality gates.")

add_current_state_slide(
    "Enterprise Architecture: Current State Findings",
    [
        ("ARB 2-Week SLA, SDRB Same-Day",
         "Well-prepared teams get faster reviews. New capabilities are slower.",
         True),
        ("Artifact Framework",
         "HLD pre-contract (all). SD/SYDD post-contract (local). PSS post-contract (enterprise/AI).",
         True),
        ("AI Tools Already in Use",
         "Cursor AI for diagrams, custom AI agent for pattern conformance scanning.",
         True),
        ("Domain Funding Constraints",
         "Architects funded by domains. Under-funding = capacity bottleneck.",
         False),
    ],
    quote="We're trying to use AI to help expedite these reviews",
    quote_source="Architecture Lead"
)

add_recommendations_slide(
    "Enterprise Architecture: Recommendations",
    [
        ("Scale Facilitator Model",
         "Governance Facilitator becomes E2E concierge template. Pre-screen all artifacts.",
         "HIGH IMPACT"),
        ("AI-Assisted Pre-Screening",
         "Leverage existing Cursor/custom agent for pattern conformance automation.",
         "NEW in v3"),
        ("Preparation Quality Gate",
         "Reject incomplete artifacts before ARB. Reduce wasted committee time.",
         "NEW in v3"),
        ("Domain Auto-Assignment",
         "4 EA leaders route by domain. Auto-assign architect by bandwidth.",
         "MEDIUM"),
    ],
    [
        ["Document ARB/SDRB SLAs", "Assess AI tools", "Domain assignment rules"],
        ["AI pre-screening pilot", "Domain auto-assignment in OB-DMN-2"],
        ["ARB preparation quality gate", "SDRB fast-track for pre-approved"],
        ["ARB cycle time", "Preparation rejection rate", "AI pre-screening accuracy"],
    ]
)

# ── Topic 7: Compliance (slides 35-38) ──────────────────────────────────────────
add_section_slide("Compliance",
                  "Regulatory alignment, phase-boundary gates, and evidence management.")

add_bpmn_desc_slide("SP3: Compliance Review Stream", "SP3_Compliance",
    "Compliance review within SP3 parallel evaluation. Phase-boundary quality gates "
    "enforce compliance checks at every transition. OneTrust integration provides "
    "automated compliance audit trail. Regulatory annotations on every task: "
    "OCC 2023-17, SOX, GDPR, DORA, EU AI Act, NIST CSF 2.0.")

add_current_state_slide(
    "Compliance: Current State Findings",
    [
        ("End-of-Year Failures",
         "Compliance failures discovered late in year. Reactive not proactive.",
         False),
        ("OneTrust Evidence Platform",
         "Assessment evidence captured in OneTrust. Control gaps documented.",
         True),
        ("No Phase-Boundary Gates",
         "Compliance checks not enforced at transitions between phases.",
         False),
        ("Regulatory Complexity",
         "OCC 2023-17, SOX, GDPR, DORA, EU AI Act all apply. No unified tracking.",
         False),
    ],
    quote=None
)

add_recommendations_slide(
    "Compliance: Recommendations",
    [
        ("Consolidated Compliance Stream",
         "Single compliance review in SP3 parallel. Phase-boundary quality gates.",
         "HIGH IMPACT"),
        ("OneTrust Integration",
         "Automated compliance audit trail. Evidence linked to process instances.",
         "NEW in v3"),
        ("Regulatory Annotation",
         "Every task tagged with applicable regulations via Camunda properties.",
         "MEDIUM"),
        ("3-Tier Review Structure",
         "Fast path (single), Committee (quorum), Board (full). DMN-3 routes.",
         "MEDIUM"),
    ],
    [
        ["Map compliance touchpoints", "OneTrust evidence capabilities"],
        ["OneTrust compliance evidence mapping"],
        ["Automated compliance audit trail via OneTrust"],
        ["Regulatory traceability dashboard", "OneTrust data feed"],
    ]
)

# ── Topic 8: AI Governance (slides 39-42) ────────────────────────────────────────
add_section_slide("AI Governance",
                  "AI risk classification, committee consolidation, and regulatory alignment.")

add_bpmn_desc_slide("SP3: AI Review Stream", "SP3_AIReview",
    "AI governance review within SP3 parallel evaluation. Single consolidated stream "
    "replaces 3 separate committees. AI risk posture classification via DMN. "
    "Fast-track pathway for AI-only reviews targeting 2-week E2E. "
    "v3 merges 3+ AI questionnaires into single dataset.")

add_current_state_slide(
    "AI Governance: Current State Findings",
    [
        ("60+ Items in Queue",
         "AI governance backlog overwhelming. No clear prioritization.",
         False),
        ("Committee Redundancy",
         "Working Committee + Decision Committee somewhat redundant with TBC.",
         False),
        ("Questionnaire Proliferation",
         '3 additional AI forms "snuck up" beyond standard DD. Being merged.',
         False),
        ("3x Vendor Contact",
         "Tech risk, cybersecurity, TPRM independently contact vendors for AI.",
         False),
    ],
    quote="I have no idea why they're there... they just snuck up and I'm really annoyed",
    quote_source="TPRM Lead"
)

add_recommendations_slide(
    "AI Governance: Recommendations",
    [
        ("Consolidate 3 Committees to 1",
         "Single AI review stream. Eliminate redundant presentations.",
         "HIGH IMPACT"),
        ("AI Questionnaire Merger (GAP-23)",
         "Merge 3+ forms into single dataset. TPRM Lead already in progress.",
         "NEW in v3"),
        ("Single Vendor Contact",
         "One AI review stream. Eliminate 3x vendor contact redundancy.",
         "NEW in v3"),
        ("AI Fast-Track (GAP-6)",
         "AI Governance + Security only. 2-week E2E target.",
         "HIGH IMPACT"),
    ],
    [
        ["Merge 3 AI questionnaires", "Map committee consolidation", "Inventory vendor contacts"],
        ["Single AI questionnaire deployed", "Committee consolidation communicated"],
        ["AI fast-track piloted", "EU AI Act mapping to DMN"],
        ["AI review cycle time (2wk target)", "Vendor contact reduction (3x to 1x)"],
    ]
)

# ── Topic 9: Privacy (slides 43-46) ─────────────────────────────────────────────
add_section_slide("Privacy",
                  "Data classification, privacy-by-design, and cross-border requirements.")

add_bpmn_desc_slide("SP1/SP3: Privacy Assessment", "SP1_SP3_Privacy",
    "Data classification at intake (SP1) drives privacy routing. "
    "Privacy review embedded in SP3 parallel evaluation. "
    "OneTrust privacy module available for automated assessments. "
    "Cross-border data flow triggers additional assessment requirements.")

add_current_state_slide(
    "Privacy: Current State Findings",
    [
        ("Embedded Across Teams",
         "Privacy review integrated into multiple team processes.",
         True),
        ("Data Classification at Intake",
         "Classification questions appear early but not enforced.",
         False),
        ("Cross-Border Complexity",
         "International data flows require additional assessment.",
         False),
        ("OneTrust Privacy Module",
         "Privacy capabilities available in existing OneTrust deployment.",
         True),
    ],
    quote=None
)

add_recommendations_slide(
    "Privacy: Recommendations",
    [
        ("Data Classification at Intake",
         "Mandatory data classification in SP1 form. Drives privacy routing.",
         "MEDIUM"),
        ("Privacy-by-Design",
         "Privacy controls embedded in architecture review checklist.",
         "MEDIUM"),
        ("OneTrust Privacy Integration",
         "Leverage existing privacy module for automated assessments.",
         "NEW in v3"),
        ("Cross-Border Assessment",
         "Automated triggers for international data flow requirements.",
         "LOW"),
    ],
    [
        ["Map privacy touchpoints", "Classification field requirements"],
        ["Data classification mandatory at intake"],
        ["OneTrust privacy integration", "Privacy-by-design checklist"],
        ["Cross-border automation", "Privacy compliance rate"],
    ]
)

# ── Topic 10: Commercial Counsel (slides 47-50) ─────────────────────────────────
add_section_slide("Commercial Counsel",
                  "Contract lifecycle, negotiation bottleneck, and deviation tracking.")

add_bpmn_desc_slide("SP4: Contracting and Build", "SP4_ContractBuild",
    "Contract negotiation, vendor selection finalization, and build/buy execution. "
    "Buy path: refine solution, PoC, technical risk, negotiate, await contract, finalize. "
    "Build path: define build requirements, PDLC execution. "
    "v3 elevates contract automation to critical priority: 2 people handling 30+ contracts/month.")

add_current_state_slide(
    "Commercial Counsel: Current State Findings",
    [
        ("Dumpster Fire #1",
         "2 people negotiating 30+ contracts monthly. 4 years unsustainable.",
         False),
        ("No Deviation Reporting",
         "No reportable format for contract deviations. Unknown compliance for older contracts.",
         False),
        ("Security Exhibits: 1.5 Years",
         "Contract negotiation for security terms can take up to 18 months.",
         False),
        ("Sourcing Doesn't Source",
         "Sourcing dept manages contract lifecycle, not actual sourcing.",
         False),
    ],
    quote="That team desperately needs automation... The firm needs them to be automated",
    quote_source="TPRM Lead"
)

add_recommendations_slide(
    "Commercial Counsel: Recommendations",
    [
        ("Contract Review Automation (GAP-20)",
         "Automated review for standard terms. Critical priority.",
         "HIGH IMPACT / NEW"),
        ("OneTrust Deviation Tracking",
         "Control gaps documented, risk acceptance tracked in OneTrust.",
         "NEW in v3"),
        ("Template Standardization",
         "Standard contract templates reduce negotiation cycles.",
         "MEDIUM"),
        ("Legacy Compliance Audit",
         "Compliance status for older contracts retroactively assessed.",
         "NEW in v3"),
    ],
    [
        ["Quantify contract backlog", "Deviation tracking requirements", "Automation criteria"],
        ["Contract template standardization", "OneTrust deviation tracking pilot"],
        ["Automated contract review for standard terms", "Deviation reporting"],
        ["Contract cycle time", "Deviation tracking coverage", "Manual review reduction"],
    ]
)

# ── Topic 11: TPRM (slides 51-54) ──────────────────────────────────────────────
add_section_slide("TPRM",
                  "Third-party risk management, vendor due diligence, and OneTrust integration.")

add_bpmn_desc_slide("SP3: Vendor Due Diligence", "SP3_VendorDD",
    "830-question vendor questionnaire with skip logic. Vendor completion avg 30 days. "
    "Internal review 75 days (down from 144 in 2019). RAE: 80 questions, 28-29 day actual "
    "vs 14-day target, 335 assessments/year with 8-person team. "
    "v3 adds OneTrust integration for assessment automation and shift-left strategy.")

add_current_state_slide(
    "TPRM: Current State Findings",
    [
        ("System Landscape Fragmented",
         "ServiceNow, Ariba, OneTrust, Oracle. Manual PDF transfers between systems.",
         False),
        ("RAE 2x Target",
         "80 questions, 28-29 day actual vs 14-day target. 335 assessments/year.",
         False),
        ("DD: 830 Questions, 75 Days",
         "Vendor completion 30 days. Internal review 75 days (down from 144).",
         False),
        ("Process Ownership Vacuum",
         "Half a person assigned. No clear accountability or communication.",
         False),
    ],
    quote="Somebody needs to be empowered at the firm to say I own this",
    quote_source="TPRM Lead"
)

add_recommendations_slide(
    "TPRM: Recommendations",
    [
        ("Empowered Process Owner (GAP-18)",
         "Dedicated owner with authority. 1 workflow manager + 2-3 PMs.",
         "HIGH IMPACT"),
        ("Shift-Left Strategy",
         "Mini-RFP tools, standard questions upfront, self-service for business users.",
         "NEW in v3"),
        ("OneTrust Integration (GAP-21)",
         "Zeebe service tasks for assessment automation. TPRM module for vendor DD.",
         "NEW in v3"),
        ("Competitor Benchmarking",
         "Target 60-90 day E2E matching less-mature competitors.",
         "MEDIUM"),
    ],
    [
        ["Map OneTrust API integration", "Shift-left question inventory", "Baseline metrics"],
        ["OneTrust-Camunda integration design", "Shift-left mini-RFP pilot"],
        ["OneTrust assessment automation in SP3", "DD skip logic optimized"],
        ["RAE cycle time (14d target)", "DD cycle time", "OneTrust automation rate"],
    ]
)


# ╔═══════════════════════════════════════════════════════════════════════════════╗
# ║  CLOSING SLIDES (55-58)                                                     ║
# ╚═══════════════════════════════════════════════════════════════════════════════╝

# ── Slide 55: OneTrust Integration Architecture ──────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, KPMG_WHITE)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
             "OneTrust Integration Architecture", font_size=24, bold=True,
             color=KPMG_BLUE, font_name='Open Sans Condensed')

# Integration flow
add_text_box(slide, Inches(0.6), Inches(1.1), Inches(12), Inches(0.3),
             "Integration Flow (Camunda 8 / Zeebe)", font_size=14, bold=True,
             color=KPMG_DARK_NAVY, font_name='Open Sans Condensed')

flow_steps = [
    ("1. Create Assessment", "Zeebe Service Task triggers OneTrust Assessment Automation API"),
    ("2. Assessor Completes", "Risk questionnaire (RAE) completed in OneTrust by assigned assessor"),
    ("3. Retrieve Results", "Zeebe Service Task retrieves assessment output via API"),
    ("4. DMN Routing", "Process variables (riskScore, riskTier, complianceFindings[]) feed DMN routing"),
]
for i, (fs, fd) in enumerate(flow_steps):
    add_card(slide, Inches(0.6 + i * 3.1), Inches(1.5), Inches(2.9), Inches(1.1), fs, fd, KPMG_LIGHT_BLUE)

# Integration points
add_text_box(slide, Inches(0.6), Inches(2.9), Inches(12), Inches(0.3),
             "Two Integration Points", font_size=14, bold=True,
             color=KPMG_DARK_NAVY, font_name='Open Sans Condensed')
int_points = [
    ("SP3: Risk Assessment", "OneTrust Assessment Automation manages questionnaires. Results pulled via API, feed DMN routing."),
    ("SP3: Vendor Due Diligence", "OneTrust TPRM module. 830-question vendor questionnaire with skip logic."),
    ("SP4: Contract Deviations", "OneTrust control gap documentation. Deviation tracking, risk acceptance."),
]
for i, (ip, desc) in enumerate(int_points):
    add_card(slide, Inches(0.6 + i * 4.1), Inches(3.3), Inches(3.9), Inches(1.1), ip, desc, KPMG_BLUE)

# Variable mapping
add_text_box(slide, Inches(0.6), Inches(4.7), Inches(12), Inches(0.3),
             "Process Variable Mapping", font_size=14, bold=True,
             color=KPMG_DARK_NAVY, font_name='Open Sans Condensed')
variables = [
    ("riskScore", "Numeric risk score from OneTrust assessment"),
    ("riskTier", "High / Limited / Minimal -- drives DD depth"),
    ("complianceFindings[]", "Array of regulatory gaps requiring remediation"),
    ("controlGapCount", "Count of deviations for contract review routing"),
]
for i, (vn, vd) in enumerate(variables):
    left = Inches(0.6 + (i % 2) * 6.2)
    top = Inches(5.1 + (i // 2) * 0.55)
    add_card(slide, left, top, Inches(5.9), Inches(0.45), vn, vd, KPMG_MEDIUM_BLUE)

# Regulatory alignment
add_text_box(slide, Inches(0.6), Inches(6.3), Inches(12), Inches(0.3),
             "Regulatory Alignment: OCC 2023-17 | NIST CSF 2.0 | DORA | SOX",
             font_size=9, bold=True, color=RGBColor(0x66, 0x66, 0x66))
add_text_box(slide, Inches(0.6), Inches(6.55), Inches(12), Inches(0.25),
             "Auth: OAuth2 (client_id/secret), tenant-specific hostname (https://{tenant}.my.onetrust.com)",
             font_size=8, color=RGBColor(0x99, 0x99, 0x99))
add_footer(slide)

# ── Slide 56: Staffing & Resource Model ─────────────────────────────────────────
add_table_slide(
    "Staffing and Resource Model",
    ["Function", "Current Staffing", "Workload", "Gap"],
    [
        ["Risk/DD Team", "8 people", "335 assessments/year", "At capacity -- doubled output while cutting timeline"],
        ["Legal/Contracts", "2 people", "30+ contracts/month", 'Critical -- "dumpster fire"'],
        ["Architecture", "2-3 people", "Recently reduced", "Funding-constrained by domains"],
        ["Security Architecture", "~1 person for AI", '"Need three of me"', "Cannot scale for AI review volume"],
        ["Vendor Management", "6 total, 2 at 50%", "Full process facilitation", "No formal allocation"],
        ["Technology Vendor Mgmt", "1 person part-time", "Full START process", "Inadequate for scope"],
    ],
    subtitle='"If we want this to really click... I can\'t have the architect review group being a critical portion with only two people" -- TPRM Lead'
)

# ── Slide 57: Measurement Dashboard ─────────────────────────────────────────────
add_table_slide(
    "Measurement Dashboard: Baseline vs. Day 120 Target",
    ["Metric", "Baseline (Current)", "Day 120 Target"],
    [
        ["E2E cycle time (standard)", "6-9 months", "60-90 days"],
        ["E2E cycle time (Enable)", "Same as standard", "30 days"],
        ["E2E cycle time (AI fast-track)", "Same as standard", "14 days"],
        ["RAE completion", "28-29 days", "14 days"],
        ["DD internal review", "75 days", "30 days"],
        ["Security review (Baseline tier)", "2 weeks", "Same-day (automated)"],
        ["Form completion rate", "Unknown", "90%+ first-pass"],
        ["Contract cycle time", "Varies (up to 1.5yr)", "90 days standard"],
        ["Business Council decision", "Monthly + email", "48-hour async SLA"],
        ["Queue transparency", "None", "Real-time dashboard"],
    ]
)

# ── Slide 58: Next Steps ────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, KPMG_MEDIUM_BLUE, KPMG_BLUE)
add_text_box(slide, Inches(1), Inches(0.8), Inches(10), Inches(0.8),
             "Next Steps", font_size=40, bold=True, color=KPMG_WHITE, font_name='Open Sans Condensed')
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.65), Inches(3), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = KPMG_GOLD
line.line.fill.background()

next_steps = [
    ("1. Executive Alignment",
     "Present discovery findings and v3 roadmap to executive sponsors. "
     "Resolve NDA timing, committee consolidation, and process ownership decisions.",
     KPMG_GOLD),
    ("2. Standards Definition",
     "Finalize 3 request types, 6 security categories, prioritization scoring formula, "
     "and 3-pathway routing rules. Name owners for every Day 60 deliverable.",
     KPMG_LIGHT_BLUE),
    ("3. Quick Wins (Day 60)",
     "Deploy unified intake form, completeness gate, deal-killer pre-screen, "
     "NDA gate, and status notifications. Pilot simultaneous engagement.",
     KPMG_EMERALD),
    ("4. Full Operation (Day 120)",
     "5 parallel streams, OneTrust integration, contract automation, "
     "distributed pod model, and executive KPI reporting with before/after metrics.",
     KPMG_WHITE),
]
for i, (ns_title, ns_body, accent) in enumerate(next_steps):
    col = i % 2
    row = i // 2
    left = Inches(1 + col * 5.8)
    top = Inches(2.2 + row * 2.3)
    # Card background
    card_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         left, top, Inches(5.4), Inches(2.0))
    card_shape.fill.solid()
    card_shape.fill.fore_color.rgb = RGBColor(0x00, 0x2A, 0x70)
    card_shape.line.fill.background()
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), Inches(2.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), Inches(5.0), Inches(0.4),
                 ns_title, font_size=14, bold=True, color=KPMG_WHITE, font_name='Open Sans Condensed')
    # Body
    add_text_box(slide, left + Inches(0.2), top + Inches(0.6), Inches(5.0), Inches(1.2),
                 ns_body, font_size=10, color=RGBColor(0xBB, 0xCC, 0xEE))

add_footer(slide, dark=True)


# ╔═══════════════════════════════════════════════════════════════════════════════╗
# ║  SAVE                                                                       ║
# ╚═══════════════════════════════════════════════════════════════════════════════╝

import os
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "v3-roadmap.pptx")
prs.save(output_path)
slide_count = len(prs.slides)
print(f"Saved {slide_count} slides to {output_path}")
