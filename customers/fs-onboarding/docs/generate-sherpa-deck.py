#!/usr/bin/env python3
"""Generate the Onboarding Sherpa Toolkit PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
NAVY = RGBColor(0x0F, 0x17, 0x2A)
DARK_BG = RGBColor(0x11, 0x18, 0x27)
SURFACE = RGBColor(0x1E, 0x29, 0x3B)
BORDER = RGBColor(0x33, 0x41, 0x55)
TEXT = RGBColor(0xE2, 0xE8, 0xF0)
TEXT2 = RGBColor(0x94, 0xA3, 0xB8)
TEXT3 = RGBColor(0x64, 0x74, 0x8B)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
WARN = RGBColor(0xF5, 0x9E, 0x0B)
DANGER = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREEN = RGBColor(0x4A, 0xDE, 0x80)
GOLD = RGBColor(0xFB, 0xBF, 0x24)

# Faint backgrounds
FAINT_GREEN = RGBColor(0x0D, 0x2A, 0x1A)
FAINT_RED = RGBColor(0x2A, 0x0D, 0x0D)
FAINT_BLUE = RGBColor(0x0D, 0x1A, 0x2A)
FAINT_YELLOW = RGBColor(0x2A, 0x24, 0x0D)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=12,
                 color=TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_shape_box(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_bullet_list(slide, left, top, width, height, items, font_size=11,
                    color=TEXT, bullet_color=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(4)
        p.level = 0
    return txBox


def add_table(slide, left, top, width, rows, cols, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                          Inches(width), Inches(0.3 * rows))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    return table


def style_table(table, header_bg=SURFACE, header_fg=TEXT2, cell_bg=DARK_BG, cell_fg=TEXT):
    for row_idx, row in enumerate(table.rows):
        for cell in row.cells:
            cell.fill.solid()
            if row_idx == 0:
                cell.fill.fore_color.rgb = header_bg
            else:
                cell.fill.fore_color.rgb = cell_bg
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.name = 'Calibri'
                if row_idx == 0:
                    paragraph.font.color.rgb = header_fg
                    paragraph.font.bold = True
                else:
                    paragraph.font.color.rgb = cell_fg
            cell.margin_left = Emu(72000)
            cell.margin_right = Emu(72000)
            cell.margin_top = Emu(36000)
            cell.margin_bottom = Emu(36000)


def add_section_badge(slide, text, top=0.4):
    shape = add_shape_box(slide, 0.5, top, 1.6, 0.28, SURFACE, BORDER)
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(8)
    shape.text_frame.paragraphs[0].font.color.rgb = ACCENT
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = 'Calibri'
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_icon_circle(slide, left, top, size, color, text):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(int(size * 18))
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = 'Calibri'
    shape.text_frame.word_wrap = False


def add_rich_text_box(slide, left, top, width, height, runs, alignment=PP_ALIGN.LEFT):
    """Add text box with mixed formatting. runs = [(text, font_size, color, bold), ...]"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    for i, (text, font_size, color, bold) in enumerate(runs):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
            run.text = text
        else:
            run = p.add_run()
            run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = 'Calibri'
    return txBox


def add_numbered_item(slide, num, x, y, title, body, title_color=WHITE, num_color=ACCENT,
                      body_color=TEXT2, card_width=5.5, card_height=1.6):
    """Add a numbered recommendation card with icon, title, and body text."""
    box = add_shape_box(slide, x, y, card_width, card_height, SURFACE, BORDER)
    add_icon_circle(slide, x + 0.15, y + 0.15, 0.4, num_color, str(num))
    add_text_box(slide, x + 0.7, y + 0.15, card_width - 0.9, 0.28, title,
                 font_size=14, color=title_color, bold=True)
    add_text_box(slide, x + 0.7, y + 0.45, card_width - 0.9, card_height - 0.55, body,
                 font_size=10, color=body_color)
    return box


# ============================================================
# BUILD PRESENTATION
# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ================================================================
# 30-DAY PLAN SLIDES (inserted before toolkit)
# ================================================================

# ---- PLAN SLIDE 1: Title ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_shape_box(slide, 0, 0, 13.333, 0.06, ACCENT)

add_text_box(slide, 1.5, 1.8, 10.3, 0.8, 'Software Onboarding',
             font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 2.6, 10.3, 0.7, '30-Day Process Improvement Plan',
             font_size=36, color=LIGHT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 2.5, 3.5, 8.3, 0.5, 'Zero technology changes  \u2014  7 organizational decisions that cut cycle time in half',
             font_size=16, color=TEXT3, alignment=PP_ALIGN.CENTER)

# Three stat boxes
stats = [
    ('6-9 months', 'Current Avg\nCycle Time', DANGER),
    ('3-4 months', 'Target After\n30 Days', GREEN),
    ('7 changes', 'Process Changes\nNo Technology', ACCENT),
]
for i, (val, label, color) in enumerate(stats):
    x = 2.5 + i * 3.0
    box = add_shape_box(slide, x, 4.5, 2.5, 1.5, SURFACE, BORDER)
    add_text_box(slide, x, 4.6, 2.5, 0.5, val,
                 font_size=28, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, 5.2, 2.5, 0.6, label,
                 font_size=11, color=TEXT3, alignment=PP_ALIGN.CENTER)

add_shape_box(slide, 0, 7.2, 13.333, 0.3, SURFACE)
add_text_box(slide, 0.5, 7.22, 12.3, 0.25, 'Financial Services  \u2022  Software Onboarding Lifecycle',
             font_size=10, color=TEXT3, alignment=PP_ALIGN.CENTER)

# ---- PLAN SLIDE 2: The Problem ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_shape_box(slide, 0, 0, 13.333, 0.06, DANGER)

add_text_box(slide, 0.5, 0.3, 10, 0.5, 'The Problem: Why Onboarding Takes 6-9 Months', font_size=28, color=WHITE, bold=True)
add_text_box(slide, 0.5, 0.85, 10, 0.3,
             'Three root causes account for 80% of elapsed time. None require technology to fix.',
             font_size=13, color=TEXT3)

# Root cause 1
box = add_shape_box(slide, 0.5, 1.4, 4.0, 3.0, FAINT_RED, RGBColor(0x50, 0x20, 0x20))
add_text_box(slide, 0.7, 1.45, 3.6, 0.3, '18 Sequential Committees', font_size=16, color=RGBColor(0xF8, 0x71, 0x71), bold=True)
add_text_box(slide, 0.7, 1.8, 3.6, 2.5,
             'Every request passes through up to 18 review committees one at a time. '
             'Security waits for Legal. Legal waits for Compliance. Compliance waits for Architecture.\n\n'
             'A 5-day review becomes 5 x 18 = 90 days of elapsed time, even though each reviewer only '
             'spends a few hours.\n\n'
             'The reviews themselves aren\'t slow \u2014 the sequencing is.',
             font_size=10, color=TEXT2)

# Root cause 2
box = add_shape_box(slide, 4.7, 1.4, 4.0, 3.0, FAINT_RED, RGBColor(0x50, 0x20, 0x20))
add_text_box(slide, 4.9, 1.45, 3.6, 0.3, 'No Single Point of Accountability', font_size=16, color=RGBColor(0xF8, 0x71, 0x71), bold=True)
add_text_box(slide, 4.9, 1.8, 3.6, 2.5,
             'Nobody owns elapsed time. Each reviewer owns their piece, '
             'but nobody tracks the overall request from intake to go-live.\n\n'
             'When a request stalls, no one notices for weeks. The requester calls around '
             'asking "where is my request?" and gets transferred between teams.\n\n'
             'The Sherpa role solves this \u2014 one person owns the clock.',
             font_size=10, color=TEXT2)

# Root cause 3
box = add_shape_box(slide, 8.9, 1.4, 4.0, 3.0, FAINT_RED, RGBColor(0x50, 0x20, 0x20))
add_text_box(slide, 9.1, 1.45, 3.6, 0.3, 'Everything Gets Full Treatment', font_size=16, color=RGBColor(0xF8, 0x71, 0x71), bold=True)
add_text_box(slide, 9.1, 1.8, 3.6, 2.5,
             'A $10K SaaS renewal with an existing vendor gets the same 9-reviewer '
             'due diligence as a $2M AI platform from a brand new vendor.\n\n'
             'Roughly 40-50% of requests could be fast-tracked with minimal review, '
             'freeing reviewer capacity for the requests that actually need scrutiny.\n\n'
             'Tiered routing fixes this on day one.',
             font_size=10, color=TEXT2)

# Bottom: timeline
add_shape_box(slide, 0.5, 4.7, 12.3, 0.06, BORDER)
add_text_box(slide, 0.5, 4.9, 12.3, 0.3, 'The 30-Day Plan: Three Phases of Change', font_size=16, color=WHITE, bold=True)

phases_30 = [
    ('Week 1-2', 'Stop the Bleeding', 'Assign Sherpas, run parallel reviews, pre-screen intake', DANGER, '#1-3'),
    ('Week 2-3', 'Create Rhythm', 'Tier routing, weekly governance stand-up', WARN, '#4-5'),
    ('Week 3-4', 'Measure & Adjust', 'Visible tracking board, enforced SLA deadlines', GREEN, '#6-7'),
]
for i, (week, title, desc, color, nums) in enumerate(phases_30):
    x = 0.5 + i * 4.2
    box = add_shape_box(slide, x, 5.3, 3.9, 1.6, SURFACE, BORDER)
    add_text_box(slide, x + 0.15, 5.35, 1.5, 0.22, week, font_size=10, color=color, bold=True)
    add_text_box(slide, x + 1.8, 5.35, 1.9, 0.22, nums, font_size=10, color=TEXT3, bold=False, alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, x + 0.15, 5.6, 3.6, 0.3, title, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, x + 0.15, 5.95, 3.6, 0.8, desc, font_size=11, color=TEXT2)

# ---- PLAN SLIDE 3: Week 1-2 - Stop the Bleeding ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_shape_box(slide, 0, 0, 13.333, 0.06, DANGER)
add_section_badge(slide, 'WEEK 1-2')

add_text_box(slide, 0.5, 0.75, 10, 0.5, 'Stop the Bleeding', font_size=28, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.25, 10, 0.3,
             'Three changes in the first two weeks that immediately impact cycle time.',
             font_size=13, color=TEXT3)

# Item 1: Appoint Sherpas
add_numbered_item(slide, 1, 0.5, 1.8, 'Appoint Onboarding Sherpas (Day 1)', (
    'Assign 2-3 people as dedicated facilitators who own each request end-to-end. '
    'The Sherpa tracks where every active request sits, chases reviewers who haven\'t responded, '
    'escalates stalls after 48 hours, and reports weekly to leadership on cycle time.\n\n'
    'Impact: Immediately surfaces where requests are stuck. Visibility alone cuts delays.'
), num_color=DANGER, card_width=6.0, card_height=2.2)

# Item 2: Kill Sequential Reviews
add_numbered_item(slide, 2, 6.8, 1.8, 'Kill Sequential Reviews \u2014 Run Parallel (Day 3)', (
    'When a request enters Due Diligence, the Sherpa sends the package to ALL reviewers simultaneously '
    'with a shared deadline (e.g., 5 business days). No reviewer waits for another to finish first.\n\n'
    'Impact: Compresses the longest phase from 8-12 weeks to 2-3 weeks. '
    'This single change is worth 2-3 months of cycle time.'
), num_color=DANGER, card_width=6.0, card_height=2.2)

# Item 3: Paper Pre-Screen
add_numbered_item(slide, 3, 0.5, 4.3, 'Implement a Paper Pre-Screen (Mini RFP on Paper)', (
    'Before a request enters the formal pipeline, require a 1-page intake form '
    'with 5 deal-killer questions:\n\n'
    '  \u2022  Is there budget approval?\n'
    '  \u2022  Is there a business sponsor (VP+ level)?\n'
    '  \u2022  Does it touch regulated data (PII/PHI/financial)?\n'
    '  \u2022  Is there an existing contract/vendor relationship?\n'
    '  \u2022  What\'s the target go-live?\n\n'
    'If the first two are "No", the request doesn\'t enter the pipeline. '
    'Today, requests that will never be approved consume weeks of committee time.\n\n'
    'Impact: Filters out 20-30% of requests before they consume any reviewer bandwidth.'
), num_color=DANGER, card_width=12.3, card_height=2.8)

# ---- PLAN SLIDE 4: Week 2-3 - Create Rhythm ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_shape_box(slide, 0, 0, 13.333, 0.06, WARN)
add_section_badge(slide, 'WEEK 2-3')

add_text_box(slide, 0.5, 0.75, 10, 0.5, 'Create Rhythm', font_size=28, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.25, 10, 0.3,
             'Establish the operating cadence that makes the process predictable.',
             font_size=13, color=TEXT3)

# Item 4: Tiered Review
add_numbered_item(slide, 4, 0.5, 1.8, 'Tiered Review \u2014 Not Everything Gets the Full Treatment', (
    'Establish 3 risk tiers on paper (maps to DMN_RiskTierClassification):'
), num_color=WARN, card_width=12.3, card_height=0.9)

# Tier table
tier_tbl = add_table(slide, 0.5, 2.85, 12.3, 4, 5, [1.2, 2.0, 4.5, 2.0, 2.6])
tier_hdrs = ['Score', 'Tier', 'Reviewers Required', 'Target SLA', 'Escalation Path']
tier_rows = [
    ['6-8 pts', 'FAST-TRACK', 'Sherpa + 1 reviewer (Security or Compliance)', '5 business days', "Sherpa's manager after Day 3"],
    ['9-13 pts', 'STANDARD', '4 key reviewers in parallel (Security, Compliance, Legal, Tech Architecture)', '15 business days', 'Governance lead after Day 10'],
    ['14-18 pts', 'ENHANCED', 'All 9 reviewers in parallel + executive briefing', '25 business days', 'CISO/CRO after Day 15'],
]
for j, h in enumerate(tier_hdrs):
    tier_tbl.cell(0, j).text = h
for i, row in enumerate(tier_rows):
    for j, v in enumerate(row):
        tier_tbl.cell(i + 1, j).text = v
style_table(tier_tbl)
tier_colors = [LIGHT_GREEN, GOLD, RGBColor(0xF8, 0x71, 0x71)]
for i, c in enumerate(tier_colors):
    for p in tier_tbl.cell(i + 1, 1).text_frame.paragraphs:
        p.font.color.rgb = c
        p.font.bold = True

add_text_box(slide, 0.5, 4.3, 12.3, 0.3,
             'Today, a $10K SaaS renewal gets the same scrutiny as a $2M AI platform. '
             '40-50% of requests move to Fast-Track, freeing reviewer capacity for what matters.',
             font_size=11, color=WARN)

# Item 5: Weekly Stand-Up
add_numbered_item(slide, 5, 0.5, 4.8, 'Weekly Governance Stand-Up (30 minutes max)', (
    'Replace monthly committee meetings with a weekly 30-minute stand-up:'
), num_color=WARN, card_width=12.3, card_height=0.9)

# Stand-up agenda
su_tbl = add_table(slide, 0.5, 5.85, 12.3, 6, 4, [1.2, 2.5, 6.6, 2.0])
su_hdrs = ['Time Box', 'Agenda Item', 'Description', 'Owner']
su_rows = [
    ('2 min', 'Pipeline Summary', 'How many requests total? How many in each phase? How many breached SLA?', 'Sherpa'),
    ('10 min', 'SLA Breaches', 'Walk through EACH breached request. What is blocking? Who needs to act? Decision: escalate, extend, or kill?', 'Sherpa + Blockers'),
    ('8 min', 'At-Risk Requests', 'Requests approaching SLA deadline (within 2 business days). What\'s needed to close on time?', 'Sherpa + Reviewers'),
    ('5 min', 'New Intake', 'New requests this week. Quick tier assignment and reviewer identification.', 'Sherpa'),
    ('5 min', 'Decisions Needed', 'Any go/no-go decisions the committee needs to make right now?', 'Governance Lead'),
]
for j, h in enumerate(su_hdrs):
    su_tbl.cell(0, j).text = h
for i, (t, item, desc, owner) in enumerate(su_rows):
    su_tbl.cell(i + 1, 0).text = t
    su_tbl.cell(i + 1, 1).text = item
    su_tbl.cell(i + 1, 2).text = desc
    su_tbl.cell(i + 1, 3).text = owner
style_table(su_tbl)
su_time_colors = [ACCENT, RGBColor(0xF8, 0x71, 0x71), GOLD, LIGHT_GREEN, TEXT2]
for i, c in enumerate(su_time_colors):
    for p in su_tbl.cell(i + 1, 0).text_frame.paragraphs:
        p.font.color.rgb = c
        p.font.bold = True

# ---- PLAN SLIDE 5: Week 3-4 - Measure & Adjust ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_shape_box(slide, 0, 0, 13.333, 0.06, GREEN)
add_section_badge(slide, 'WEEK 3-4')

add_text_box(slide, 0.5, 0.75, 10, 0.5, 'Measure & Adjust', font_size=28, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.25, 10, 0.3,
             'Make progress visible and create accountability through deadlines.',
             font_size=13, color=TEXT3)

# Item 6: Visible Kanban Board
add_numbered_item(slide, 6, 0.5, 1.8, 'Publish a Visible Kanban Board', (
    'Put every active request on a physical board or simple spreadsheet with columns:\n\n'
    '    Intake  \u2192  Triage  \u2192  Due Diligence  \u2192  Governance Review  \u2192  Contracting  \u2192  Go-Live\n\n'
    'Color-code by age: GREEN (on track), YELLOW (approaching SLA), RED (breached).\n'
    'Update weekly. Post it where leadership walks by.\n\n'
    'Impact: Social pressure and visibility. When leadership can see that 6 requests have been '
    'stuck in "Due Diligence" for 40+ days, things move.'
), num_color=GREEN, card_width=12.3, card_height=2.6)

# Visual example of the kanban
phases_kanban = [
    ('Intake', GREEN, ['ONB-007\nSlack Enterprise\nDay 1']),
    ('Due Diligence', DANGER, ['ONB-003\nAI Analytics\nDay 18 (BREACH!)', 'ONB-005\nCRM Tool\nDay 8']),
    ('Gov. Review', WARN, ['ONB-002\nSnowflake\nDay 4']),
    ('Contracting', GREEN, ['ONB-001\nServiceNow\nDay 3']),
    ('Go-Live', GREEN, []),
]
for i, (phase, color, cards) in enumerate(phases_kanban):
    x = 0.5 + i * 2.5
    box = add_shape_box(slide, x, 4.6, 2.2, 2.6, SURFACE, BORDER)
    add_text_box(slide, x + 0.1, 4.65, 2.0, 0.22, phase, font_size=10, color=color, bold=True)
    for j, card_text in enumerate(cards):
        card_y = 4.95 + j * 0.75
        is_breach = 'BREACH' in card_text
        card_bg = FAINT_RED if is_breach else DARK_BG
        card_border = RGBColor(0x80, 0x30, 0x30) if is_breach else BORDER
        add_shape_box(slide, x + 0.1, card_y, 2.0, 0.65, card_bg, card_border)
        add_text_box(slide, x + 0.15, card_y + 0.05, 1.9, 0.55, card_text,
                     font_size=8, color=RGBColor(0xF8, 0x71, 0x71) if is_breach else TEXT2)
    if not cards:
        add_text_box(slide, x + 0.1, 5.3, 2.0, 0.4, '(empty)',
                     font_size=9, color=TEXT3, alignment=PP_ALIGN.CENTER)

# Item 7: SLA Deadlines
add_numbered_item(slide, 7, 0.5, 4.6, '', '', num_color=GREEN, card_width=0.01, card_height=0.01)
# Just the number circle, table carries the content

# ---- PLAN SLIDE 6: SLA Deadlines (full slide) ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_shape_box(slide, 0, 0, 13.333, 0.06, GREEN)
add_section_badge(slide, 'WEEK 3-4')

add_numbered_item(slide, 7, 0.5, 0.6, 'Set and Enforce SLA Deadlines (with Escalation)', (
    'Publish target timelines per phase. The escalation path is the enforcement mechanism \u2014 '
    'no technology needed, just the Sherpa sending an email: "This review is 3 days past SLA. Escalating per policy."'
), num_color=GREEN, card_width=12.3, card_height=1.0)

sla_tbl = add_table(slide, 0.5, 1.8, 12.3, 8, 5, [2.0, 1.5, 2.0, 3.4, 3.4])
sla_hdrs = ['Phase', 'SLA Target', 'Warning At', 'Level 1 Escalation', 'Level 2 Escalation']
sla_rows = [
    ('Intake / Triage', '2 business days', 'Day 1 (50%)', "Sherpa's manager (Day 2)", 'Governance Lead (Day 3)'),
    ('Planning', '5 business days', 'Day 3 (60%)', 'Governance Lead (Day 5)', 'CRO (Day 7)'),
    ('Due Diligence', '10 business days', 'Day 7 (70%)', 'CISO / CLO (Day 10)', 'CRO + CFO (Day 14)'),
    ('Individual Review', '5 business days', 'Day 3 (60%)', "Reviewer's manager (Day 5)", 'Domain head (Day 7)'),
    ('Governance Review', '5 business days', 'Day 3 (60%)', 'CRO (Day 5)', 'COO (Day 7)'),
    ('Contracting', '7 business days', 'Day 5 (70%)', 'CLO / CPO (Day 7)', 'CFO (Day 10)'),
    ('UAT / Go-Live', '5 business days', 'Day 3 (60%)', 'Project Sponsor (Day 5)', 'CTO (Day 7)'),
]
for j, h in enumerate(sla_hdrs):
    sla_tbl.cell(0, j).text = h
for i, row in enumerate(sla_rows):
    for j, v in enumerate(row):
        sla_tbl.cell(i + 1, j).text = v
style_table(sla_tbl)

# Impact statement
add_text_box(slide, 0.5, 5.3, 12.3, 0.3,
             'Impact: Creates accountability where none exists today. '
             'Decisions that currently wait 2-4 weeks for the next committee meeting happen within 5 days.',
             font_size=12, color=GREEN)

# ---- PLAN SLIDE 7: Expected 30-Day Results ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_shape_box(slide, 0, 0, 13.333, 0.06, ACCENT)

add_text_box(slide, 0.5, 0.3, 10, 0.5, 'Expected 30-Day Results', font_size=28, color=WHITE, bold=True)
add_text_box(slide, 0.5, 0.85, 10, 0.3,
             'Conservative estimates based on peer financial services implementations.',
             font_size=13, color=TEXT3)

# Results table
res_tbl = add_table(slide, 0.5, 1.3, 12.3, 6, 4, [3.0, 2.5, 2.5, 4.3])
res_hdrs = ['Metric', 'Current State', 'After 30 Days', 'Driver']
res_rows = [
    ('Avg Cycle Time', '6-9 months', '3-4 months (new requests)', 'Parallel reviews + tiered routing'),
    ('Requests in Limbo', 'Unknown', '100% visible', 'Tracking board + weekly stand-up'),
    ('Reviewer Bottleneck', 'Hidden', 'Identified by name/team', 'SLA escalation ladder'),
    ('Low-Value Requests Consuming Capacity', '~30%', 'Filtered at intake', 'Pre-screen deal-killers'),
    ('Decision Latency', '2-4 weeks (next committee)', '<5 days (weekly stand-up)', 'Weekly governance cadence'),
]
for j, h in enumerate(res_hdrs):
    res_tbl.cell(0, j).text = h
for i, (metric, current, after, driver) in enumerate(res_rows):
    res_tbl.cell(i + 1, 0).text = metric
    res_tbl.cell(i + 1, 1).text = current
    res_tbl.cell(i + 1, 2).text = after
    res_tbl.cell(i + 1, 3).text = driver
style_table(res_tbl)
# Red for current state, green for after
for i in range(1, 6):
    for p in res_tbl.cell(i, 1).text_frame.paragraphs:
        p.font.color.rgb = RGBColor(0xF8, 0x71, 0x71)
    for p in res_tbl.cell(i, 2).text_frame.paragraphs:
        p.font.color.rgb = LIGHT_GREEN
        p.font.bold = True

# Big three wins
add_text_box(slide, 0.5, 4.5, 12.3, 0.3, 'The Three Biggest Wins', font_size=18, color=WHITE, bold=True)

wins = [
    ('#2', 'Parallel Reviews', DANGER,
     'Compresses Due Diligence from 8-12 weeks to 2-3 weeks. '
     'This single change recovers 2-3 months of cycle time per request. '
     'The reviews themselves aren\'t slow \u2014 the sequential scheduling is.'),
    ('#4', 'Tiered Routing', WARN,
     'Moves 40-50% of requests to Fast-Track (5-day SLA vs 25-day). '
     'Frees reviewer capacity so Enhanced-tier requests get faster attention too. '
     'Net effect: everyone moves faster.'),
    ('#1', 'Sherpa Role', GREEN,
     'One person owns elapsed time per request. '
     'Chases reviewers, escalates stalls, reports weekly. '
     'Turns "nobody\'s watching" into "someone is always watching."'),
]
for i, (num, title, color, desc) in enumerate(wins):
    x = 0.5 + i * 4.2
    box = add_shape_box(slide, x, 4.95, 3.9, 2.2, SURFACE, BORDER)
    add_text_box(slide, x + 0.15, 5.0, 0.6, 0.25, num, font_size=12, color=color, bold=True)
    add_text_box(slide, x + 0.55, 5.0, 3.2, 0.25, title, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, x + 0.15, 5.3, 3.6, 1.8, desc, font_size=10, color=TEXT2)

# ---- PLAN SLIDE 8: Transition to Toolkit ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_shape_box(slide, 0, 0, 13.333, 0.06, GREEN)

add_text_box(slide, 1.5, 2.0, 10.3, 0.6, 'The Sherpa Toolkit',
             font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 2.0, 2.7, 9.3, 0.8,
             '8 ready-to-use forms, templates, and checklists that give the Sherpa\n'
             'a mechanistic, repeatable process from intake to close-out.',
             font_size=16, color=TEXT2, alignment=PP_ALIGN.CENTER)

# Eight form icons in a row
form_names = [
    ('1', 'Intake\nPre-Screen', GREEN),
    ('2', 'Risk Tier\nClassification', ACCENT),
    ('3', 'Request\nTracker', ACCENT),
    ('4', 'Parallel Review\nKickoff', WARN),
    ('5', 'Weekly\nStand-Up', WARN),
    ('6', 'SLA\nEscalation', DANGER),
    ('7', 'Phase\nChecklists', GREEN),
    ('8', 'Close-Out\nForm', ACCENT),
]
for i, (num, name, color) in enumerate(form_names):
    x = 0.8 + i * 1.55
    add_icon_circle(slide, x + 0.3, 3.8, 0.55, color, num)
    add_text_box(slide, x, 4.5, 1.15, 0.6, name,
                 font_size=9, color=TEXT2, alignment=PP_ALIGN.CENTER)

add_shape_box(slide, 0, 7.2, 13.333, 0.3, SURFACE)
add_text_box(slide, 0.5, 7.22, 12.3, 0.25,
             'These forms are the manual version of what we\'ll automate in Camunda. '
             'The process improvement comes first \u2014 the technology enforces it later.',
             font_size=10, color=TEXT3, alignment=PP_ALIGN.CENTER)

# ================================================================
# TOOLKIT SLIDES (existing)
# ================================================================

# ---- SLIDE: Toolkit Title ----
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, NAVY)

# Green accent bar at top
add_shape_box(slide, 0, 0, 13.333, 0.06, GREEN)

# Logo circle
add_icon_circle(slide, 5.9, 1.5, 1.0, GREEN, 'S')

add_text_box(slide, 1.5, 2.8, 10.3, 0.8, 'Onboarding Sherpa Toolkit',
             font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 2.5, 3.6, 8.3, 0.5, 'Manual Process Forms, Templates & Checklists',
             font_size=20, color=LIGHT_GREEN, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 2.5, 4.3, 8.3, 0.5, 'Software Onboarding Lifecycle  |  Financial Services',
             font_size=14, color=TEXT3, alignment=PP_ALIGN.CENTER)

# Bottom bar
add_shape_box(slide, 0, 7.2, 13.333, 0.3, SURFACE)
add_text_box(slide, 0.5, 7.2, 12.3, 0.3, 'Zero technology required  \u2022  Implementable in 30 days  \u2022  8 ready-to-use forms',
             font_size=10, color=TEXT3, alignment=PP_ALIGN.CENTER)

# ---- SLIDE 2: Table of Contents ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_shape_box(slide, 0, 0, 13.333, 0.06, GREEN)

add_text_box(slide, 0.5, 0.3, 6, 0.5, 'Toolkit Contents', font_size=28, color=WHITE, bold=True)
add_text_box(slide, 0.5, 0.85, 8, 0.3, '8 forms and templates that give the Sherpa a mechanistic, repeatable process',
             font_size=12, color=TEXT3)

toc_items = [
    ('1', 'Intake Pre-Screen Form', 'Filter requests before they enter the pipeline', GREEN),
    ('2', 'Risk Tier Classification', 'Score and route to Fast-Track / Standard / Enhanced', ACCENT),
    ('3', 'Request Tracking Board', 'Kanban + register for all in-flight requests', ACCENT),
    ('4', 'Parallel Review Kickoff', '9-reviewer assignment matrix and kickoff email', WARN),
    ('5', 'Weekly Governance Stand-Up', '30-minute agenda, notes template, summary email', WARN),
    ('6', 'SLA Escalation Playbook', 'Escalation ladder with 3-level email templates', DANGER),
    ('7', 'Phase Transition Checklists', 'Gate checklists for every phase boundary', GREEN),
    ('8', 'Request Close-Out Form', 'SLA performance capture and lessons learned', ACCENT),
]

for i, (num, title, desc, color) in enumerate(toc_items):
    y = 1.5 + i * 0.7
    add_icon_circle(slide, 1.0, y + 0.05, 0.45, color, num)
    add_text_box(slide, 1.7, y, 4, 0.3, title, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, 1.7, y + 0.3, 5, 0.25, desc, font_size=11, color=TEXT2)

# Right side: impact summary
box = add_shape_box(slide, 7.5, 1.5, 5.3, 5.0, SURFACE, BORDER)
add_text_box(slide, 7.8, 1.6, 4.7, 0.35, 'Expected 30-Day Impact', font_size=16, color=LIGHT_GREEN, bold=True)

impact_data = [
    ('Avg Cycle Time', '6-9 months', '3-4 months', 'Parallel reviews + tiered routing'),
    ('Requests in Limbo', 'Unknown', '100% visible', 'Tracking board + weekly stand-up'),
    ('Reviewer Bottleneck', 'Hidden', 'Named', 'SLA escalation ladder'),
    ('Low-Value Filtering', '~0%', '20-30%', 'Pre-screen deal-killers'),
    ('Decision Latency', '2-4 weeks', '<5 days', 'Weekly governance stand-up'),
]

table = add_table(slide, 7.8, 2.1, 4.8, 6, 4, [1.2, 1.0, 1.0, 1.6])
table.cell(0, 0).text = 'Metric'
table.cell(0, 1).text = 'Current'
table.cell(0, 2).text = 'After 30d'
table.cell(0, 3).text = 'How'
for i, (metric, current, after, how) in enumerate(impact_data):
    table.cell(i + 1, 0).text = metric
    table.cell(i + 1, 1).text = current
    table.cell(i + 1, 2).text = after
    table.cell(i + 1, 3).text = how
style_table(table)
# Color the "After 30d" column green
for i in range(1, 6):
    for p in table.cell(i, 2).text_frame.paragraphs:
        p.font.color.rgb = LIGHT_GREEN

# ---- SLIDE 3: Intake Pre-Screen (Part 1) ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_shape_box(slide, 0, 0, 13.333, 0.06, GREEN)
add_section_badge(slide, 'FORM 1 OF 8')

add_text_box(slide, 0.5, 0.75, 8, 0.5, 'Intake Pre-Screen Form', font_size=28, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.25, 8, 0.3, 'Complete BEFORE entering a request into the formal pipeline. ~10 minutes.',
             font_size=12, color=TEXT3)

# Left column: Requester info
box = add_shape_box(slide, 0.5, 1.8, 6.0, 2.2, SURFACE, BORDER)
add_text_box(slide, 0.7, 1.85, 3, 0.3, 'Requester Information', font_size=14, color=WHITE, bold=True)
fields1 = [
    'Requester Name *                    Department / BU',
    'Email                               Date Submitted',
    '',
    'Software / Vendor Name *',
    'Request Type *    (New Vendor | Existing Vendor | Renewal | Upgrade | Build | Forced Update)',
    'Target Go-Live Date',
    'Business Problem / Need *',
    'Existing Solutions Considered',
]
add_bullet_list(slide, 0.7, 2.2, 5.6, 1.7, fields1, font_size=10, color=TEXT2)

# Right column: Deal killers
box = add_shape_box(slide, 6.8, 1.8, 6.0, 2.2, FAINT_RED, RGBColor(0x60, 0x20, 0x20))
add_text_box(slide, 7.0, 1.85, 5, 0.3, 'Deal-Killer Screen', font_size=14, color=RGBColor(0xF8, 0x71, 0x71), bold=True)
add_text_box(slide, 7.0, 2.15, 5.5, 0.2, 'If ANY is "No", request cannot proceed', font_size=10, color=TEXT3)
killers = [
    '\u2610  Budget Approval: Has a budget been identified or pre-approved?',
    '\u2610  Business Sponsor: Is there an executive sponsor (VP+ level)?',
    '\u2610  Not a Duplicate: Confirmed not already in-flight or in portfolio?',
    '\u2610  Organizational Alignment: Aligns with strategic initiative?',
    '\u2610  Realistic Timeline: Go-live at least 30 days out (60 for new)?',
]
add_bullet_list(slide, 7.0, 2.45, 5.5, 1.5, killers, font_size=10, color=RGBColor(0xF8, 0x71, 0x71))

# Bottom: Data classification
box = add_shape_box(slide, 0.5, 4.2, 6.0, 2.5, SURFACE, BORDER)
add_text_box(slide, 0.7, 4.25, 5, 0.3, 'Data Classification (Quick Screen)', font_size=14, color=WHITE, bold=True)
add_text_box(slide, 0.7, 4.55, 5.5, 0.2, 'Check all that apply \u2014 determines risk tier', font_size=10, color=TEXT3)
data_items = [
    '\u2610  PII (Personally Identifiable Information)',
    '\u2610  PHI (Protected Health Information)',
    '\u2610  Financial transaction data (SOX-relevant)',
    '\u2610  AI / machine learning models',
    '\u2610  Critical infrastructure (core banking, payments)',
    '\u2610  Cloud-hosted by third party',
    '\u2610  Cross-border data transfer',
    '\u2610  None of the above',
]
add_bullet_list(slide, 0.7, 4.8, 5.5, 2.0, data_items, font_size=10, color=TEXT2)

# Sherpa decision
box = add_shape_box(slide, 6.8, 4.2, 6.0, 2.5, SURFACE, BORDER)
add_text_box(slide, 7.0, 4.25, 5, 0.3, 'Sherpa Decision', font_size=14, color=WHITE, bold=True)
decisions = [
    'PROCEED \u2014 Enter into pipeline',
    'RETURN \u2014 Missing information (specify)',
    'REJECT \u2014 Deal-killer not met',
    'REDIRECT \u2014 Existing solution available',
]
add_text_box(slide, 7.0, 4.55, 5.5, 0.2, 'Decision options:', font_size=10, color=TEXT3)
add_bullet_list(slide, 7.0, 4.8, 5.5, 1.0, decisions, font_size=11, color=TEXT)

add_text_box(slide, 7.0, 5.6, 5.5, 0.2, 'Sherpa Name: ________________________', font_size=11, color=TEXT2)
add_text_box(slide, 7.0, 5.9, 5.5, 0.2, 'Notes / Conditions:', font_size=10, color=TEXT3)
add_shape_box(slide, 7.0, 6.15, 5.5, 0.5, DARK_BG, BORDER)

# ---- SLIDE 4: Intake Email Templates ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_shape_box(slide, 0, 0, 13.333, 0.06, GREEN)
add_section_badge(slide, 'FORM 1 \u2022 EMAILS')

add_text_box(slide, 0.5, 0.75, 8, 0.5, 'Intake Email Templates', font_size=28, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.25, 10, 0.3, 'Copy and customize for common intake communications',
             font_size=12, color=TEXT3)

# Left: Rejection email
box = add_shape_box(slide, 0.5, 1.7, 6.0, 5.0, SURFACE, BORDER)
add_text_box(slide, 0.7, 1.75, 5.5, 0.3, 'Pre-Screen Rejection', font_size=14, color=DANGER, bold=True)
rejection = """Subject: Software Request Pre-Screen - Action Required

Hi [Requester Name],

Thank you for submitting your request for [Software/Vendor].

After our pre-screen review, we've identified the following
items that need to be addressed before we can enter this
into the formal onboarding pipeline:

  [  ] Budget approval from cost center owner
  [  ] Executive sponsor identified (VP+ level)
  [  ] Confirmation no existing solution covers this need
  [  ] Strategic alignment documentation
  [  ] Realistic timeline (minimum 30 days)

Please provide the missing items and resubmit.

If this is an emergency or regulatory-mandated change,
let me know and we'll route it through our expedited path.

[Sherpa Name], Onboarding Facilitator"""
add_text_box(slide, 0.7, 2.1, 5.5, 4.5, rejection, font_size=8, color=TEXT2, font_name='Consolas')

# Right: Acceptance email
box = add_shape_box(slide, 6.8, 1.7, 6.0, 5.0, SURFACE, BORDER)
add_text_box(slide, 7.0, 1.75, 5.5, 0.3, 'Request Accepted - Next Steps', font_size=14, color=GREEN, bold=True)
accepted = """Subject: Software Onboarding Request Accepted - [Tracking ID]

Hi [Requester Name],

Your request for [Software/Vendor Name] has passed
pre-screening and is now in our onboarding pipeline.

  Tracking ID:    [ONB-2026-XXX]
  Risk Tier:      [Fast-Track / Standard / Enhanced]
  Target SLA:     [5 / 15 / 25] business days
  Your Sherpa:    [Sherpa Name]

Next Steps:
1. I will coordinate with the required review teams
2. You may be contacted for additional details
3. I'll send you weekly status updates every [Day]

You can reach me directly with any questions. I'm your
single point of contact throughout this process.

[Sherpa Name], Onboarding Facilitator"""
add_text_box(slide, 7.0, 2.1, 5.5, 4.5, accepted, font_size=8, color=TEXT2, font_name='Consolas')

# ---- SLIDE 5: Risk Tier Classification ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_shape_box(slide, 0, 0, 13.333, 0.06, ACCENT)
add_section_badge(slide, 'FORM 2 OF 8')

add_text_box(slide, 0.5, 0.75, 8, 0.5, 'Risk Tier Classification Worksheet', font_size=28, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.25, 10, 0.3, 'Score each dimension 1-3. Total determines Fast-Track / Standard / Enhanced review level.',
             font_size=12, color=TEXT3)

# Scoring matrix table
table = add_table(slide, 0.5, 1.7, 12.3, 8, 5, [2.0, 2.8, 3.0, 3.0, 1.5])
headers = ['Dimension', '1 - Low', '2 - Medium', '3 - High', 'Score']
data = [
    ['Data Sensitivity', 'No PII/PHI, public data', 'Internal data, limited PII', 'Regulated (PII/PHI/SOX)', '____'],
    ['Vendor Criticality', 'Existing vendor, proven', 'Known vendor, new product', 'New vendor, no history', '____'],
    ['Financial Exposure', 'Under $50K annual', '$50K - $500K annual', 'Over $500K annual', '____'],
    ['Integration Complexity', 'Standalone SaaS', 'API with non-critical systems', 'Core banking / payments', '____'],
    ['AI / ML Usage', 'No AI components', 'AI features (analytics)', 'AI in decisions (credit/fraud)', '____'],
    ['Regulatory Exposure', 'No regulatory reqs', 'Standard (GDPR, SOX)', 'Multi-reg (OCC, DORA, EU AI)', '____'],
    ['', '', '', 'TOTAL:', '____'],
]
for j, h in enumerate(headers):
    table.cell(0, j).text = h
for i, row_data in enumerate(data):
    for j, val in enumerate(row_data):
        table.cell(i + 1, j).text = val
style_table(table)
# Bold the total row
for j in range(5):
    for p in table.cell(7, j).text_frame.paragraphs:
        p.font.bold = True
        p.font.color.rgb = WHITE

# Tier routing guide
box = add_shape_box(slide, 0.5, 5.3, 12.3, 1.8, SURFACE, BORDER)
add_text_box(slide, 0.7, 5.35, 5, 0.3, 'Tier Routing Guide', font_size=14, color=WHITE, bold=True)

tier_table = add_table(slide, 0.7, 5.7, 11.9, 4, 5, [1.0, 1.5, 3.5, 1.5, 4.4])
tier_headers = ['Score', 'Tier', 'Reviewers Required', 'Target SLA', 'Escalation Path']
tier_data = [
    ['6-8', 'FAST-TRACK', 'Sherpa + 1 reviewer (Security or Compliance)', '5 days', "Sherpa's manager after Day 3"],
    ['9-13', 'STANDARD', '4 reviewers parallel (Security, Compliance, Legal, Tech Arch)', '15 days', 'Governance lead after Day 10'],
    ['14-18', 'ENHANCED', 'All 9 reviewers parallel + executive briefing', '25 days', 'CISO/CRO after Day 15'],
]
for j, h in enumerate(tier_headers):
    tier_table.cell(0, j).text = h
for i, row_data in enumerate(tier_data):
    for j, val in enumerate(row_data):
        tier_table.cell(i + 1, j).text = val
style_table(tier_table)
# Color-code tier names
colors = [LIGHT_GREEN, GOLD, RGBColor(0xF8, 0x71, 0x71)]
for i, c in enumerate(colors):
    for p in tier_table.cell(i + 1, 1).text_frame.paragraphs:
        p.font.color.rgb = c
        p.font.bold = True

# Override note
add_text_box(slide, 0.7, 7.0, 11.5, 0.3,
             'Override: If ANY single dimension = 3, minimum tier is Standard. AI + Regulatory both = 3 \u2192 Enhanced automatically.',
             font_size=10, color=WARN)

# ---- SLIDE 6: Request Tracking Board ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_shape_box(slide, 0, 0, 13.333, 0.06, ACCENT)
add_section_badge(slide, 'TOOL 3 OF 8')

add_text_box(slide, 0.5, 0.75, 8, 0.5, 'Request Tracking Board', font_size=28, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.25, 10, 0.3, 'Visual tracker for all active requests. Update weekly. Color-code by SLA status.',
             font_size=12, color=TEXT3)

# Kanban columns
phases = ['Intake', 'Due Diligence', 'Gov. Review', 'Contracting', 'Go-Live']
phase_colors = [GREEN, ACCENT, WARN, ACCENT, GREEN]
for i, (phase, pc) in enumerate(zip(phases, phase_colors)):
    x = 0.5 + i * 2.5
    box = add_shape_box(slide, x, 1.7, 2.2, 2.0, SURFACE, BORDER)
    add_text_box(slide, x + 0.1, 1.75, 2.0, 0.25, phase, font_size=11, color=pc, bold=True)
    # Example card placeholder
    add_shape_box(slide, x + 0.1, 2.1, 2.0, 0.45, DARK_BG, BORDER)
    add_text_box(slide, x + 0.15, 2.12, 1.9, 0.2, '[Request Name]', font_size=9, color=TEXT, bold=True)
    add_text_box(slide, x + 0.15, 2.32, 1.9, 0.15, 'Day X of Y  |  Tier', font_size=8, color=TEXT3)

# Register table
add_text_box(slide, 0.5, 4.0, 5, 0.3, 'Active Request Register', font_size=14, color=WHITE, bold=True)
add_text_box(slide, 0.5, 4.3, 10, 0.2, 'Master list \u2014 print weekly for the stand-up', font_size=10, color=TEXT3)

reg_table = add_table(slide, 0.5, 4.6, 12.3, 5, 9, [1.0, 1.6, 1.2, 0.8, 1.4, 1.0, 1.0, 1.0, 3.3])
reg_headers = ['ID', 'Software/Vendor', 'Requester', 'Tier', 'Current Phase', 'Days In', 'SLA Target', 'Status', 'Blocker']
for j, h in enumerate(reg_headers):
    reg_table.cell(0, j).text = h
# Fill example rows
examples = [
    ['ONB-001', 'Snowflake', 'J. Smith', 'Std', 'Due Diligence', '7', 'Day 15', 'On Track', ''],
    ['ONB-002', 'Custom CRM', 'A. Jones', 'Enh', 'Gov. Review', '12', 'Day 25', 'At Risk', 'Waiting on security assessment'],
    ['ONB-003', 'Slack Enterprise', 'M. Chen', 'Fast', 'Contracting', '3', 'Day 5', 'On Track', ''],
    ['ONB-004', 'AI Analytics', 'R. Patel', 'Enh', 'Due Diligence', '18', 'Day 25', 'Breached', 'Legal review overdue 3 days'],
]
for i, row_data in enumerate(examples):
    for j, val in enumerate(row_data):
        reg_table.cell(i + 1, j).text = val
style_table(reg_table)
# Color-code status
status_colors = {'On Track': LIGHT_GREEN, 'At Risk': GOLD, 'Breached': RGBColor(0xF8, 0x71, 0x71)}
for i in range(1, 5):
    status = reg_table.cell(i, 7).text
    if status in status_colors:
        for p in reg_table.cell(i, 7).text_frame.paragraphs:
            p.font.color.rgb = status_colors[status]
            p.font.bold = True

# ---- SLIDE 7: Parallel Review Kickoff ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_shape_box(slide, 0, 0, 13.333, 0.06, WARN)
add_section_badge(slide, 'FORM 4 OF 8')

add_text_box(slide, 0.5, 0.75, 10, 0.5, 'Parallel Review Kickoff Package', font_size=28, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.25, 10, 0.3, 'Send to ALL required reviewers simultaneously. Do NOT wait for sequential reviews.',
             font_size=12, color=TEXT3)

# Review assignment matrix
rev_table = add_table(slide, 0.5, 1.7, 12.3, 10, 6, [0.5, 2.5, 2.0, 1.5, 1.5, 4.3])
rev_headers = ["Req'd", 'Review Domain', 'Reviewer', 'Sent Date', 'Due Date', 'Status']
domains = [
    ('Always', 'Security Assessment', 'Pen testing, vuln scan, architecture'),
    ('', 'Tech Architecture', 'Integration, scalability, infrastructure'),
    ('', 'Risk Assessment', 'Operational, financial, reputational risk'),
    ('', 'Compliance Review', 'Regulatory mapping, control requirements'),
    ('', 'Privacy Assessment', 'GDPR/CCPA, data classification, cross-border'),
    ('', 'Legal Review', 'Contract terms, IP, liability, indemnification'),
    ('', 'Financial Analysis', 'TCO, ROI, budget impact, NPV'),
    ('', 'Vendor Landscape', 'Alternative vendors, market analysis'),
    ('', 'AI Governance', 'Model risk, EU AI Act, SR 11-7'),
]
for j, h in enumerate(rev_headers):
    rev_table.cell(0, j).text = h
for i, (req, domain, desc) in enumerate(domains):
    rev_table.cell(i + 1, 0).text = '\u2610' if not req else '\u2611'
    rev_table.cell(i + 1, 1).text = f'{domain}\n{desc}'
    rev_table.cell(i + 1, 2).text = ''
    rev_table.cell(i + 1, 3).text = '____/____'
    rev_table.cell(i + 1, 4).text = '____/____'
    rev_table.cell(i + 1, 5).text = 'Not Sent / Sent / In Review / Complete / Overdue'
style_table(rev_table)
# Make domain names white
for i in range(1, 10):
    for p in rev_table.cell(i, 1).text_frame.paragraphs:
        if p.text and not p.text.startswith(' '):
            p.font.color.rgb = WHITE

# ---- SLIDE 8: Review Kickoff Email + Findings Summary ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_shape_box(slide, 0, 0, 13.333, 0.06, WARN)
add_section_badge(slide, 'FORM 4 \u2022 EMAILS')

add_text_box(slide, 0.5, 0.75, 8, 0.5, 'Review Kickoff & Findings', font_size=28, color=WHITE, bold=True)

# Left: Kickoff email
box = add_shape_box(slide, 0.5, 1.3, 6.0, 5.5, SURFACE, BORDER)
add_text_box(slide, 0.7, 1.35, 5, 0.3, 'Review Kickoff Email', font_size=14, color=WARN, bold=True)
kickoff_email = """Subject: REVIEW REQUESTED - [Software] - Due [Date]

Hi [Reviewer Name],

A software onboarding request requires your
[Domain] review.

  Request ID:     [ONB-2026-XXX]
  Software:       [Name and brief description]
  Vendor:         [Name, existing/new]
  Risk Tier:      [Fast-Track / Standard / Enhanced]
  Your Deadline:  [Date - X business days]

YOUR REVIEW SCOPE:
  [Specific questions for this domain]

DELIVERABLE:
Pass / Conditional Pass / Fail with findings.

NOTE: This review runs IN PARALLEL with [N]
other reviews. You do NOT need to wait.

If you cannot meet the deadline, let me know
by [Date + 1] so I can arrange coverage.

[Sherpa Name]"""
add_text_box(slide, 0.7, 1.7, 5.5, 5.0, kickoff_email, font_size=8, color=TEXT2, font_name='Consolas')

# Right: Findings summary
box = add_shape_box(slide, 6.8, 1.3, 6.0, 5.5, SURFACE, BORDER)
add_text_box(slide, 7.0, 1.35, 5, 0.3, 'Reviewer Findings Summary', font_size=14, color=WHITE, bold=True)
add_text_box(slide, 7.0, 1.65, 5.5, 0.2, 'Consolidate before presenting to governance committee', font_size=10, color=TEXT3)

find_table = add_table(slide, 7.0, 1.95, 5.6, 10, 3, [1.5, 1.0, 3.1])
find_table.cell(0, 0).text = 'Domain'
find_table.cell(0, 1).text = 'Result'
find_table.cell(0, 2).text = 'Key Findings / Conditions'
find_domains = ['Security', 'Tech Arch', 'Risk', 'Compliance', 'Privacy', 'Legal', 'Financial', 'Vendor Landscape', 'AI Governance']
for i, d in enumerate(find_domains):
    find_table.cell(i + 1, 0).text = d
    find_table.cell(i + 1, 1).text = 'Pass / Cond / Fail'
    find_table.cell(i + 1, 2).text = ''
style_table(find_table)

# ---- SLIDE 9: Weekly Stand-Up ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_shape_box(slide, 0, 0, 13.333, 0.06, WARN)
add_section_badge(slide, 'TEMPLATE 5 OF 8')

add_text_box(slide, 0.5, 0.75, 10, 0.5, 'Weekly Governance Stand-Up', font_size=28, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.25, 10, 0.3, '30-minute maximum. Same time every week. Sherpa facilitates. Status and decisions only \u2014 no deep dives.',
             font_size=12, color=TEXT3)

# Agenda table
agenda_table = add_table(slide, 0.5, 1.7, 12.3, 6, 4, [1.2, 2.5, 6.6, 2.0])
agenda_table.cell(0, 0).text = 'Time Box'
agenda_table.cell(0, 1).text = 'Agenda Item'
agenda_table.cell(0, 2).text = 'Description'
agenda_table.cell(0, 3).text = 'Owner'

agenda = [
    ('2 min', 'Pipeline Summary', 'How many requests total? How many in each phase? How many breached SLA?', 'Sherpa'),
    ('10 min', 'SLA Breaches', 'Walk through EACH breached request. What is blocking? Who needs to act? Decision: escalate, extend, or kill?', 'Sherpa + Blockers'),
    ('8 min', 'At-Risk Requests', 'Requests approaching SLA deadline (within 2 days). What\'s needed to close on time?', 'Sherpa + Reviewers'),
    ('5 min', 'New Intake', 'Any new requests entering the pipeline this week? Quick tier assignment and reviewer identification.', 'Sherpa'),
    ('5 min', 'Decisions Needed', 'Any go/no-go decisions the committee needs to make right now? Governance approvals waiting?', 'Governance Lead'),
]
for i, (time, item, desc, owner) in enumerate(agenda):
    agenda_table.cell(i + 1, 0).text = time
    agenda_table.cell(i + 1, 1).text = item
    agenda_table.cell(i + 1, 2).text = desc
    agenda_table.cell(i + 1, 3).text = owner
style_table(agenda_table)
# Color the time boxes
time_colors = [ACCENT, RGBColor(0xF8, 0x71, 0x71), GOLD, LIGHT_GREEN, TEXT2]
for i, c in enumerate(time_colors):
    for p in agenda_table.cell(i + 1, 0).text_frame.paragraphs:
        p.font.color.rgb = c
        p.font.bold = True

# Ground rules
box = add_shape_box(slide, 0.5, 4.5, 12.3, 0.6, FAINT_BLUE, RGBColor(0x20, 0x40, 0x80))
add_text_box(slide, 0.7, 4.55, 11.8, 0.5,
             'Ground Rules: No deep-dive discussions. If a topic needs > 2 minutes, schedule separately. '
             'Sherpa publishes action items within 1 hour.',
             font_size=11, color=RGBColor(0x60, 0xA5, 0xFA))

# Notes template
box = add_shape_box(slide, 0.5, 5.3, 6.0, 2.0, SURFACE, BORDER)
add_text_box(slide, 0.7, 5.35, 5, 0.3, 'Meeting Notes Template', font_size=14, color=WHITE, bold=True)
notes = [
    'Meeting Date: ____/____/____    Attendees: ________',
    'SLA Breaches: ID | Blocker | Decision | Owner | Deadline',
    'At-Risk Items: ID | Risk | Mitigation | Owner',
    'Decisions Made: Decision | Rationale | Who needs to know',
    'Action Items: Action | Owner | Due Date',
]
add_bullet_list(slide, 0.7, 5.7, 5.5, 1.5, notes, font_size=10, color=TEXT2)

# Summary email
box = add_shape_box(slide, 6.8, 5.3, 6.0, 2.0, SURFACE, BORDER)
add_text_box(slide, 7.0, 5.35, 5, 0.3, 'Post-Meeting Summary Email', font_size=14, color=WHITE, bold=True)
summary = """Subject: Governance Stand-Up Summary - [Date]

Pipeline: [N] active | [N] breached | [N] new
SLA Breaches: [IDs and decisions]
At-Risk: [IDs and needed actions]
Decisions: [Made this week]
Action Items: [Owner - action - due date]

Next stand-up: [Date/Time]"""
add_text_box(slide, 7.0, 5.7, 5.5, 1.5, summary, font_size=8, color=TEXT2, font_name='Consolas')

# ---- SLIDE 10: SLA Escalation Playbook ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_shape_box(slide, 0, 0, 13.333, 0.06, DANGER)
add_section_badge(slide, 'TEMPLATE 6 OF 8')

add_text_box(slide, 0.5, 0.75, 10, 0.5, 'SLA Escalation Playbook', font_size=28, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.25, 10, 0.3, 'The Sherpa owns escalation. Do not wait for someone else to notice a breach.',
             font_size=12, color=TEXT3)

# Escalation ladder
esc_table = add_table(slide, 0.5, 1.7, 12.3, 8, 5, [2.0, 1.5, 2.0, 3.0, 3.8])
esc_headers = ['Phase', 'SLA Target', 'Warning Sent', 'Level 1 Escalation', 'Level 2 Escalation']
esc_data = [
    ['Intake / Triage', '2 business days', 'Day 1 (50%)', "Sherpa's manager (Day 2)", 'Governance Lead (Day 3)'],
    ['Planning', '5 business days', 'Day 3 (60%)', 'Governance Lead (Day 5)', 'CRO (Day 7)'],
    ['Due Diligence', '10 business days', 'Day 7 (70%)', 'CISO / CLO (Day 10)', 'CRO + CFO (Day 14)'],
    ['Individual Review', '5 business days', 'Day 3 (60%)', "Reviewer's manager (Day 5)", 'Domain head (Day 7)'],
    ['Governance Review', '5 business days', 'Day 3 (60%)', 'CRO (Day 5)', 'COO (Day 7)'],
    ['Contracting', '7 business days', 'Day 5 (70%)', 'CLO / CPO (Day 7)', 'CFO (Day 10)'],
    ['UAT / Go-Live', '5 business days', 'Day 3 (60%)', 'Project Sponsor (Day 5)', 'CTO (Day 7)'],
]
for j, h in enumerate(esc_headers):
    esc_table.cell(0, j).text = h
for i, row_data in enumerate(esc_data):
    for j, val in enumerate(row_data):
        esc_table.cell(i + 1, j).text = val
style_table(esc_table)

# Email template summary
add_text_box(slide, 0.5, 5.2, 12, 0.3, 'Three Escalation Email Templates', font_size=14, color=WHITE, bold=True)

# Three boxes side by side
for idx, (title, color, desc) in enumerate([
    ('Level 0: Friendly Reminder', LIGHT_GREEN, 'At warning threshold\n\n"Quick reminder that your [domain]\nreview is due in [N] days.\n\nIf you need info or are blocked,\nlet me know so I can help."'),
    ('Level 1: Formal Escalation', GOLD, 'At SLA breach\n\n"I\'m escalating [ONB-XXX] which\nhas exceeded its SLA in [Phase].\n\nSLA Target: [N] days\nActual: [N] days (overdue [N])\nBlocker: [specific reason]"'),
    ('Level 2: Executive', RGBColor(0xF8, 0x71, 0x71), 'After Level 1 fails\n\n"[ONB-XXX] has been stalled for\n[N] days beyond SLA. Level 1 on\n[Date] has not resolved.\n\nOPTIONS:\n1. [Recommended path]\n2. [Alternative]\n3. Kill the request"'),
]):
    x = 0.5 + idx * 4.2
    box = add_shape_box(slide, x, 5.6, 3.9, 1.7, SURFACE, BORDER)
    add_text_box(slide, x + 0.15, 5.65, 3.6, 0.25, title, font_size=11, color=color, bold=True)
    add_text_box(slide, x + 0.15, 5.9, 3.6, 1.3, desc, font_size=8, color=TEXT2, font_name='Consolas')

# ---- SLIDE 11: Phase Transition Checklists (Part 1) ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_shape_box(slide, 0, 0, 13.333, 0.06, GREEN)
add_section_badge(slide, 'CHECKLISTS 7 OF 8')

add_text_box(slide, 0.5, 0.75, 10, 0.5, 'Phase Transition Checklists', font_size=28, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.25, 10, 0.3, 'A request cannot advance until ALL items checked. Sherpa verifies and obtains sign-off.',
             font_size=12, color=TEXT3)

checklists = [
    ('Phase 1: Intake \u2192 Planning', [
        ('\u2610 Pre-screen form completed (all deal-killers pass)', 'Sherpa'),
        ('\u2610 Existing portfolio checked (no duplicate)', 'Sherpa'),
        ('\u2610 Data classification completed', 'Sherpa'),
        ('\u2610 Risk tier assigned (Fast-Track / Standard / Enhanced)', 'Sherpa'),
        ('\u2610 Executive sponsor confirmed', 'Requester'),
        ('\u2610 Request entered in tracking register', 'Sherpa'),
        ('\u2610 Requester notified with tracking ID and SLA', 'Sherpa'),
    ]),
    ('Phase 2: Planning \u2192 Due Diligence', [
        ('\u2610 Buy vs Build pathway determined', 'Governance'),
        ('\u2610 NDA executed (if new vendor)', 'Legal'),
        ('\u2610 Preliminary business analysis documented', 'Business'),
        ('\u2610 Review team identified based on tier', 'Sherpa'),
        ('\u2610 Vendor documentation package assembled', 'Sherpa'),
        ('\u2610 Review kickoff package prepared', 'Sherpa'),
    ]),
    ('Phase 3: Due Diligence \u2192 Gov. Review', [
        ('\u2610 ALL required reviews completed (per tier)', 'Sherpa'),
        ('\u2610 Findings consolidated (pass/cond/fail per domain)', 'Sherpa'),
        ('\u2610 No unresolved "Fail" findings', 'Sherpa'),
        ('\u2610 Conditional findings have remediation plans', 'Reviewers'),
        ('\u2610 Vendor DD response received and evaluated', 'Procurement'),
        ('\u2610 Governance committee briefing prepared', 'Sherpa'),
    ]),
]

for i, (title, items) in enumerate(checklists):
    x = 0.5 + i * 4.2
    box = add_shape_box(slide, x, 1.7, 3.9, 3.5, SURFACE, BORDER)
    add_text_box(slide, x + 0.15, 1.75, 3.6, 0.3, title, font_size=12, color=LIGHT_GREEN, bold=True)
    item_texts = [f'{item}  [{owner}]' for item, owner in items]
    add_bullet_list(slide, x + 0.15, 2.1, 3.6, 3.0, item_texts, font_size=8, color=TEXT2)

# Bottom row
checklists2 = [
    ('Phase 4: Gov. Review \u2192 Contracting', [
        ('\u2610 Committee vote recorded', 'Governance'),
        ('\u2610 Conditions documented (if conditional)', 'Governance'),
        ('\u2610 Rejection communicated (if rejected)', 'Sherpa'),
        ('\u2610 Approval memo signed', 'Gov. Lead'),
        ('\u2610 Budget authorization confirmed', 'Finance'),
        ('\u2610 Contract requirements identified', 'Legal'),
    ]),
    ('Phase 5: Contracting \u2192 Go-Live', [
        ('\u2610 Contract fully executed (both parties)', 'Legal'),
        ('\u2610 SLA terms documented in contract', 'Legal'),
        ('\u2610 Security controls contractually required', 'Security'),
        ('\u2610 UAT plan with acceptance criteria', 'Business'),
        ('\u2610 Implementation timeline agreed', 'Sherpa'),
        ('\u2610 Conditional items remediated/tracked', 'Sherpa'),
    ]),
    ('Close-Out: Go-Live \u2192 Operations', [
        ('\u2610 UAT completed and signed off', 'Business'),
        ('\u2610 Production deployment completed', 'IT'),
        ('\u2610 Monitoring and alerting configured', 'IT/Security'),
        ('\u2610 Onboarding record archived', 'Sherpa'),
        ('\u2610 Requester notified of completion', 'Sherpa'),
        ('\u2610 Monitoring cadence assigned (per tier)', 'Governance'),
        ('\u2610 Contract renewal date flagged', 'Sherpa'),
    ]),
]

for i, (title, items) in enumerate(checklists2):
    x = 0.5 + i * 4.2
    box = add_shape_box(slide, x, 5.4, 3.9, 1.9, SURFACE, BORDER)
    add_text_box(slide, x + 0.15, 5.45, 3.6, 0.3, title, font_size=12, color=LIGHT_GREEN, bold=True)
    item_texts = [f'{item}  [{owner}]' for item, owner in items]
    add_bullet_list(slide, x + 0.15, 5.75, 3.6, 1.5, item_texts, font_size=8, color=TEXT2)

# ---- SLIDE 12: Request Close-Out ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_shape_box(slide, 0, 0, 13.333, 0.06, ACCENT)
add_section_badge(slide, 'FORM 8 OF 8')

add_text_box(slide, 0.5, 0.75, 8, 0.5, 'Request Close-Out Form', font_size=28, color=WHITE, bold=True)
add_text_box(slide, 0.5, 1.25, 10, 0.3, 'Complete when a request reaches a terminal state. This is the permanent record.',
             font_size=12, color=TEXT3)

# Close-out summary
box = add_shape_box(slide, 0.5, 1.7, 6.0, 2.0, SURFACE, BORDER)
add_text_box(slide, 0.7, 1.75, 5, 0.3, 'Close-Out Summary', font_size=14, color=WHITE, bold=True)
fields = [
    'Request ID: ONB-____-____          Software: ____________________',
    'Outcome:  \u2610 Onboarded   \u2610 Rejected   \u2610 Withdrawn   \u2610 Redirected',
    'Risk Tier: \u2610 Fast-Track   \u2610 Standard   \u2610 Enhanced',
    'Date Submitted: ____/____   Date Closed: ____/____',
    'Total Elapsed Days: ____',
]
add_bullet_list(slide, 0.7, 2.1, 5.5, 1.5, fields, font_size=10, color=TEXT2)

# SLA Performance
box = add_shape_box(slide, 6.8, 1.7, 6.0, 2.0, SURFACE, BORDER)
add_text_box(slide, 7.0, 1.75, 5, 0.3, 'SLA Performance', font_size=14, color=WHITE, bold=True)

perf_table = add_table(slide, 7.0, 2.1, 5.6, 7, 4, [1.5, 1.0, 1.0, 2.1])
perf_headers = ['Phase', 'Target', 'Actual', 'Delay Reason']
perf_phases = ['Intake', 'Planning', 'Due Diligence', 'Gov. Review', 'Contracting', 'Go-Live']
perf_targets = ['2d', '5d', '10d', '5d', '7d', '5d']
for j, h in enumerate(perf_headers):
    perf_table.cell(0, j).text = h
for i, (phase, target) in enumerate(zip(perf_phases, perf_targets)):
    perf_table.cell(i + 1, 0).text = phase
    perf_table.cell(i + 1, 1).text = target
    perf_table.cell(i + 1, 2).text = '____'
    perf_table.cell(i + 1, 3).text = ''
style_table(perf_table)

# Lessons learned
box = add_shape_box(slide, 0.5, 3.9, 6.0, 3.2, SURFACE, BORDER)
add_text_box(slide, 0.7, 3.95, 5, 0.3, 'Lessons Learned', font_size=14, color=WHITE, bold=True)
add_text_box(slide, 0.7, 4.25, 5.5, 0.15, 'Review quarterly to improve the process', font_size=10, color=TEXT3)
lessons = [
    'What went well?',
    '________________________________________',
    '',
    'What caused delays?',
    '________________________________________',
    '',
    'Process improvement suggestions:',
    '________________________________________',
]
add_bullet_list(slide, 0.7, 4.5, 5.5, 2.5, lessons, font_size=10, color=TEXT2)

# Close-out email
box = add_shape_box(slide, 6.8, 3.9, 6.0, 3.2, SURFACE, BORDER)
add_text_box(slide, 7.0, 3.95, 5, 0.3, 'Close-Out Notification Email', font_size=14, color=WHITE, bold=True)
close_email = """Subject: Onboarding Complete - [Software Name]

Hi [Requester],

  Request ID:     [ONB-2026-XXX]
  Outcome:        [Onboarded / Rejected / Withdrawn]
  Total Duration: [N] business days
  SLA Met:        [Yes / No - overdue by N days]

If onboarded:
  Monitoring:     [Team/Person]
  Review cadence: [Quarterly / Semi-annual / Annual]
  Contract renewal: [Date]

If rejected:
  Rationale: [Summary]
  Alternatives: [If applicable]

[Sherpa Name], Onboarding Facilitator"""
add_text_box(slide, 7.0, 4.3, 5.5, 2.7, close_email, font_size=8, color=TEXT2, font_name='Consolas')

# ---- SLIDE 13: Quick Reference Card ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)
add_shape_box(slide, 0, 0, 13.333, 0.06, GREEN)

add_text_box(slide, 0.5, 0.3, 10, 0.5, 'Sherpa Quick Reference Card', font_size=28, color=WHITE, bold=True)
add_text_box(slide, 0.5, 0.8, 10, 0.3, 'Print this page and keep at your desk. One card for the entire lifecycle.',
             font_size=12, color=TEXT3)

# Left: Lifecycle at a glance
box = add_shape_box(slide, 0.5, 1.3, 4.0, 5.8, SURFACE, BORDER)
add_text_box(slide, 0.7, 1.35, 3.5, 0.3, 'Lifecycle at a Glance', font_size=14, color=LIGHT_GREEN, bold=True)

lifecycle = [
    ('1. INTAKE', '2 days', 'Pre-screen, deal-killers, tier assign'),
    ('2. PLANNING', '5 days', 'Pathway, NDA, review team ID'),
    ('3. DUE DILIGENCE', '10 days', 'All reviews IN PARALLEL'),
    ('4. GOV. REVIEW', '5 days', 'Committee vote'),
    ('5. CONTRACTING', '7 days', 'Contract execution'),
    ('6. GO-LIVE', '5 days', 'UAT, deploy, monitor setup'),
]
for i, (phase, sla, desc) in enumerate(lifecycle):
    y = 1.75 + i * 0.85
    add_text_box(slide, 0.7, y, 1.8, 0.2, phase, font_size=10, color=WHITE, bold=True)
    add_text_box(slide, 2.6, y, 1.0, 0.2, sla, font_size=10, color=WARN, bold=True)
    add_text_box(slide, 0.7, y + 0.22, 3.5, 0.35, desc, font_size=9, color=TEXT3)

# Center: Key rules
box = add_shape_box(slide, 4.7, 1.3, 4.2, 5.8, SURFACE, BORDER)
add_text_box(slide, 4.9, 1.35, 3.8, 0.3, 'Sherpa Rules', font_size=14, color=LIGHT_GREEN, bold=True)

rules = [
    'YOU own elapsed time. Nobody else is watching.',
    'NEVER wait for sequential reviews. Send all at once.',
    'ALWAYS set a deadline. No open-ended requests.',
    'ESCALATE at threshold. Not after. Not before.',
    'WEEKLY stand-up is sacred. 30 min, no exceptions.',
    'TRACK everything. If it\'s not on the board, it doesn\'t exist.',
    'FILTER early. 1 rejected pre-screen saves 10 reviewer-hours.',
    'CLOSE the loop. Requester hears from you, not silence.',
    'MEASURE cycle time. What gets measured gets improved.',
    'DOCUMENT lessons learned. Every close-out form, every time.',
]
for i, rule in enumerate(rules):
    y = 1.7 + i * 0.5
    num = str(i + 1)
    add_text_box(slide, 4.9, y, 0.3, 0.2, num + '.', font_size=10, color=WARN, bold=True)
    add_text_box(slide, 5.2, y, 3.5, 0.4, rule, font_size=9, color=TEXT2)

# Right: SLA cheat sheet
box = add_shape_box(slide, 9.1, 1.3, 3.7, 2.7, SURFACE, BORDER)
add_text_box(slide, 9.3, 1.35, 3.3, 0.3, 'Tier Cheat Sheet', font_size=14, color=LIGHT_GREEN, bold=True)

tier_data = [
    ('FAST-TRACK', '6-8 pts', '5 days', '1 reviewer', LIGHT_GREEN),
    ('STANDARD', '9-13 pts', '15 days', '4 reviewers', GOLD),
    ('ENHANCED', '14-18 pts', '25 days', '9 reviewers', RGBColor(0xF8, 0x71, 0x71)),
]
for i, (tier, pts, sla, rev, c) in enumerate(tier_data):
    y = 1.75 + i * 0.7
    add_text_box(slide, 9.3, y, 1.5, 0.2, tier, font_size=10, color=c, bold=True)
    add_text_box(slide, 9.3, y + 0.2, 3.3, 0.2, f'{pts}  |  {sla}  |  {rev}', font_size=9, color=TEXT3)

# Right bottom: Escalation cheat sheet
box = add_shape_box(slide, 9.1, 4.2, 3.7, 2.9, SURFACE, BORDER)
add_text_box(slide, 9.3, 4.25, 3.3, 0.3, 'Escalation Cheat Sheet', font_size=14, color=LIGHT_GREEN, bold=True)

esc_cheat = [
    'At 50-70% of SLA:', 'Send friendly reminder', '',
    'At 100% of SLA:', 'Level 1 \u2192 reviewer\'s manager', '',
    'At 140% of SLA:', 'Level 2 \u2192 executive (CISO/CRO/CFO)', '',
    'Key principle:', 'Be factual, not emotional.', 'Options, not complaints.',
]
for i in range(0, len(esc_cheat), 3):
    y = 4.65 + (i // 3) * 0.6
    add_text_box(slide, 9.3, y, 3.3, 0.2, esc_cheat[i], font_size=9, color=WARN, bold=True)
    add_text_box(slide, 9.3, y + 0.18, 3.3, 0.2, esc_cheat[i + 1], font_size=9, color=TEXT2)
    if esc_cheat[i + 2]:
        add_text_box(slide, 9.3, y + 0.34, 3.3, 0.15, esc_cheat[i + 2], font_size=8, color=TEXT3)

# ---- Save ----
output_path = os.path.join(os.path.dirname(__file__), 'sherpa-toolkit.pptx')
prs.save(output_path)
print(f'Saved: {output_path}')
print(f'Slides: {len(prs.slides)}')
