#!/usr/bin/env python3
"""Generate Vendor Onboarding Governance PowerPoint from presentation data."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
import os

# Brand colors
NAVY = RGBColor(0x0D, 0x21, 0x37)
DEEP_BLUE = RGBColor(0x1A, 0x3A, 0x6B)
BLUE = RGBColor(0x25, 0x63, 0xEB)
TEAL = RGBColor(0x00, 0xB4, 0xD8)
GOLD = RGBColor(0xD9, 0x77, 0x06)
EMERALD = RGBColor(0x05, 0x96, 0x69)
ROSE = RGBColor(0xDC, 0x26, 0x26)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
GRAY = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
BORDER_COLOR = RGBColor(0xE2, 0xE8, 0xF0)
LIGHT_BLUE_BG = RGBColor(0xEF, 0xF6, 0xFF)
LIGHT_GREEN_BG = RGBColor(0xF0, 0xFD, 0xF4)
LIGHT_ORANGE_BG = RGBColor(0xFF, 0xF7, 0xED)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_gradient_bg(slide, color1=NAVY, color2=DEEP_BLUE):
    """Add a solid dark background (gradient approximation)."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color1


def add_light_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or WHITE
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text="", font_size=12, color=DARK_TEXT,
                 bold=False, font_name="Calibri", alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=12, color=DARK_TEXT, bold=False,
                  font_name="Calibri", alignment=PP_ALIGN.LEFT, space_before=Pt(0), space_after=Pt(4)):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_bullet_list(slide, left, top, width, height, items, font_size=11, color=DARK_TEXT, bullet_color=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.level = 0
        p.space_after = Pt(4)
        pf = p._pPr
        if pf is None:
            from pptx.oxml.ns import qn
            pf = p._p.get_or_add_pPr()
        # Add bullet
        from lxml import etree
        from pptx.oxml.ns import qn, nsmap
        buNone = pf.find(qn('a:buNone'))
        if buNone is not None:
            pf.remove(buNone)
        buChar = etree.SubElement(pf, qn('a:buChar'))
        buChar.set('char', '\u2022')
        if bullet_color:
            buClr = etree.SubElement(pf, qn('a:buClr'))
            srgb = etree.SubElement(buClr, qn('a:srgbClr'))
            srgb.set('val', '%02X%02X%02X' % (bullet_color[0], bullet_color[1], bullet_color[2]))
    return txBox


def add_card(slide, left, top, width, height, title, body_lines, accent_color=BLUE, title_size=13, body_size=10):
    """Add a card with left accent border."""
    # Card background
    card = add_shape(slide, left, top, width, height, fill_color=WHITE)
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    bar.shadow.inherit = False
    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.12), width - Inches(0.3), Inches(0.35),
                 title, font_size=title_size, color=DEEP_BLUE, bold=True)
    # Body
    y_offset = top + Inches(0.5)
    for line in body_lines:
        add_text_box(slide, left + Inches(0.2), y_offset, width - Inches(0.3), Inches(0.25),
                     line, font_size=body_size, color=GRAY)
        y_offset += Inches(0.22)
    return card


def add_metric_box(slide, left, top, width, height, value, label, value_color=BLUE):
    box = add_shape(slide, left, top, width, height, fill_color=WHITE)
    add_text_box(slide, left, top + Inches(0.15), width, Inches(0.5),
                 value, font_size=28, color=value_color, bold=True,
                 font_name="Calibri", alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(0.6), width, Inches(0.3),
                 label, font_size=9, color=GRAY, alignment=PP_ALIGN.CENTER)


def add_badge(slide, left, top, text, bg_color=DEEP_BLUE, text_color=WHITE, width=None):
    w = width or Inches(max(1.0, len(text) * 0.09 + 0.3))
    h = Inches(0.3)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_callout(slide, left, top, width, height, text, accent_color=GOLD, bg_alpha=None):
    """Add a callout box with accent border."""
    # Use a lighter version of the accent color for background
    bg_map = {
        GOLD: RGBColor(0xFF, 0xF7, 0xED),
        BLUE: RGBColor(0xEF, 0xF6, 0xFF),
        EMERALD: RGBColor(0xF0, 0xFD, 0xF4),
    }
    bg = bg_map.get(accent_color, RGBColor(0xF8, 0xFA, 0xFC))
    box = add_shape(slide, left, top, width, height, fill_color=bg)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.05), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    bar.shadow.inherit = False
    add_text_box(slide, left + Inches(0.15), top + Inches(0.08), width - Inches(0.25), height - Inches(0.16),
                 text, font_size=10, color=DARK_TEXT)


def add_footer(slide, dark=False):
    color = RGBColor(0xAA, 0xAA, 0xAA) if dark else GRAY
    add_text_box(slide, Inches(0.5), H - Inches(0.4), Inches(2), Inches(0.3),
                 "Confidential", font_size=7, color=color)
    add_text_box(slide, W - Inches(3.5), H - Inches(0.4), Inches(3), Inches(0.3),
                 "Vendor Onboarding Governance", font_size=7, color=color, alignment=PP_ALIGN.RIGHT)


def add_table(slide, left, top, width, rows_data, col_widths=None, header_color=DEEP_BLUE):
    """Add a styled table."""
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, Inches(rows * 0.35))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r_idx, row in enumerate(rows_data):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.name = "Calibri"
                if r_idx == 0:
                    paragraph.font.color.rgb = WHITE
                    paragraph.font.bold = True
                else:
                    paragraph.font.color.rgb = DARK_TEXT
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return table_shape


# ============================================================
# SLIDE 1: HERO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_gradient_bg(slide)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.2),
             "Vendor Onboarding Governance", font_size=44, color=WHITE, bold=True,
             font_name="Calibri")
add_text_box(slide, Inches(0.8), Inches(2.6), Inches(11), Inches(0.8),
             "Ideal State \u2014 Procurement & Legal Executive Briefing", font_size=24, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(3.5), Inches(9), Inches(0.8),
             "A 3-pool, 38-task governance model with end-to-end regulatory traceability, "
             "automated evidence collection, and contract lifecycle management built for financial services.",
             font_size=13, color=RGBColor(0xCC, 0xCC, 0xCC))

# Badges
badges = ["OCC 2023-17", "GDPR / CCPA", "DORA Art. 28-30", "SOX 302/404",
          "SEC 17a-4", "NIST CSF 2.0", "ISO 27001", "BPMN 2.0", "DMN 1.3"]
x = Inches(0.8)
for b in badges:
    add_badge(slide, x, Inches(4.6), b, bg_color=RGBColor(0x2A, 0x4A, 0x7A), text_color=WHITE)
    x += Inches(max(1.1, len(b) * 0.09 + 0.4))

add_footer(slide, dark=True)

# ============================================================
# SLIDE 2: AGENDA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.6),
             "Agenda", font_size=30, color=DEEP_BLUE, bold=True)

# Left column
sections_left = [
    ("Overview", [
        "The Challenge: Current State Pain Points",
        "Key Metrics & Target Outcomes",
        "3-Pool Architecture Overview",
        "End-to-End Process Flow"
    ]),
    ("Procurement Journey", [
        "Vendor Qualification & Intake",
        "Vendor Evaluation Triad",
        "Contract Negotiation & Review",
        "Contract Execution & Compliance"
    ])
]

y = Inches(1.2)
for section_title, items in sections_left:
    add_text_box(slide, Inches(0.8), y, Inches(5), Inches(0.35),
                 section_title, font_size=16, color=NAVY, bold=True)
    y += Inches(0.4)
    for item in items:
        add_text_box(slide, Inches(1.2), y, Inches(5), Inches(0.25),
                     "\u2022  " + item, font_size=11, color=BLUE)
        y += Inches(0.28)
    y += Inches(0.15)

# Right column
sections_right = [
    ("Legal & Regulatory Deep Dive", [
        "Regulatory Framework Coverage",
        "DPA Lifecycle",
        "Audit Rights & Evidence Chain",
        "Risk-Tiered Governance"
    ]),
    ("Automation & Next Steps", [
        "Bottleneck Elimination",
        "Evidence Automation Matrix",
        "Ongoing Monitoring & Close-out",
        "Call to Action"
    ])
]

y = Inches(1.2)
for section_title, items in sections_right:
    add_text_box(slide, Inches(7), y, Inches(5), Inches(0.35),
                 section_title, font_size=16, color=NAVY, bold=True)
    y += Inches(0.4)
    for item in items:
        add_text_box(slide, Inches(7.4), y, Inches(5), Inches(0.25),
                     "\u2022  " + item, font_size=11, color=BLUE)
        y += Inches(0.28)
    y += Inches(0.15)

add_footer(slide)

# ============================================================
# SLIDE 3: THE CHALLENGE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "The Challenge: Current State vs. Target State", font_size=28, color=DEEP_BLUE, bold=True)

# Current state card
add_card(slide, Inches(0.8), Inches(1.3), Inches(5.2), Inches(3.2),
         "Current State Pain Points",
         ["90+ day vendor onboarding cycles",
          "Manual evidence collection across siloed teams",
          "Contract review bottlenecks (avg. 3\u20134 revision cycles)",
          "Fragmented regulatory compliance tracking",
          "No standardized vendor qualification criteria",
          "DPA negotiations handled ad-hoc",
          "Audit findings from inconsistent documentation",
          "Exit planning deferred until contract renewal"],
         accent_color=ROSE, body_size=10)

# Arrow
add_text_box(slide, Inches(6.1), Inches(2.5), Inches(0.8), Inches(0.6),
             "\u2192", font_size=36, color=EMERALD, bold=True, alignment=PP_ALIGN.CENTER)

# Target state card
add_card(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(3.2),
         "Ideal State Target",
         ["Under 30-day onboarding (risk-tiered)",
          "Automated evidence collection with audit trail",
          "Parallel contract and compliance review",
          "7 regulatory frameworks traced to every task",
          "DMN-driven vendor qualification scoring",
          "Templated DPA with GDPR Art. 28(3) coverage",
          "Complete evidence chain from intake to close-out",
          "Exit plan required before contract execution"],
         accent_color=EMERALD, body_size=10)

# Callouts
add_callout(slide, Inches(0.8), Inches(4.7), Inches(5.8), Inches(0.55),
            "Procurement: Vendor qualification standardized with weighted scoring, parallel evaluation tracks eliminate sequential bottlenecks.",
            accent_color=GOLD)
add_callout(slide, Inches(6.8), Inches(4.7), Inches(5.7), Inches(0.55),
            "Legal: Every task carries regulatory annotations (OCC, GDPR, SOX, DORA). DPA lifecycle is embedded, not bolted on.",
            accent_color=BLUE)

add_footer(slide)

# ============================================================
# SLIDE 4: KEY METRICS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "Key Metrics & Target Outcomes", font_size=28, color=DEEP_BLUE, bold=True)

# Row 1 metrics
metrics_r1 = [("<30", "Days \u2014 Target Onboarding Cycle", EMERALD),
              ("38", "Governed Tasks", BLUE),
              ("7", "Regulatory Frameworks", BLUE),
              ("3", "Orchestrated Pools", BLUE)]
for i, (val, lbl, clr) in enumerate(metrics_r1):
    add_metric_box(slide, Inches(0.8 + i * 3.05), Inches(1.3), Inches(2.8), Inches(1.0), val, lbl, clr)

# Row 2 metrics
metrics_r2 = [("65%", "Evidence Automation Target", GOLD),
              ("4", "Parallel Evaluation Tracks", TEAL),
              ("5", "Message Flows (Cross-Pool)", ROSE),
              ("100%", "Regulatory Traceability", EMERALD)]
for i, (val, lbl, clr) in enumerate(metrics_r2):
    add_metric_box(slide, Inches(0.8 + i * 3.05), Inches(2.5), Inches(2.8), Inches(1.0), val, lbl, clr)

# Cards
add_card(slide, Inches(0.8), Inches(3.8), Inches(5.8), Inches(2.0),
         "Procurement Wins",
         ["\u2022 Vendor landscape assessment with concentration risk analysis",
          "\u2022 Weighted evaluation matrix applied consistently",
          "\u2022 Contract negotiation SLA timers prevent drift",
          "\u2022 Exit cost analysis required before execution"],
         accent_color=GOLD, body_size=10)

add_card(slide, Inches(7), Inches(3.8), Inches(5.5), Inches(2.0),
         "Legal Wins",
         ["\u2022 OCC 2023-17 mandatory provisions checklist embedded",
          "\u2022 GDPR Art. 28(3) DPA provisions automated",
          "\u2022 SEC 17a-4 records retention from day one",
          "\u2022 DORA Art. 30 ICT provider registration built in"],
         accent_color=BLUE, body_size=10)

add_footer(slide)

# ============================================================
# SLIDE 5: 3-POOL ARCHITECTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "3-Pool Architecture", font_size=28, color=DEEP_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.4),
             "The onboarding model orchestrates three independent organizational boundaries connected by 5 message flows.",
             font_size=12, color=GRAY)

# Pool boxes
pools = [
    ("Software Requester", "5 tasks", BLUE, LIGHT_BLUE_BG,
     ["Review existing solutions", "Leverage or proceed", "Gather documentation", "Submit software request"]),
    ("Product Management", "24 tasks", EMERALD, LIGHT_GREEN_BG,
     ["Triage, analysis, pathway routing", "4 parallel evaluations", "Due diligence, PoC, negotiation", "UAT, approval, onboarding"]),
    ("Vendor / Third Party", "9 tasks", GOLD, LIGHT_ORANGE_BG,
     ["Vendor intake & qualification", "Proposal & demo", "Security & compliance review", "Contract review & signing"]),
]

for i, (name, count, color, bg, items) in enumerate(pools):
    left = Inches(0.8 + i * 4.1)
    add_shape(slide, left, Inches(1.6), Inches(3.8), Inches(2.8), fill_color=bg, border_color=color, border_width=Pt(2))
    add_text_box(slide, left, Inches(1.7), Inches(3.8), Inches(0.35), name,
                 font_size=15, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(2.1), Inches(3.8), Inches(0.45), count,
                 font_size=28, color=color, bold=True, alignment=PP_ALIGN.CENTER, font_name="Calibri")
    y = Inches(2.7)
    for item in items:
        add_text_box(slide, left + Inches(0.3), y, Inches(3.2), Inches(0.22),
                     "\u2022 " + item, font_size=9, color=GRAY)
        y += Inches(0.22)

# Message flows table
add_text_box(slide, Inches(0.8), Inches(4.55), Inches(4), Inches(0.3),
             "Cross-Pool Message Flows", font_size=14, color=NAVY, bold=True)

msg_data = [
    ["#", "Message Flow", "From", "To"],
    ["1", "Software Request", "Requester", "Product Mgmt"],
    ["2", "Due Diligence Request", "Product Mgmt", "Vendor"],
    ["3", "Vendor Response", "Vendor", "Product Mgmt"],
    ["4", "Contract Draft", "Product Mgmt", "Vendor"],
    ["5", "Signed Contract", "Vendor", "Product Mgmt"],
]
add_table(slide, Inches(0.8), Inches(4.9), Inches(11.5), msg_data)

add_footer(slide)

# ============================================================
# SLIDE 6: END-TO-END PROCESS FLOW
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "End-to-End Process Flow", font_size=28, color=DEEP_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
             "All 38 tasks across 3 pools with procurement and legal touchpoints highlighted.",
             font_size=12, color=GRAY)

# Phase pipeline
phases = [("1", "Request &\nTriage", RGBColor(0x1D, 0x4E, 0xD8)),
          ("2", "Analysis &\nRouting", BLUE),
          ("3", "Parallel\nEvaluation", RGBColor(0x3B, 0x82, 0xF6)),
          ("4", "Vendor Due\nDiligence", RGBColor(0x0E, 0xA5, 0xE9)),
          ("5", "Contract &\nExecution", RGBColor(0x08, 0x91, 0xB2)),
          ("6", "UAT &\nOnboarding", EMERALD)]

for i, (num, label, color) in enumerate(phases):
    left = Inches(0.8 + i * 2.05)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(1.95), Inches(0.75))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = num + ". " + label.replace("\n", " ")
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

# Flow table
flow_data = [
    ["Stage", "Pool", "Key Tasks", "Procurement", "Legal"],
    ["Request", "Requester", "Review existing, gather docs, submit request", "\u2014", "\u2014"],
    ["Triage", "Product Mgmt", "Initial triage, preliminary analysis, backlog", "\u25cf", "\u2014"],
    ["Routing", "Product Mgmt", "Pathway routing (DMN), Buy vs. Build, vendor landscape", "\u25cf", "\u2014"],
    ["Evaluation", "Product Mgmt", "Tech Architecture, Security, Risk & Compliance, Financial", "\u25cf", "\u25cf"],
    ["Vendor DD", "PM + Vendor", "Due diligence, vendor proposal, evaluation, PoC", "\u25cf", "\u25cf"],
    ["Contract", "PM + Vendor", "Negotiate, vendor review, vendor sign, finalize", "\u25cf", "\u25cf"],
    ["Onboarding", "Product Mgmt", "UAT, final approval, software onboarding, close", "\u25cf", "\u2014"],
]
add_table(slide, Inches(0.8), Inches(2.5), Inches(11.5), flow_data)

add_callout(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.55),
            "Key insight: Procurement and Legal are jointly engaged in 4 of 7 stages. The parallel evaluation stage is where the highest efficiency gains are achieved \u2014 4 tracks run simultaneously instead of sequentially.",
            accent_color=EMERALD)

add_footer(slide)

# ============================================================
# SLIDE 7: PROCUREMENT SECTION DIVIDER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.8),
             "Your Procurement Lifecycle", font_size=38, color=WHITE, bold=True)
# Gold underline
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.35), Inches(4), Inches(0.05))
bar.fill.solid()
bar.fill.fore_color.rgb = GOLD
bar.line.fill.background()
bar.shadow.inherit = False

add_text_box(slide, Inches(0.8), Inches(2.6), Inches(10), Inches(0.6),
             "Automated, governed, and traceable \u2014 from vendor qualification through contract execution.",
             font_size=18, color=RGBColor(0xCC, 0xCC, 0xCC))

# Stats boxes
stats = [("4", "Parallel Evaluation Tracks"),
         ("DMN", "Pathway Routing (Buy/Build)"),
         ("SLA", "Timer-Enforced Contract Review"),
         ("100%", "Vendor Qualification Documented")]
for i, (val, lbl) in enumerate(stats):
    left = Inches(0.8 + i * 3.05)
    box = add_shape(slide, left, Inches(3.8), Inches(2.8), Inches(1.2),
                    fill_color=RGBColor(0x1A, 0x3A, 0x6B))
    add_text_box(slide, left, Inches(3.95), Inches(2.8), Inches(0.5),
                 val, font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(4.45), Inches(2.8), Inches(0.4),
                 lbl, font_size=10, color=RGBColor(0xBB, 0xBB, 0xBB), alignment=PP_ALIGN.CENTER)

add_footer(slide, dark=True)

# ============================================================
# SLIDE 8: VENDOR QUALIFICATION & INTAKE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "Vendor Qualification & Intake", font_size=28, color=DEEP_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
             "Standardized intake with automated screening ensures only qualified vendors enter the evaluation pipeline.",
             font_size=12, color=GRAY)

# Left - Vendor Pool Tasks
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.3),
             "Vendor Pool Tasks", font_size=15, color=NAVY, bold=True)

vendor_tasks = [
    ["Task", "Regulatory Controls"],
    ["Vendor Intake\nSanctions screening, conflict-of-interest", "OCC, NIST"],
    ["Vendor Qualification Gate\nPass/fail against minimum criteria", "OCC, DORA"],
    ["Vendor Proposal\nStructured response to evaluation reqs", "OCC"],
]
add_table(slide, Inches(0.8), Inches(1.9), Inches(5.5), vendor_tasks)

add_callout(slide, Inches(0.8), Inches(3.6), Inches(5.5), Inches(0.65),
            "Procurement: Vendor intake now includes automated sanctions screening, financial stability checks, and concentration risk analysis. Disqualified vendors exit early.",
            accent_color=GOLD)

# Right - PM Tasks
add_text_box(slide, Inches(7), Inches(1.5), Inches(5), Inches(0.3),
             "Product Management Tasks", font_size=15, color=NAVY, bold=True)

pm_tasks = [
    ["Task", "Regulatory Controls"],
    ["Initial Triage\nFeasibility, strategic alignment, duplicates", "OCC, SOX"],
    ["Assess Vendor Landscape\nMarket research, shortlisting", "OCC, DORA"],
    ["Pathway Routing (DMN)\nAutomated Buy vs. Build determination", "NIST"],
]
add_table(slide, Inches(7), Inches(1.9), Inches(5.5), pm_tasks)

add_callout(slide, Inches(7), Inches(3.6), Inches(5.5), Inches(0.65),
            "Legal: Sole-source justification required when < 3 vendors evaluated. OCC 2023-17 \u00a735 mandates documented, objective selection criteria.",
            accent_color=BLUE)

add_footer(slide)

# ============================================================
# SLIDE 9: PARALLEL EVALUATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "Parallel Evaluation \u2014 The Efficiency Multiplier", font_size=28, color=DEEP_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
             "Four evaluation tracks run simultaneously, collapsing what was sequential into parallel execution.",
             font_size=12, color=GRAY)

eval_cards = [
    ("Technical Architecture Review", "Scalability, integration compatibility, enterprise standards, DR architecture.",
     ["NIST PR.PS", "DORA Art. 9", "OCC \u00a750"], BLUE),
    ("Security Assessment", "Vulnerability posture, threat modeling, encryption, access control, incident response.",
     ["OCC \u00a750-55", "NIST CSF", "ISO 27001", "DORA Art. 9-11"], ROSE),
    ("Risk & Compliance Evaluation", "GDPR/CCPA, SOX controls, OCC TPRM, DORA resilience, cross-border transfer.",
     ["OCC \u00a735-40", "GDPR Art. 28,35", "SOX \u00a7302,404", "DORA Art. 5-7"], GOLD),
    ("Financial Analysis", "TCO (3yr/5yr), ROI, NPV, exit cost analysis, market pricing benchmarks.",
     ["SOX \u00a7302,404", "OCC \u00a740"], EMERALD),
]

for i, (title, desc, regs, color) in enumerate(eval_cards):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 6.15)
    top = Inches(1.5 + row * 1.9)
    add_card(slide, left, top, Inches(5.9), Inches(1.7), title,
             [desc, "Regulations: " + ", ".join(regs)],
             accent_color=color, body_size=10)

# Bottom metrics
add_text_box(slide, Inches(0.8), Inches(5.4), Inches(10), Inches(0.3),
             "On the vendor side, 3 parallel tracks run simultaneously: Security Questionnaire, Compliance Documentation, and Technical Demo.",
             font_size=10, color=GRAY)

metrics = [("4+3", "Parallel Tracks (PM + Vendor)", ROSE),
           ("~70%", "Time Saved vs. Sequential", EMERALD),
           ("Weighted", "Evaluation Scoring Matrix", BLUE)]
for i, (val, lbl, clr) in enumerate(metrics):
    add_metric_box(slide, Inches(0.8 + i * 4.1), Inches(5.8), Inches(3.8), Inches(0.95), val, lbl, clr)

add_footer(slide)

# ============================================================
# SLIDE 10: CONTRACT NEGOTIATION & REVIEW
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "Contract Negotiation & Review", font_size=28, color=DEEP_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
             "Dual-track contract lifecycle between Product Management and Vendor with mandatory provisions enforcement.",
             font_size=12, color=GRAY)

# Left: Internal
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.3),
             "Internal (Product Management)", font_size=15, color=NAVY, bold=True)

int_tasks = [
    ["Task", "Description"],
    ["Negotiate Contract", "Draft terms; enforce mandatory provisions per OCC 2023-17"],
    ["Finalize Contract", "Incorporate signed agreement; verify signatory; archive per SEC 17a-4"],
]
add_table(slide, Inches(0.8), Inches(1.9), Inches(5.5), int_tasks)

# OCC Mandatory Provisions
add_text_box(slide, Inches(0.8), Inches(3.0), Inches(5), Inches(0.3),
             "OCC 2023-17 Mandatory Provisions", font_size=13, color=NAVY, bold=True)

provisions = [
    "Performance benchmarks and remedies",
    "Right to audit and examination",
    "Subcontracting limitations",
    "Business continuity and DR requirements",
    "Data protection and confidentiality",
    "Regulatory compliance obligations",
    "Default and termination clauses",
    "Dispute resolution mechanisms",
    "Insurance requirements",
    "Indemnification provisions"
]
y = Inches(3.3)
for p in provisions:
    add_text_box(slide, Inches(1.0), y, Inches(5), Inches(0.2),
                 "\u2713 " + p, font_size=9, color=EMERALD)
    y += Inches(0.2)

# Right: External
add_text_box(slide, Inches(7), Inches(1.5), Inches(5), Inches(0.3),
             "External (Vendor Pool)", font_size=15, color=NAVY, bold=True)

ext_tasks = [
    ["Task", "Description"],
    ["Vendor Contract Review", "Vendor reviews draft; redlines; counter-proposals tracked"],
    ["Vendor Contract Sign", "Authorized signatory executes; digital signature with timestamp"],
]
add_table(slide, Inches(7), Inches(1.9), Inches(5.5), ext_tasks)

add_callout(slide, Inches(7), Inches(3.1), Inches(5.5), Inches(0.8),
            "Legal: Every contract negotiation task enforces a mandatory provisions checklist derived from OCC 2023-17 \u00a760-75. Missing provisions trigger a compliance gate failure.",
            accent_color=BLUE)

add_callout(slide, Inches(7), Inches(4.1), Inches(5.5), Inches(0.7),
            "Procurement: SLA timer on contract negotiation prevents indefinite redline cycles. Vendor response deadlines enforced by BPMN boundary timer events.",
            accent_color=GOLD)

add_footer(slide)

# ============================================================
# SLIDE 11: CONTRACT EXECUTION & COMPLIANCE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "Contract Execution & Compliance", font_size=28, color=DEEP_BLUE, bold=True)

# Left: Execution Controls
add_text_box(slide, Inches(0.8), Inches(1.2), Inches(5), Inches(0.3),
             "Execution Controls", font_size=15, color=NAVY, bold=True)

exec_cards = [
    ("Authorized Signatory Verification", "Digital signature validation against delegated authority matrix. Only pre-approved signatories can execute.", BLUE, "SOX \u00a7302, OCC \u00a760"),
    ("SEC 17a-4 Records Retention", "All contract artifacts archived in tamper-evident, WORM-compliant storage with minimum 6-year retention.", GOLD, "SEC 17a-4, SOX"),
    ("DORA Art. 30 Registration", "ICT third-party service provider registered per DORA. Critical/important functions flagged for enhanced oversight.", TEAL, "DORA Art. 28-30"),
]

y = Inches(1.6)
for title, desc, color, regs in exec_cards:
    add_card(slide, Inches(0.8), y, Inches(5.5), Inches(1.2), title, [desc, regs], accent_color=color, body_size=9)
    y += Inches(1.35)

# Right: DPA & Post-execution
add_text_box(slide, Inches(7), Inches(1.2), Inches(5), Inches(0.3),
             "DPA Execution Flow", font_size=15, color=NAVY, bold=True)

dpa_steps = ["Draft DPA", "Vendor Review", "Execute DPA", "Archive"]
for i, step in enumerate(dpa_steps):
    left = Inches(7 + i * 1.4)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.65), Inches(1.2), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE if step == "Execute DPA" else WHITE
    shape.line.color.rgb = BLUE
    shape.line.width = Pt(2)
    shape.shadow.inherit = False
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = step
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE if step == "Execute DPA" else BLUE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    if i < len(dpa_steps) - 1:
        add_text_box(slide, left + Inches(1.2), Inches(1.72), Inches(0.2), Inches(0.3),
                     "\u2192", font_size=14, color=BLUE, bold=True)

add_text_box(slide, Inches(7), Inches(2.4), Inches(5), Inches(0.3),
             "Post-Execution Compliance", font_size=15, color=NAVY, bold=True)

compliance_data = [
    ["Control", "Framework", "Frequency"],
    ["Contract compliance review", "OCC 2023-17", "Annual"],
    ["Performance monitoring", "SLA Framework", "Continuous"],
    ["DPA adequacy review", "GDPR Art. 28", "Annual"],
    ["Exit plan validation", "DORA Art. 28", "Annual"],
    ["Records retention audit", "SEC 17a-4", "Semi-annual"],
]
add_table(slide, Inches(7), Inches(2.8), Inches(5.5), compliance_data)

add_footer(slide)

# ============================================================
# SLIDE 12: LEGAL SECTION DIVIDER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.8),
             "Regulatory Compliance Built Into Every Step", font_size=36, color=WHITE, bold=True)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.35), Inches(5), Inches(0.05))
bar.fill.solid()
bar.fill.fore_color.rgb = GOLD
bar.line.fill.background()
bar.shadow.inherit = False

add_text_box(slide, Inches(0.8), Inches(2.6), Inches(10), Inches(0.6),
             "Not an afterthought. Not a separate workstream. Compliance woven into the fabric of every task.",
             font_size=18, color=RGBColor(0xCC, 0xCC, 0xCC))

stats = [("7", "Regulatory Frameworks Mapped"),
         ("38", "Tasks with Regulatory Annotations"),
         ("Art. 28(3)", "GDPR DPA Provisions Embedded"),
         ("6yr+", "SEC 17a-4 Retention")]
for i, (val, lbl) in enumerate(stats):
    left = Inches(0.8 + i * 3.05)
    box = add_shape(slide, left, Inches(3.8), Inches(2.8), Inches(1.2),
                    fill_color=RGBColor(0x1A, 0x3A, 0x6B))
    add_text_box(slide, left, Inches(3.95), Inches(2.8), Inches(0.5),
                 val, font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(4.45), Inches(2.8), Inches(0.4),
                 lbl, font_size=10, color=RGBColor(0xBB, 0xBB, 0xBB), alignment=PP_ALIGN.CENTER)

add_footer(slide, dark=True)

# ============================================================
# SLIDE 13: REGULATORY FRAMEWORK COVERAGE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "Regulatory Framework Coverage Matrix", font_size=28, color=DEEP_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
             "Every task in the onboarding model is mapped to its applicable regulatory controls.",
             font_size=12, color=GRAY)

# Heatmap as table
tasks_short = ["Review\nExisting", "Gather\nDocs", "Initial\nTriage", "Prelim\nAnalysis",
               "Vendor\nLandscpe", "Tech Arch", "Security\nAssess", "Risk &\nCompl",
               "Financial\nAnalysis", "Vendor\nDD", "Negotiate\nContract", "Finalize\nContract",
               "Final\nApproval", "Onboard\nSoftware"]
regs = ["OCC 2023-17", "GDPR/CCPA", "SOX", "DORA", "NIST CSF", "ISO 27001", "SEC 17a-4"]

# Build heatmap data: 2=primary, 1=secondary, 0=N/A
heat = [
    [1,0,1,1,1,1,1,1,1,1,2,2,1,0],  # OCC
    [0,1,0,1,0,0,0,2,0,1,2,1,0,0],  # GDPR
    [1,1,1,1,0,0,0,1,2,0,1,1,1,0],  # SOX
    [0,0,0,0,1,1,1,1,0,1,1,1,0,0],  # DORA
    [1,0,1,1,1,2,2,0,0,1,0,0,0,1],  # NIST
    [0,0,0,0,0,0,2,0,0,0,0,0,0,0],  # ISO
    [0,0,0,0,0,0,0,0,0,0,1,2,0,0],  # SEC
]

header = [""] + [t.replace("\n", " ") for t in tasks_short]
rows = [header]
symbols = {2: "\u25a0", 1: "\u25a1", 0: "\u2014"}
for r_idx, reg in enumerate(regs):
    row = [reg] + [symbols[v] for v in heat[r_idx]]
    rows.append(row)

tbl_shape = add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), rows)
# Color the cells
tbl = tbl_shape.table
for r_idx in range(1, len(rows)):
    for c_idx in range(1, len(rows[0])):
        val = heat[r_idx - 1][c_idx - 1]
        cell = tbl.cell(r_idx, c_idx)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
        if val == 2:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x16, 0x65, 0x34)
            for p in cell.text_frame.paragraphs:
                p.font.color.rgb = WHITE
        elif val == 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x65, 0xA3, 0x0D)
            for p in cell.text_frame.paragraphs:
                p.font.color.rgb = WHITE

# Legend
add_text_box(slide, Inches(0.8), Inches(5.6), Inches(10), Inches(0.3),
             "\u25a0 Primary     \u25a1 Secondary     \u2014 Not Applicable", font_size=10, color=GRAY)

add_callout(slide, Inches(0.8), Inches(5.95), Inches(11.5), Inches(0.5),
            "Legal: This matrix provides exam-ready evidence mapping. When regulators ask which OCC 2023-17 sections are addressed, the answer is traceable to specific BPMN tasks with documented evidence artifacts.",
            accent_color=BLUE)

add_footer(slide)

# ============================================================
# SLIDE 14: DPA LIFECYCLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "DPA Lifecycle \u2014 GDPR Article 28(3)", font_size=28, color=DEEP_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
             "End-to-end Data Processing Agreement journey embedded in the onboarding workflow.",
             font_size=12, color=GRAY)

# DPA flow
dpa_steps = ["Identify Data\nProcessing", "DPIA\nScreening", "Draft\nDPA",
             "Cross-Border\nAssessment", "Vendor\nReview", "Execute\n& Archive"]
for i, step in enumerate(dpa_steps):
    left = Inches(0.8 + i * 2.05)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(1.7), Inches(0.65))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = step.replace("\n", " ")
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    if i < len(dpa_steps) - 1:
        add_text_box(slide, left + Inches(1.7), Inches(1.6), Inches(0.35), Inches(0.4),
                     "\u2192", font_size=18, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# Left: GDPR provisions checklist
add_text_box(slide, Inches(0.8), Inches(2.4), Inches(5.5), Inches(0.3),
             "GDPR Art. 28(3) Mandatory Provisions", font_size=14, color=NAVY, bold=True)

provisions = [
    "Process data only on documented instructions",
    "Ensure confidentiality obligations for personnel",
    "Implement appropriate technical & organizational measures",
    "Sub-processor authorization and notification",
    "Assist controller with data subject requests",
    "Assist with security, breach notification, DPIA",
    "Delete or return data after processing ends",
    "Make available all info to demonstrate compliance"
]
y = Inches(2.8)
for p in provisions:
    add_text_box(slide, Inches(1.0), y, Inches(5), Inches(0.2),
                 "\u2713 " + p, font_size=9, color=EMERALD)
    y += Inches(0.22)

# Right: Enforcement points
add_text_box(slide, Inches(7), Inches(2.4), Inches(5.5), Inches(0.3),
             "Where DPA Provisions Are Enforced", font_size=14, color=NAVY, bold=True)

dpa_enforcement = [
    ["Provision", "Enforced At"],
    ["Processing scope", "Gather Documentation (Requester)"],
    ["DPIA screening", "Risk & Compliance Evaluation"],
    ["Cross-border transfer", "Risk & Compliance Evaluation"],
    ["Sub-processor mgmt", "Vendor Due Diligence"],
    ["Technical measures", "Security Assessment"],
    ["Breach notification", "Contract Negotiation"],
    ["Data return/deletion", "Contract Negotiation"],
    ["Audit access", "Contract Negotiation"],
]
add_table(slide, Inches(7), Inches(2.8), Inches(5.5), dpa_enforcement)

add_footer(slide)

# ============================================================
# SLIDE 15: AUDIT RIGHTS & EVIDENCE CHAIN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "Audit Rights & Evidence Chain", font_size=28, color=DEEP_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
             "Complete, tamper-evident audit trail from initial request through close-out.",
             font_size=12, color=GRAY)

# Left: Evidence collection table
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.3),
             "Evidence Collection Points", font_size=14, color=NAVY, bold=True)

evidence_data = [
    ["Stage", "Evidence Artifacts", "Retention"],
    ["Request", "Intake form, portfolio analysis, sign-off", "6 yr"],
    ["Triage", "Triage scorecard, duplicate detection", "6 yr"],
    ["Analysis", "TCO model, DPIA screening, pathway rationale", "6 yr"],
    ["Evaluation", "Arch review, security scan, compliance gaps", "6 yr"],
    ["Due Diligence", "DD report, financial stability, sanctions", "6 yr"],
    ["Contract", "All drafts, redlines, executed agreement, DPA", "6 yr+"],
    ["Onboarding", "UAT results, approval decision, deployment", "6 yr"],
    ["Close-out", "Lessons learned, performance baselines, KRIs", "6 yr"],
]
add_table(slide, Inches(0.8), Inches(1.9), Inches(5.8), evidence_data)

# Right: Audit rights cards
add_text_box(slide, Inches(7), Inches(1.5), Inches(5), Inches(0.3),
             "Audit Access Rights", font_size=14, color=NAVY, bold=True)

add_card(slide, Inches(7), Inches(1.9), Inches(5.5), Inches(1.0),
         "OCC 2023-17 \u00a765: Right to Audit",
         ["Contracts must include provisions allowing the institution", "and its regulators to access vendor's books and records."],
         accent_color=BLUE, body_size=9)

add_card(slide, Inches(7), Inches(3.1), Inches(5.5), Inches(1.0),
         "DORA Art. 28(7): Regulatory Access",
         ["Financial entities must ensure ICT contracts enable", "competent authorities to exercise audit rights."],
         accent_color=GOLD, body_size=9)

add_card(slide, Inches(7), Inches(4.3), Inches(5.5), Inches(1.0),
         "Evidence Integrity",
         ["All artifacts timestamped, version-controlled, stored in", "immutable audit log. SEC 17a-4 / SOX \u00a7802 compliant."],
         accent_color=TEAL, body_size=9)

add_footer(slide)

# ============================================================
# SLIDE 16: RISK-TIERED GOVERNANCE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "Risk-Tiered Governance", font_size=28, color=DEEP_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
             "DMN-driven risk classification determines due diligence depth, monitoring cadence, and governance oversight.",
             font_size=12, color=GRAY)

tier_data = [
    ["Risk Tier", "DD Depth", "Monitoring", "Review Authority", "Contract Provisions"],
    ["Critical", "Enhanced (all)", "Continuous + quarterly", "Advisory Board + Regulators", "Full mandatory + enhanced audit"],
    ["High", "Comprehensive", "Monthly + quarterly", "Governance Committee", "Full mandatory provisions"],
    ["Limited", "Standard", "Quarterly", "Department Head", "Standard provisions"],
    ["Minimal", "Streamlined", "Semi-annual", "Manager sign-off", "Simplified provisions"],
]
tbl_shape = add_table(slide, Inches(0.8), Inches(1.5), Inches(11.5), tier_data)

# Color risk tier cells
tbl = tbl_shape.table
tier_colors = [ROSE, GOLD, BLUE, EMERALD]
for i, color in enumerate(tier_colors):
    cell = tbl.cell(i + 1, 0)
    cell.fill.solid()
    cell.fill.fore_color.rgb = color
    for p in cell.text_frame.paragraphs:
        p.font.color.rgb = WHITE
        p.font.bold = True

# Two columns: inputs and outputs
add_text_box(slide, Inches(0.8), Inches(3.7), Inches(5.5), Inches(0.3),
             "DMN Risk Tier Inputs", font_size=14, color=NAVY, bold=True)
inputs = ["Data sensitivity classification (PII, PHI, financial)",
          "Operational criticality assessment",
          "Vendor concentration risk score",
          "Geographic / jurisdictional exposure",
          "Regulatory framework applicability count",
          "Fourth-party dependency depth"]
y = Inches(4.05)
for item in inputs:
    add_text_box(slide, Inches(1.0), y, Inches(5.3), Inches(0.2),
                 "\u2022 " + item, font_size=10, color=DARK_TEXT)
    y += Inches(0.22)

add_text_box(slide, Inches(7), Inches(3.7), Inches(5.5), Inches(0.3),
             "DMN Monitoring Cadence Outputs", font_size=14, color=NAVY, bold=True)
outputs = ["KRI monitoring frequency",
           "Performance review schedule",
           "Compliance assessment schedule",
           "Financial stability review schedule",
           "Escalation thresholds by metric",
           "Board reporting requirements"]
y = Inches(4.05)
for item in outputs:
    add_text_box(slide, Inches(7.2), y, Inches(5.3), Inches(0.2),
                 "\u2022 " + item, font_size=10, color=DARK_TEXT)
    y += Inches(0.22)

add_callout(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.5),
            "Procurement + Legal: Risk tier drives the contract clause library. Critical-tier vendors require enhanced audit rights, sub-contractor approval requirements, and DORA-specific ICT resilience testing provisions.",
            accent_color=EMERALD)

add_footer(slide)

# ============================================================
# SLIDE 17: AUTOMATION SECTION DIVIDER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.8),
             "From 90 Days to Under 30", font_size=38, color=WHITE, bold=True)
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.35), Inches(3.5), Inches(0.05))
bar.fill.solid()
bar.fill.fore_color.rgb = GOLD
bar.line.fill.background()
bar.shadow.inherit = False

add_text_box(slide, Inches(0.8), Inches(2.6), Inches(10), Inches(0.6),
             "Automation, parallelization, and standardization eliminate the bottlenecks that define today\u2019s vendor onboarding.",
             font_size=18, color=RGBColor(0xCC, 0xCC, 0xCC))

stats = [("90 \u2192 30", "Day Cycle Reduction Target"),
         ("65%", "Evidence Automation"),
         ("7", "Parallel Evaluation Tracks"),
         ("SLA", "Timer-Enforced at Every Stage")]
for i, (val, lbl) in enumerate(stats):
    left = Inches(0.8 + i * 3.05)
    box = add_shape(slide, left, Inches(3.8), Inches(2.8), Inches(1.2),
                    fill_color=RGBColor(0x1A, 0x3A, 0x6B))
    add_text_box(slide, left, Inches(3.95), Inches(2.8), Inches(0.5),
                 val, font_size=26, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(4.45), Inches(2.8), Inches(0.4),
                 lbl, font_size=10, color=RGBColor(0xBB, 0xBB, 0xBB), alignment=PP_ALIGN.CENTER)

add_footer(slide, dark=True)

# ============================================================
# SLIDE 18: BOTTLENECK ELIMINATION (Chart)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "Bottleneck Elimination", font_size=28, color=DEEP_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
             "Before/after comparison across procurement and legal touchpoints.",
             font_size=12, color=GRAY)

# Chart
chart_data = CategoryChartData()
chart_data.categories = ['Vendor Evidence\nCollection', 'Contract Review\nCycle',
                         'Committee\nScheduling', 'Compliance Gap\nAnalysis',
                         'DPA\nNegotiation', 'Risk Tier\nClassification']
chart_data.add_series('Before (days)', (15, 20, 12, 10, 8, 5))
chart_data.add_series('After (days)', (4, 7, 5, 3, 3, 1))

chart_shape = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.8), Inches(1.5),
    Inches(11.5), Inches(3.8), chart_data
)
chart = chart_shape.chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(10)

# Style series
plot = chart.plots[0]
plot.gap_width = 80
series_before = plot.series[0]
series_after = plot.series[1]
series_before.format.fill.solid()
series_before.format.fill.fore_color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
series_after.format.fill.solid()
series_after.format.fill.fore_color.rgb = BLUE

chart.value_axis.visible = True
chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
chart.category_axis.tick_labels.font.size = Pt(9)
chart.value_axis.tick_labels.font.size = Pt(9)

add_callout(slide, Inches(0.8), Inches(5.5), Inches(5.5), Inches(0.55),
            "Procurement: Vendor evidence collection drops from 15+ days to ~4 days through automated intake and parallel vendor-side tracks.",
            accent_color=GOLD)
add_callout(slide, Inches(6.8), Inches(5.5), Inches(5.5), Inches(0.55),
            "Legal: Contract review cycle collapses from 20+ days to ~7 days with mandatory provisions templates and SLA-enforced deadlines.",
            accent_color=BLUE)

add_footer(slide)

# ============================================================
# SLIDE 19: EVIDENCE AUTOMATION MATRIX
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "Evidence Automation \u2014 Human Judgment vs. Automation", font_size=28, color=DEEP_BLUE, bold=True)

auto_data = [
    ["Task", "Automation Tier", "Procurement", "Legal", "Rationale"],
    ["Sanctions Screening", "Fully Auto", "\u25cf", "\u25cf", "OFAC/sanctions list lookup"],
    ["Concentration Risk Analysis", "Fully Auto", "\u25cf", "\u2014", "Portfolio database query"],
    ["Pathway Routing (DMN)", "Fully Auto", "\u25cf", "\u2014", "DMN decision table"],
    ["Financial Stability Check", "Partially Auto", "\u25cf", "\u2014", "Auto-pull financials, human review"],
    ["Security Assessment", "Partially Auto", "\u2014", "\u25cf", "Auto-scan + expert review"],
    ["Risk & Compliance Eval", "Partially Auto", "\u2014", "\u25cf", "Auto-map regs, human gap analysis"],
    ["Vendor Evaluation Scoring", "Partially Auto", "\u25cf", "\u2014", "Weighted matrix, human override"],
    ["Contract Negotiation", "Human-Led", "\u25cf", "\u25cf", "Strategic negotiation, legal judgment"],
    ["DPA Drafting & Review", "Human-Led", "\u2014", "\u25cf", "GDPR Art. 28(3) requires legal judgment"],
    ["Governance Committee", "Human-Led", "\u25cf", "\u25cf", "Strategic risk acceptance"],
    ["Signatory Verification", "Human-Led", "\u2014", "\u25cf", "Delegated authority matrix"],
]
tbl_shape = add_table(slide, Inches(0.5), Inches(1.2), Inches(12.3), auto_data)

# Color automation tier cells
tbl = tbl_shape.table
tier_map = {"Fully Auto": EMERALD, "Partially Auto": BLUE, "Human-Led": GOLD}
for r in range(1, len(auto_data)):
    tier = auto_data[r][1]
    cell = tbl.cell(r, 1)
    cell.fill.solid()
    cell.fill.fore_color.rgb = tier_map.get(tier, GRAY)
    for p in cell.text_frame.paragraphs:
        p.font.color.rgb = WHITE
        p.font.bold = True
    # Center procurement and legal columns
    for c in [2, 3]:
        for p in tbl.cell(r, c).text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER

# Metrics
metrics = [("~35%", "Fully Automatable", EMERALD),
           ("~30%", "Partially Automatable", BLUE),
           ("~35%", "Human Judgment Required", GOLD)]
for i, (val, lbl, clr) in enumerate(metrics):
    add_metric_box(slide, Inches(0.8 + i * 4.1), Inches(5.8), Inches(3.8), Inches(0.95), val, lbl, clr)

add_footer(slide)

# ============================================================
# SLIDE 20: ONGOING MONITORING & CLOSE-OUT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "Ongoing Monitoring & Close-out", font_size=28, color=DEEP_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
             "Post-onboarding governance ensures vendor performance is tracked, risks are monitored, and knowledge is captured.",
             font_size=12, color=GRAY)

# Left: Post-onboarding tasks
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.3),
             "Post-Onboarding Tasks", font_size=14, color=NAVY, bold=True)

post_data = [
    ["Task", "Description"],
    ["Software Onboarding", "Deploy to production, configure monitoring, establish baselines"],
    ["Performance Baselines", "Set KRI thresholds based on vendor risk tier and contract SLAs"],
    ["Close Request", "Archive all evidence, lessons learned, update vendor registry"],
]
add_table(slide, Inches(0.8), Inches(1.9), Inches(5.5), post_data)

add_callout(slide, Inches(0.8), Inches(3.4), Inches(5.5), Inches(0.6),
            "Procurement: Vendor performance baselines are set during onboarding \u2014 not 6 months later. KRIs calibrated to contract SLAs. Breaches trigger automated escalation.",
            accent_color=GOLD)

# Right: Knowledge capture
add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.3),
             "Knowledge Capture", font_size=14, color=NAVY, bold=True)

add_card(slide, Inches(7), Inches(1.9), Inches(5.5), Inches(1.0),
         "Lessons Learned Archive",
         ["Every onboarding contributes to the organizational knowledge", "base. Patterns, bottlenecks, and strategies captured."],
         accent_color=EMERALD, body_size=9)

add_card(slide, Inches(7), Inches(3.1), Inches(5.5), Inches(1.0),
         "Vendor Registry Update",
         ["Vendor profile enriched with evaluation scores, risk tier,", "contract terms, SLA thresholds, and key contacts."],
         accent_color=BLUE, body_size=9)

add_card(slide, Inches(7), Inches(4.3), Inches(5.5), Inches(1.0),
         "Cumulative Knowledge Effect",
         ["30\u201350% effort reduction on similar vendor assessments", "over time as the knowledge base accumulates."],
         accent_color=GOLD, body_size=9)

add_footer(slide)

# ============================================================
# SLIDE 21: CALL TO ACTION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "Next Steps", font_size=44, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.5),
             "Implementation Roadmap", font_size=22, color=WHITE, bold=True)

# Three phase boxes
phases = [
    ("Phase 1", "Weeks 1\u20134", TEAL,
     ["Validate BPMN model with stakeholders",
      "Map existing processes to ideal state",
      "Identify gap remediation priorities",
      "Define regulatory annotation requirements"]),
    ("Phase 2", "Weeks 5\u201310", GOLD,
     ["Implement contract clause library",
      "Deploy DMN pathway routing",
      "Configure SLA timers and escalations",
      "Build evidence automation pipelines"]),
    ("Phase 3", "Weeks 11\u201316", EMERALD,
     ["Pilot with 3\u20135 vendor onboardings",
      "Measure cycle time reduction",
      "Validate regulatory exam readiness",
      "Scale to full production"]),
]

for i, (phase, timeline, color, items) in enumerate(phases):
    left = Inches(0.8 + i * 4.1)
    box = add_shape(slide, left, Inches(2.5), Inches(3.8), Inches(3.0),
                    fill_color=RGBColor(0x1A, 0x3A, 0x6B))
    add_text_box(slide, left + Inches(0.2), Inches(2.65), Inches(3.4), Inches(0.35),
                 phase, font_size=20, color=color, bold=True)
    add_text_box(slide, left + Inches(0.2), Inches(3.0), Inches(3.4), Inches(0.3),
                 timeline, font_size=14, color=WHITE, bold=True)
    y = Inches(3.45)
    for item in items:
        add_text_box(slide, left + Inches(0.3), y, Inches(3.3), Inches(0.22),
                     "\u2022 " + item, font_size=10, color=RGBColor(0xCC, 0xCC, 0xCC))
        y += Inches(0.25)

# CTA box
cta_box = add_shape(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.8),
                    fill_color=RGBColor(0x1A, 0x3A, 0x6B))
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.8), Inches(0.05), Inches(0.8))
bar.fill.solid()
bar.fill.fore_color.rgb = GOLD
bar.line.fill.background()
bar.shadow.inherit = False

add_text_box(slide, Inches(1.1), Inches(5.9), Inches(11), Inches(0.6),
             "Immediate ask: Schedule 2-hour working session with Procurement and Legal leads to walk through "
             "the detailed BPMN model in Camunda Modeler and validate task-level regulatory annotations.",
             font_size=12, color=WHITE)

add_footer(slide, dark=True)

# ============================================================
# SAVE
# ============================================================
output_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                           "docs", "presentations", "onboarding",
                           "Vendor-Onboarding-Governance.pptx")
prs.save(output_path)
print(f"Saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
