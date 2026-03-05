#!/usr/bin/env python3
"""Generate Vendor / Third Party swimlane task documentation as PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import re

# Brand colors
DARK_BLUE = RGBColor(0x00, 0x33, 0x8D)
MED_BLUE = RGBColor(0x00, 0x53, 0xAE)
LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x2D, 0x2D, 0x2D)
GRAY = RGBColor(0x5F, 0x63, 0x68)
ACCENT = RGBColor(0x00, 0x7B, 0x5F)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
BORDER_GRAY = RGBColor(0xDD, 0xDD, 0xDD)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=12, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_header_bar(slide):
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), fill_color=DARK_BLUE)
    add_shape(slide, Inches(0), Inches(0.9), Inches(13.333), Inches(0.06), fill_color=ACCENT)

def add_footer(slide, slide_num, total):
    add_shape(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), fill_color=LIGHT_GRAY)
    add_text_box(slide, Inches(0.5), Inches(7.15), Inches(6), Inches(0.3),
                 "Vendor / Third Party — External Software or Contracting",
                 font_size=8, color=GRAY)
    add_text_box(slide, Inches(10), Inches(7.15), Inches(3), Inches(0.3),
                 f"Slide {slide_num} of {total}", font_size=8, color=GRAY, alignment=PP_ALIGN.RIGHT)


# ── Task data extracted from BPMN ──

tasks = [
    {
        "num": 1,
        "name": "Vendor Intake and Qualification",
        "id": "Task_VendorIntake",
        "description": "Screen vendor eligibility and establish baseline qualification before any substantive information exchange begins. Collect vendor legal entity details, ownership structure, financial stability indicators, geographic footprint, and existing certifications. Perform conflict-of-interest screening and preliminary risk tier classification. Verify the vendor is not on any sanctions or exclusion lists. Establish the lawful basis for data collection under applicable privacy regulations.",
        "controls": [
            ("OCC 2023-17 Sections 30-35", "Pre-engagement due diligence requirements — verify vendor identity, assess concentration risk, determine if vendor is critical or significant."),
            ("NIST CSF 2.0 GV.OC-05", "Supply chain risk identification at initial onboarding — document risk factors in the vendor risk register."),
            ("GDPR Article 6 / CCPA", "Establish lawful basis before collecting any personal data from or about the vendor."),
        ],
        "evidence": [
            "Completed vendor qualification form with legal entity name, DUNS number, jurisdiction of incorporation",
            "Risk tier pre-assessment worksheet (feeds into DMN_RiskTierClassification)",
            "Signed conflict-of-interest declaration from the requesting business unit",
            "Sanctions screening result (OFAC, EU consolidated list)",
            "Vendor financial stability indicators (credit rating, annual revenue, years in operation)",
            "Proof of insurance certificates (E&O, cyber liability)",
        ],
        "data_points": "Legal entity name, tax ID, DUNS, ownership structure (UBO disclosure for >25% holders), primary contact, geographic presence, industry certifications held, number of employees, annual revenue range, existing client references in financial services.",
    },
    {
        "num": 2,
        "name": "Vendor Proposal",
        "id": "Task_VendorProposal",
        "description": "Vendor prepares and submits a formal commercial and technical proposal in response to the organization's requirements. The proposal must address functional requirements, technical architecture, SLA commitments, pricing model, implementation timeline, and support structure. Vendor should disclose financial stability information, sub-contractor dependencies, and data handling practices.",
        "controls": [
            ("OCC 2023-17 Section 40", "RFP/RFI requirements — vendor must demonstrate capability, financial stability, and willingness to comply with regulatory examination requirements."),
            ("SOX Section 404", "Proposed financial controls and reporting capabilities must be documented to ensure auditability of financial data flows."),
            ("DORA Article 28", "ICT third-party provider must assert resilience capabilities, including RTO, availability targets, and incident notification commitments."),
        ],
        "evidence": [
            "Signed proposal document with commercial terms, technical specifications, and SLA commitments",
            "Vendor capability matrix mapping proposed features to stated requirements",
            "Financial stability disclosures (audited financials, credit references, or parent company guarantee)",
            "Sub-contractor dependency disclosure listing all fourth parties involved in service delivery",
            "Data handling and residency commitments",
            "Implementation timeline with milestones and resource commitments",
        ],
        "data_points": "Proposed pricing model, TCO estimate over 3-5 years, SLA targets (uptime, response time, resolution time), proposed technology stack, integration approach, data residency locations, sub-processor list, key personnel assigned, reference customers in regulated industries.",
    },
    {
        "num": 3,
        "name": "Security Questionnaire",
        "id": "Task_VendorSecurityReview",
        "description": "Vendor completes a standardized security questionnaire (HECVAT, SIG, or equivalent) covering information security program maturity, access controls, encryption practices, incident response capabilities, vulnerability management, and business continuity planning. Vendor provides supporting evidence including SOC 2 Type II reports, penetration test summaries, and vulnerability disclosure policies.",
        "controls": [
            ("OCC 2023-17 Sections 50-55", "Information security program assessment — evaluate vendor access controls, data handling, encryption, and incident response readiness."),
            ("NIST CSF 2.0 PR.AA / PR.DS / DE.CM", "Assess authentication, encryption standards, key management, data classification, logging, monitoring, and alerting capabilities."),
            ("ISO 27001 Annex A", "Map vendor controls to Annex A domains — A.9 Access Control, A.10 Cryptography, A.12 Operations Security, A.16 Incident Management, A.17 Business Continuity."),
            ("DORA Articles 9, 28", "ICT security requirements — verify vendor incident reporting capability within 24-hour notification windows."),
        ],
        "evidence": [
            "Completed HECVAT or SIG security questionnaire with all sections addressed",
            "SOC 2 Type II report (current year) or SOC 1 if financial data processing",
            "Most recent penetration test executive summary (within 12 months)",
            "Vulnerability disclosure and patch management policy",
            "Business continuity and disaster recovery plan summary with RTO/RPO commitments",
            "Data encryption standards documentation (at rest, in transit, key management)",
            "Security incident history disclosure (breaches in past 3 years)",
        ],
        "data_points": "Security certifications held (ISO 27001, SOC 2, PCI DSS, FedRAMP), encryption algorithms and key lengths, MFA enforcement %, MTTP critical vulnerabilities, incident response SLA, backup frequency, DR test frequency, security awareness training completion rate.",
    },
    {
        "num": 4,
        "name": "Compliance Documentation",
        "id": "Task_VendorComplianceReview",
        "description": "Vendor submits regulatory certifications, compliance attestations, and data protection documentation. This includes proof of regulatory compliance, audit rights confirmation, Data Processing Agreement (DPA) terms, sub-processor disclosures, and data residency documentation. Vendor must demonstrate compliance posture across all applicable regulatory frameworks.",
        "controls": [
            ("OCC 2023-17 Section 55", "Vendor must confirm willingness to submit to regulatory examination, provide audit rights, and maintain compliance with applicable banking regulations."),
            ("GDPR Article 28 / CCPA", "DPA requirements covering processing purpose limitation, data minimization, security measures, sub-processor authorization, breach notification, and data return/deletion on termination."),
            ("SOX", "External audit independence confirmation if vendor touches financial reporting data; financial controls attestation."),
            ("DORA Article 30", "Vendor must disclose full sub-contracting chain, confirm exit strategy, and support regulatory reporting."),
            ("NIST CSF 2.0 GV.SC-07", "Supplier compliance posture documentation and supply chain risk management practices."),
        ],
        "evidence": [
            "Regulatory certifications (ISO 27001, SOC 2 Type II, PCI DSS AOC, FedRAMP as applicable)",
            "Signed or draft DPA with all GDPR Article 28(3) provisions addressed",
            "Complete sub-processor list with names, locations, and services provided",
            "Audit rights confirmation letter",
            "Data residency documentation",
            "Privacy impact assessment or DPIA if high-risk processing",
            "Regulatory examination history disclosure (past 5 years)",
        ],
        "data_points": "Certifications held and expiry dates, DPA execution status, sub-processor count and jurisdictions, data residency countries, audit rights scope, regulatory examination history, privacy officer contact, breach notification timeline, data retention and deletion policies.",
    },
    {
        "num": 5,
        "name": "Technical Demonstration",
        "id": "Task_VendorTechDemo",
        "description": "Vendor conducts a structured technical demonstration of the proposed solution against predefined evaluation criteria. The demo must cover functional capabilities, integration points, security controls in action, performance under load, failover and recovery mechanisms, and administrative workflows. Evaluators score using a standardized scorecard.",
        "controls": [
            ("OCC 2023-17 Section 40", "Technology capability validation — verify the vendor's solution meets technical requirements without introducing unacceptable risk."),
            ("NIST CSF 2.0 PR.PS", "Demonstrate secure configuration management, hardening practices, and change management controls."),
            ("DORA Article 9", "ICT resilience capabilities — demonstrate failover mechanisms, recovery procedures, and availability targets."),
        ],
        "evidence": [
            "Pre-defined demo script with evaluation criteria mapped to requirements",
            "Completed evaluation scorecard signed by all assessors",
            "Architecture diagram reviewed and annotated during the demonstration",
            "RTO/RPO verification results from demonstrated failover",
            "Integration compatibility assessment (API standards, auth methods, data formats)",
            "Performance benchmark results (response times, throughput, concurrent users)",
        ],
        "data_points": "Demo completion % against script, evaluator scores by category, identified gaps or limitations, integration complexity assessment, estimated customization effort, accessibility compliance status (WCAG 2.1 AA).",
    },
    {
        "num": 6,
        "name": "Vendor Contract Review",
        "id": "Task_VendorContractReview",
        "description": "Vendor receives the draft contract and performs a thorough review of all terms, conditions, and obligations. Reviews service level commitments, liability provisions, indemnification clauses, audit rights, data protection obligations, termination and exit provisions, intellectual property rights, and regulatory compliance requirements. Vendor provides redlined comments and negotiation positions.",
        "controls": [
            ("OCC 2023-17 Sections 60-70", "Required contractual provisions — audit rights, subcontracting restrictions, data ownership, BC/DR commitments, termination rights, confidentiality, and regulatory compliance representations."),
            ("SOX Section 302/404", "Audit rights and financial controls provisions must be explicitly stated."),
            ("GDPR Article 28(3) / CCPA", "Mandatory DPA provisions including documented instructions, confidentiality, security measures, sub-processor management, breach notification within 72 hours, and data deletion."),
            ("DORA Article 30", "Mandatory ICT contractual provisions — SLA definitions, availability guarantees, audit rights, exit strategy, incident notification, and sub-outsourcing governance."),
        ],
        "evidence": [
            "Redlined contract document with vendor comments and negotiation positions",
            "Legal review checklist cross-referenced against OCC 2023-17 Section 60",
            "DPA review sign-off confirming all GDPR Article 28(3) requirements",
            "SLA schedule review confirming measurable targets with remedies",
            "Insurance requirements confirmation",
            "Escalation matrix and governance structure for contract management",
        ],
        "data_points": "Number of redlined clauses, categories of requested changes, negotiation risk assessment (high/medium/low), estimated time to resolution, vendor authorized negotiator identity and authority level.",
    },
    {
        "num": 7,
        "name": "Contract Execution",
        "id": "Task_VendorContractSign",
        "description": "Vendor executes (counter-signs) the final negotiated contract, including the master services agreement, all schedules (SLA, DPA, SOW), and any side letters. The authorized signatory must be verified against corporate authorization records. Execution must occur before any data processing, system access provisioning, or service delivery begins.",
        "controls": [
            ("OCC 2023-17 Section 60", "No services may commence until the contract is fully executed with all mandatory OCC provisions present."),
            ("SOX Section 302", "Authorized signatory confirmation — verify corporate authority to bind the vendor."),
            ("GDPR Article 28 / CCPA", "Executed DPA is a legal prerequisite before any personal data processing begins."),
            ("DORA Article 30", "Contract must be registered in the ICT third-party provider register upon execution."),
            ("SEC 17a-4", "Executed contract must be stored in a compliant, tamper-evident repository (7+ years)."),
        ],
        "evidence": [
            "Fully executed contract (wet signature or qualified e-signature with audit trail)",
            "Executed Data Processing Agreement",
            "Authorized signatory verification (corporate resolution, power of attorney)",
            "Contract reference number assigned and logged",
            "Insurance certificates confirmed as current",
            "Executed NDA if not covered by master agreement",
            "Board/committee approval if contract exceeds delegated authority thresholds",
        ],
        "data_points": "Execution date, contract effective date, contract term, total contract value, authorized signatories, contract reference number, DPA execution confirmation, insurance coverage verification date, repository location.",
    },
    {
        "num": 8,
        "name": "Vendor Onboarding",
        "id": "Task_VendorOnboarding",
        "description": "Vendor is formally onboarded into the organization's ecosystem. Includes provisioning system access on least-privilege basis, establishing data flows per the DPA, adding vendor to the third-party risk register and ICT provider inventory, configuring monitoring baselines, and completing knowledge transfer. Vendor completes security awareness acknowledgment and acceptable use policy acceptance.",
        "controls": [
            ("OCC 2023-17 Section 80", "Ongoing monitoring baseline established — vendor added to third-party inventory with risk tier, criticality designation, and monitoring cadence."),
            ("NIST CSF 2.0 PR.AA-05", "Vendor access provisioned on least-privilege principle with MFA enforced."),
            ("DORA Article 28", "Vendor added to the ICT third-party provider register with all mandatory fields."),
            ("GDPR / CCPA", "Data flows documented in ROPA. DPA controls activated and verified."),
            ("ISO 27001 A.15.1", "Vendor formally enrolled in the supplier management program."),
        ],
        "evidence": [
            "Vendor entry in the third-party risk register with complete metadata",
            "Access provisioning record (systems, roles, permissions — least-privilege verified)",
            "Updated data flow diagram showing vendor data ingress/egress points",
            "Supplier onboarding checklist signed off by Risk and Governance team",
            "Vendor security awareness acknowledgment and acceptable use policy acceptance",
            "Emergency contact and escalation path registration",
            "Monitoring baseline configuration (KRIs, thresholds, alert recipients)",
        ],
        "data_points": "Systems accessed, access level per system, data categories processed, data flow endpoints, monitoring cadence assigned, KRI thresholds, onboarding completion date, first scheduled review date, vendor relationship manager assigned.",
    },
    {
        "num": 9,
        "name": "Deployment Support",
        "id": "Task_VendorDeploySupport",
        "description": "Vendor actively supports deployment and go-live within the organization's environment. Includes executing the deployment runbook, configuring production environments, performing data migration (if applicable), conducting smoke testing, establishing configuration baselines, verifying rollback procedures, and providing on-call support during stabilization.",
        "controls": [
            ("OCC 2023-17 Section 80", "Change management oversight — all vendor deployment activities must follow the organization's change management process."),
            ("NIST CSF 2.0 PR.PS-04", "Audit logs for all configuration changes, access events, and deployment actions. Configuration baseline established post-deployment."),
            ("DORA Article 12", "ICT change management — vendor changes covered under resilience testing. Rollback verified before go-live."),
            ("ISO 27001 A.12.1.2", "Change management controls — formal change request, impact assessment, approval, implementation, and post-implementation review."),
        ],
        "evidence": [
            "Vendor-provided deployment runbook with step-by-step procedures",
            "Go-live sign-off record (Technical Assessment, Business Owner, Vendor)",
            "Configuration baseline snapshot captured post-deployment",
            "Rollback plan confirmed — tested and verified before production deployment",
            "Deployment log with timestamps, actions, personnel, and issues",
            "Smoke test results confirming core functionality in production",
            "Data migration validation report (if applicable)",
        ],
        "data_points": "Deployment date/duration, deployment method (blue-green, rolling, big-bang), environments deployed to, config items changed, data migration volume/validation, smoke test results, issues and resolution, rollback test result, stabilization period, on-call contact.",
    },
    {
        "num": 10,
        "name": "Close Request",
        "id": "Task_VendorCloseRequest",
        "description": "Vendor confirms engagement completion and provides final deliverables for handoff to ongoing operations and monitoring. Includes archiving all security and compliance evidence, establishing vendor performance baselines, completing final documentation, confirming data handling obligations are met, and transitioning to the ongoing monitoring cadence.",
        "controls": [
            ("OCC 2023-17 Section 80", "Transition to ongoing monitoring — vendor performance baseline documented and monitoring cadence formally established."),
            ("NIST CSF 2.0 GV.SC-09", "Supplier performance debrief — lessons-learned review of the onboarding process."),
            ("GDPR / CCPA", "Confirm all data processing activities are within contracted scope. Verify no unauthorized residual data retained."),
            ("SEC 17a-4", "Engagement records archived in compliant, tamper-evident repository. Complete audit trail preserved."),
        ],
        "evidence": [
            "Vendor closure confirmation signed by both parties",
            "Performance baseline scorecard with initial KRI values and targets",
            "Data handling confirmation letter from vendor attesting DPA compliance",
            "Complete engagement summary filed to the audit record",
            "Transition to monitoring cadence documented (reviewer, schedule, escalation path)",
            "Lessons learned document from the onboarding process",
            "Updated third-party risk register reflecting steady-state risk tier",
        ],
        "data_points": "Onboarding completion date, total duration (days), onboarding cost, initial vendor performance score, monitoring cadence assigned, next scheduled review date, vendor relationship manager, open action items with owners and due dates, process improvement recommendations.",
    },
]

TOTAL_SLIDES = 1 + len(tasks) * 2  # title + 2 slides per task (overview + details)

# ── Title Slide ──
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(3.2), fill_color=DARK_BLUE)
add_shape(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.08), fill_color=ACCENT)

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(1),
             "Vendor / Third Party", font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)
add_text_box(slide, Inches(1), Inches(1.7), Inches(11), Inches(0.6),
             "External Software or Contracting", font_size=24, color=RGBColor(0xA0, 0xC4, 0xFF), alignment=PP_ALIGN.LEFT)
add_text_box(slide, Inches(1), Inches(2.4), Inches(11), Inches(0.5),
             "Swimlane User Task Documentation  |  Onboarding Ideal State v2", font_size=14, color=RGBColor(0x80, 0xA8, 0xE0), alignment=PP_ALIGN.LEFT)

# Task overview grid
add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.5),
             "10 User Tasks in the Vendor Lifecycle", font_size=20, bold=True, color=DARK_BLUE)

grid_y = Inches(4.5)
for i, task in enumerate(tasks):
    col = i % 5
    row = i // 5
    x = Inches(1 + col * 2.3)
    y = grid_y + Emu(int(row * Inches(1.1)))
    box = add_shape(slide, x, y, Inches(2.1), Inches(0.85), fill_color=LIGHT_BLUE, line_color=MED_BLUE)
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.text = f"{task['num']}. {task['name']}"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    box.text_frame.paragraphs[0].space_before = Pt(8)

add_footer(slide, 1, TOTAL_SLIDES)

# ── Task Slides (2 per task: Overview + Detail) ──
slide_num = 2
for task in tasks:
    # ── Slide A: Task Overview ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_header_bar(slide)

    # Header text
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(1), Inches(0.6),
                 f"{task['num']:02d}", font_size=32, bold=True, color=WHITE, font_name="Calibri")
    add_text_box(slide, Inches(1.4), Inches(0.15), Inches(10), Inches(0.35),
                 task["name"], font_size=22, bold=True, color=WHITE)
    add_text_box(slide, Inches(1.4), Inches(0.5), Inches(10), Inches(0.3),
                 f"BPMN ID: {task['id']}  |  Candidate Group: vendor-response", font_size=10, color=RGBColor(0xA0, 0xC4, 0xFF))

    # Description box
    add_shape(slide, Inches(0.5), Inches(1.15), Inches(12.3), Inches(1.6), fill_color=LIGHT_BLUE, line_color=MED_BLUE)
    add_text_box(slide, Inches(0.7), Inches(1.2), Inches(11.9), Inches(0.3),
                 "DESCRIPTION", font_size=10, bold=True, color=MED_BLUE)
    add_text_box(slide, Inches(0.7), Inches(1.5), Inches(11.9), Inches(1.2),
                 task["description"], font_size=11, color=BLACK)

    # Regulatory Controls
    add_text_box(slide, Inches(0.5), Inches(2.95), Inches(12), Inches(0.3),
                 "REGULATORY CONTROLS", font_size=12, bold=True, color=DARK_BLUE)
    add_shape(slide, Inches(0.5), Inches(3.25), Inches(12.3), Inches(0.03), fill_color=ACCENT)

    ctrl_y = 3.35
    for reg, desc in task["controls"]:
        # Regulation name
        add_text_box(slide, Inches(0.7), Inches(ctrl_y), Inches(11.9), Inches(0.25),
                     f"{reg}", font_size=10, bold=True, color=ACCENT)
        ctrl_y += 0.22
        # Description - truncate if needed for space
        desc_text = desc if len(desc) < 200 else desc[:197] + "..."
        add_text_box(slide, Inches(0.9), Inches(ctrl_y), Inches(11.7), Inches(0.35),
                     desc_text, font_size=9, color=GRAY)
        ctrl_y += 0.35

    add_footer(slide, slide_num, TOTAL_SLIDES)
    slide_num += 1

    # ── Slide B: Evidence & Data Points ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_header_bar(slide)

    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(1), Inches(0.6),
                 f"{task['num']:02d}", font_size=32, bold=True, color=WHITE, font_name="Calibri")
    add_text_box(slide, Inches(1.4), Inches(0.15), Inches(10), Inches(0.35),
                 f"{task['name']} — Evidence & Data", font_size=22, bold=True, color=WHITE)
    add_text_box(slide, Inches(1.4), Inches(0.5), Inches(10), Inches(0.3),
                 "Evidence Collection Requirements and Key Data Points", font_size=10, color=RGBColor(0xA0, 0xC4, 0xFF))

    # Evidence section
    add_text_box(slide, Inches(0.5), Inches(1.15), Inches(7), Inches(0.3),
                 "EVIDENCE COLLECTION", font_size=12, bold=True, color=DARK_BLUE)
    add_shape(slide, Inches(0.5), Inches(1.45), Inches(7.3), Inches(0.03), fill_color=ACCENT)

    ev_y = 1.6
    for ev in task["evidence"]:
        # Bullet marker
        add_text_box(slide, Inches(0.6), Inches(ev_y), Inches(0.3), Inches(0.3),
                     "\u2714", font_size=11, color=ACCENT)
        add_text_box(slide, Inches(0.95), Inches(ev_y), Inches(6.7), Inches(0.4),
                     ev, font_size=10, color=BLACK)
        ev_y += 0.42

    # Data Points section (right column)
    add_shape(slide, Inches(8.2), Inches(1.15), Inches(4.6), Inches(5.5), fill_color=LIGHT_BLUE, line_color=MED_BLUE)
    add_text_box(slide, Inches(8.4), Inches(1.25), Inches(4.2), Inches(0.3),
                 "KEY DATA POINTS", font_size=12, bold=True, color=DARK_BLUE)
    add_shape(slide, Inches(8.4), Inches(1.55), Inches(4.2), Inches(0.03), fill_color=ACCENT)

    # Split data points by comma and display
    dp_items = [d.strip() for d in task["data_points"].replace(".,", ",").split(",") if d.strip()]
    dp_y = 1.7
    for dp in dp_items:
        dp_clean = dp.strip().rstrip(".")
        if dp_clean:
            add_text_box(slide, Inches(8.5), Inches(dp_y), Inches(0.25), Inches(0.25),
                         "\u2022", font_size=10, color=ACCENT)
            add_text_box(slide, Inches(8.75), Inches(dp_y), Inches(3.8), Inches(0.3),
                         dp_clean, font_size=9, color=BLACK)
            dp_y += 0.28

    add_footer(slide, slide_num, TOTAL_SLIDES)
    slide_num += 1

# Save
output = "/Users/proth/repos/sla/docs/vendor-third-party-tasks.pptx"
prs.save(output)
print(f"Saved to {output}")
print(f"Total slides: {TOTAL_SLIDES}")
