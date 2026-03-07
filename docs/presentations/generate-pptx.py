#!/usr/bin/env python3
"""Generate KPMG-branded PowerPoint from v3 discovery presentation data."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

KPMG_BLUE = RGBColor(0x00, 0x33, 0x8D)
KPMG_MEDIUM_BLUE = RGBColor(0x00, 0x5E, 0xB8)
KPMG_LIGHT_BLUE = RGBColor(0x00, 0x91, 0xDA)
KPMG_DARK_NAVY = RGBColor(0x00, 0x1D, 0x48)
KPMG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
KPMG_LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
KPMG_DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
KPMG_GOLD = RGBColor(0xD4, 0xA8, 0x43)
KPMG_EMERALD = RGBColor(0x00, 0xA8, 0x6B)
KPMG_ROSE = RGBColor(0xC6, 0x28, 0x28)
KPMG_AMBER = RGBColor(0xE6, 0x7E, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_gradient_bg(slide, c1, c2):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[1].color.rgb = c2
    fill.gradient_angle = 135


def add_text_box(slide, left, top, width, height, text, font_size=12,
                 bold=False, color=KPMG_DARK_GRAY, font_name='Open Sans',
                 alignment=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_footer(slide, text="KPMG Confidential", dark=False):
    color = RGBColor(0x99, 0x99, 0x99) if not dark else RGBColor(0x99, 0xAA, 0xCC)
    add_text_box(slide, Inches(0.5), Inches(7.0), Inches(4), Inches(0.4),
                 text, font_size=8, color=color)
    add_text_box(slide, Inches(9), Inches(7.0), Inches(4), Inches(0.4),
                 "March 2026", font_size=8, color=color, alignment=PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, title, body, accent_color=KPMG_BLUE):
    # Card background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = KPMG_WHITE
    shape.line.fill.background()
    shape.shadow.inherit = False
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.3), Inches(0.35),
                 title, font_size=11, bold=True, color=KPMG_BLUE)
    # Body
    add_text_box(slide, left + Inches(0.2), top + Inches(0.45), width - Inches(0.3), height - Inches(0.55),
                 body, font_size=9, color=RGBColor(0x66, 0x66, 0x66))


def add_metric(slide, left, top, value, label, value_color=KPMG_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = KPMG_WHITE
    shape.line.fill.background()
    add_text_box(slide, left, top + Inches(0.1), Inches(2), Inches(0.5),
                 value, font_size=24, bold=True, color=value_color,
                 font_name='Open Sans Condensed', alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(0.65), Inches(2), Inches(0.3),
                 label.upper(), font_size=7, color=RGBColor(0x66, 0x66, 0x66),
                 alignment=PP_ALIGN.CENTER)


def add_section_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide, KPMG_MEDIUM_BLUE, KPMG_BLUE)
    add_text_box(slide, Inches(1), Inches(2.5), Inches(10), Inches(1.2),
                 title, font_size=40, bold=True, color=KPMG_WHITE, font_name='Open Sans Condensed')
    if subtitle:
        add_text_box(slide, Inches(1), Inches(3.8), Inches(9), Inches(0.8),
                     subtitle, font_size=14, color=RGBColor(0xCC, 0xDD, 0xFF))
    # Gold accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.55), Inches(3), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = KPMG_GOLD
    line.line.fill.background()
    add_footer(slide, dark=True)
    return slide


def add_current_state_slide(title, findings, quote=None, quote_source=None, raci=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, KPMG_WHITE)
    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
                 title, font_size=24, bold=True, color=KPMG_BLUE, font_name='Open Sans Condensed')

    cols = 2
    rows = (len(findings) + 1) // 2
    for i, (ftitle, fbody, is_positive) in enumerate(findings):
        col = i % cols
        row = i // cols
        left = Inches(0.6 + col * 6.2)
        top = Inches(1.2 + row * 1.15)
        accent = KPMG_EMERALD if is_positive else KPMG_ROSE
        add_card(slide, left, top, Inches(5.9), Inches(1.0), ftitle, fbody, accent)

    if quote:
        y = Inches(1.2 + rows * 1.15 + 0.1)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), y, Inches(0.06), Inches(0.7))
        bar.fill.solid()
        bar.fill.fore_color.rgb = KPMG_GOLD
        bar.line.fill.background()
        add_text_box(slide, Inches(0.8), y, Inches(11), Inches(0.5),
                     f'"{quote}"', font_size=10, color=KPMG_DARK_GRAY)
        if quote_source:
            add_text_box(slide, Inches(0.8), y + Inches(0.45), Inches(11), Inches(0.25),
                         f"-- {quote_source}", font_size=8, color=RGBColor(0x66, 0x66, 0x66))

    if raci:
        add_text_box(slide, Inches(0.6), Inches(6.7), Inches(12), Inches(0.25),
                     raci, font_size=8, color=RGBColor(0x66, 0x66, 0x66))
    add_footer(slide)
    return slide


def add_recommendations_slide(title, recs, roadmap):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, KPMG_WHITE)
    add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
                 title, font_size=24, bold=True, color=KPMG_BLUE, font_name='Open Sans Condensed')

    for i, (rtitle, rbody, impact) in enumerate(recs[:4]):
        col = i % 2
        row = i // 2
        left = Inches(0.6 + col * 6.2)
        top = Inches(1.1 + row * 1.3)
        add_card(slide, left, top, Inches(5.9), Inches(1.15), rtitle,
                 f"{rbody}\n{impact}", KPMG_EMERALD)

    # Roadmap row
    y_road = Inches(3.9)
    colors = [KPMG_EMERALD, KPMG_LIGHT_BLUE, KPMG_MEDIUM_BLUE, KPMG_BLUE]
    labels = ["30 Days", "60 Days", "90 Days", "120 Days"]
    for i, (items, label) in enumerate(zip(roadmap, labels)):
        left = Inches(0.6 + i * 3.1)
        # Header bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, y_road, Inches(2.9), Inches(0.05))
        bar.fill.solid()
        bar.fill.fore_color.rgb = colors[i]
        bar.line.fill.background()
        add_text_box(slide, left, y_road + Inches(0.08), Inches(2.9), Inches(0.3),
                     label, font_size=10, bold=True, color=KPMG_BLUE)
        text = "\n".join(f"  {item}" for item in items)
        add_text_box(slide, left, y_road + Inches(0.4), Inches(2.9), Inches(2.8),
                     text, font_size=8, color=RGBColor(0x66, 0x66, 0x66))

    add_footer(slide)
    return slide


# ===== SLIDE 0: HERO =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, KPMG_BLUE, KPMG_MEDIUM_BLUE)
add_text_box(slide, Inches(1), Inches(1.5), Inches(10), Inches(1.5),
             "Software Onboarding\nProcess Transformation", font_size=44, bold=True,
             color=KPMG_WHITE, font_name='Open Sans Condensed')
add_text_box(slide, Inches(1), Inches(3.3), Inches(10), Inches(0.6),
             "Discovery Findings, Recommendations, and Implementation Roadmap",
             font_size=18, bold=True, color=RGBColor(0xCC, 0xDD, 0xFF))
add_text_box(slide, Inches(1), Inches(4.2), Inches(9), Inches(1),
             "Based on 14 stakeholder interviews across Architecture, Product, Security, Risk Management, "
             "Vendor Management, Finance, Legal, and Compliance teams. Covering 11 governance domains "
             "with actionable 30/60/90/120-day implementation plans.",
             font_size=11, color=RGBColor(0xBB, 0xCC, 0xEE))
# Badge row
badges = ["14 Stakeholder Sessions", "11 Governance Domains", "16 Gap Findings", "120-Day Roadmap"]
for i, b in enumerate(badges):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(1 + i * 2.6), Inches(5.5), Inches(2.3), Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x00, 0x2A, 0x70)
    shape.line.color.rgb = RGBColor(0x33, 0x55, 0xAA)
    shape.line.width = Pt(1)
    add_text_box(slide, Inches(1 + i * 2.6), Inches(5.53), Inches(2.3), Inches(0.35),
                 b, font_size=9, color=RGBColor(0xCC, 0xDD, 0xFF), alignment=PP_ALIGN.CENTER)
add_footer(slide, dark=True)

# ===== SLIDE 1: EXECUTIVE SUMMARY =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, KPMG_WHITE)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
             "Executive Summary", font_size=28, bold=True, color=KPMG_BLUE, font_name='Open Sans Condensed')
add_text_box(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.7),
             "The software onboarding process spans 6 to 9 months end-to-end, driven by 18 sequential "
             "committees, 5+ disconnected intake channels, and critical resource bottlenecks in Security and Legal. "
             "The START initiative (9 months old) created centralized awareness but did not integrate underlying team processes.",
             font_size=11, color=KPMG_DARK_GRAY)

# Metrics row
metrics = [("6-9 mo", "Current E2E Cycle", KPMG_ROSE), ("18", "Committees", KPMG_ROSE),
           ("335", "Assessments/Year", KPMG_BLUE), ("75 days", "DD Internal Review", KPMG_ROSE),
           ("60+", "AI Queue Items", KPMG_ROSE)]
for i, (val, lab, clr) in enumerate(metrics):
    add_metric(slide, Inches(0.6 + i * 2.5), Inches(1.8), val, lab, clr)

add_text_box(slide, Inches(0.6), Inches(3.1), Inches(6), Inches(0.3),
             "Biggest Challenges", font_size=14, bold=True, color=KPMG_DARK_NAVY, font_name='Open Sans Condensed')
challenges = [
    ("Sequential Reviews", "Requesters present to ARB, TBC, AI Governance, and DART sequentially. Each committee 5+ weeks apart."),
    ("Requester Burden", "DART formation falls entirely on the requester, who must contact 5-6 teams independently."),
    ("Resource Crisis", "2 people negotiate 30+ contracts/month. Security is #1 SLA bottleneck. \"Half a person\" owns START."),
]
for i, (ct, cb) in enumerate(challenges):
    add_card(slide, Inches(0.6 + i * 4.1), Inches(3.45), Inches(3.9), Inches(1.1), ct, cb, KPMG_ROSE)

add_text_box(slide, Inches(0.6), Inches(4.75), Inches(6), Inches(0.3),
             "Highest-ROI Investment Areas", font_size=14, bold=True, color=KPMG_DARK_NAVY, font_name='Open Sans Condensed')
investments = [
    ("Parallel Evaluation", "Replace 18 sequential committees with 5 parallel evaluation streams."),
    ("Unified Intake + Deal-Killer", "Consolidate 5+ channels. Block non-starters at day 1."),
    ("Contract Automation", "Automate contract review for the 2-person team handling 30+/month."),
]
for i, (it, ib) in enumerate(investments):
    add_card(slide, Inches(0.6 + i * 4.1), Inches(5.1), Inches(3.9), Inches(1.0), it, ib, KPMG_EMERALD)

add_footer(slide)

# ===== SLIDE 2: E2E WORKFLOW =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, KPMG_WHITE)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
             "End-to-End Workflow: Current Pain Points", font_size=24, bold=True,
             color=KPMG_BLUE, font_name='Open Sans Condensed')

# Phase pipeline
phases = [("1 Intake", "5+ channels", KPMG_ROSE), ("2 Prioritize", "No formula", KPMG_ROSE),
          ("3 Funding", "Locked forms", KPMG_AMBER), ("4 Sourcing", "75d review", KPMG_ROSE),
          ("5 Cyber", "#1 bottleneck", KPMG_ROSE), ("6 EA", "2wk ARB", KPMG_AMBER),
          ("7 AI Gov", "60+ queue", KPMG_AMBER), ("8 Legal", "2 people", KPMG_ROSE),
          ("9 Contract", "Up to 1.5yr", KPMG_AMBER)]
for i, (name, issue, color) in enumerate(phases):
    left = Inches(0.4 + i * 1.38)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.2), Inches(1.3), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_text_box(slide, left, Inches(1.22), Inches(1.3), Inches(0.4),
                 name, font_size=9, bold=True, color=KPMG_WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(1.58), Inches(1.3), Inches(0.3),
                 issue, font_size=7, color=RGBColor(0xFF, 0xDD, 0xDD), alignment=PP_ALIGN.CENTER)

# Bottleneck table
add_text_box(slide, Inches(0.6), Inches(2.3), Inches(6), Inches(0.3),
             "Critical Bottlenecks", font_size=14, bold=True, color=KPMG_DARK_NAVY, font_name='Open Sans Condensed')
bottlenecks = [
    "Contract Negotiation: 2 people / 30+ contracts monthly [CRITICAL]",
    "Security Review: \"Biggest bottleneck... reason our SLA takes 2 weeks\" [CRITICAL]",
    "Sequential Committees (18): Same presentation repeated 3-4 times [CRITICAL]",
    "DART Formation: Requester manages 5-6 teams independently [CRITICAL]",
    "AI Governance Queue: 60+ items, 3 separate committees [HIGH]",
    "Business Council Quorum: 2-3 of 8-10 members attend [HIGH]",
]
for i, b in enumerate(bottlenecks):
    add_text_box(slide, Inches(0.8), Inches(2.7 + i * 0.32), Inches(6.5), Inches(0.3),
                 b, font_size=9, color=KPMG_DARK_GRAY)

# Strengths
add_text_box(slide, Inches(7.5), Inches(2.3), Inches(5), Inches(0.3),
             "What the Organization Does Well", font_size=14, bold=True, color=KPMG_DARK_NAVY, font_name='Open Sans Condensed')
strengths = [
    ("Architecture Governance", "Dedicated Facilitator pre-screens artifacts, manages follow-ups. Most disciplined team."),
    ("TPRM Due Diligence", "Reduced DD from 144 to 75 days. Output doubled to 335 assessments/year."),
    ("Acquisition 2.0", "Brings all teams together at start. Collective go/no-go at first tollgate."),
]
for i, (st, sb) in enumerate(strengths):
    add_card(slide, Inches(7.5), Inches(2.7 + i * 1.2), Inches(5.2), Inches(1.05), st, sb, KPMG_EMERALD)

add_footer(slide)

# ===== SLIDE 3: ROADMAP =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, KPMG_WHITE)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
             "Implementation Roadmap: 30 / 60 / 90 / 120 Days", font_size=24, bold=True,
             color=KPMG_BLUE, font_name='Open Sans Condensed')

roadmap_data = [
    ("Days 1-30: Consolidate", KPMG_EMERALD, [
        "Unified intake form (replace 5+ channels)", "Deal-killer pre-screen gate",
        "Completeness quality gate", "Committee scope definition",
        "3-pathway routing: Buy/Build/Enable", "Prioritization scoring formula",
        "Status notification at transitions", "Concierge role definition"]),
    ("Days 31-60: Automate", KPMG_LIGHT_BLUE, [
        "Parallel evaluation (replace sequential)", "Automated DART formation",
        "Progressive forms", "Contract review automation pilot",
        "Security tiered assessment", "AI governance consolidation",
        "Finance rework loop", "Workload visibility dashboard"]),
    ("Days 61-90: Optimize", KPMG_MEDIUM_BLUE, [
        "AI fast-track pathway (2-week target)", "Enable pathway live",
        "Time-bound conditional approvals", "Mandatory ownership assignment",
        "NDA-first gate enforcement", "Automated security baselines",
        "Self-service mini-RFP tools", "SLA enforcement + escalation"]),
    ("Days 91-120: Scale", KPMG_BLUE, [
        "Pre-onboarding idea funnel", "Full workload dashboard",
        "Exception routing for rapid RA", "Post-onboarding utilization tracking",
        "Distributed pod model pilot", "Annual ownership validation",
        "Process mining / CI", "Executive KPI reporting"]),
]

for i, (title, color, items) in enumerate(roadmap_data):
    left = Inches(0.5 + i * 3.15)
    # Top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.2), Inches(2.95), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text_box(slide, left, Inches(1.3), Inches(2.95), Inches(0.35),
                 title, font_size=11, bold=True, color=KPMG_BLUE)
    text = "\n".join(f"  {item}" for item in items)
    add_text_box(slide, left, Inches(1.7), Inches(2.95), Inches(4.5),
                 text, font_size=8, color=RGBColor(0x66, 0x66, 0x66))

# ROI metrics
roi = [("60-70%", "Cycle Time Reduction"), ("18 to 5", "Sequential to Parallel"),
       ("3 Pathways", "Buy/Build/Enable"), ("Day 1", "Non-Starter ID")]
for i, (v, l) in enumerate(roi):
    left = Inches(0.6 + i * 3.15)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(6.2), Inches(2.9), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = KPMG_BLUE
    shape.line.fill.background()
    add_text_box(slide, left, Inches(6.22), Inches(2.9), Inches(0.4),
                 v, font_size=20, bold=True, color=KPMG_GOLD,
                 font_name='Open Sans Condensed', alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(6.6), Inches(2.9), Inches(0.3),
                 l, font_size=9, color=RGBColor(0xCC, 0xDD, 0xFF), alignment=PP_ALIGN.CENTER)

add_footer(slide)

# ===== GOVERNANCE DOMAIN DEEP DIVES =====

# Define all 11 domains with their data
domains = [
    {
        "name": "Intake",
        "subtitle": "Initial request capture, portfolio review, and solution discovery.",
        "bpmn_image": "sp1-intake.png",
        "bpmn_title": "SP1: Request and Triage",
        "bpmn_desc": "Request capture, existing solution review, bypass routing for leveraged software, initial triage with 2-day SLA timer, and DMN-driven pathway routing.",
        "findings": [
            ("5+ disconnected entry points", "ServiceNow START, AI Use Case form, AI Governance form, Rapid Risk Assessment (Power Apps), email/chat.", False),
            ("Requester confusion", "First-time requesters struggle with complexity. Forms discovered ad hoc.", False),
            ("Forms require multi-disciplinary expertise", "RAE form has 80 questions. Business partners frequently cannot complete accurately.", False),
            ("Process bypasses are common", "Technology teams especially prone to bypassing standard processes.", False),
            ("No completeness validation", "Incomplete submissions cascade, causing rework when requirements surface late.", False),
            ("START created awareness", "ServiceNow central intake (9 months old) prevents work on unapproved initiatives.", True),
        ],
        "quote": "The process never solved the step of a requester who's never been through onboarding having to work with 5-6 teams through independent processes.",
        "quote_source": "Vendor Management Lead",
        "raci": "RACI: R/A Business | C Governance, Compliance",
        "recs": [
            ("Unified Intake Gateway", "Single entry point absorbing all channels. Dynamic routing by request type.", "Eliminates 4 redundant channels"),
            ("Deal-Killer Pre-Screen", "DMN-driven no-go check at submission. Block non-starters immediately.", "Saves 100% downstream effort on non-starters"),
            ("Completeness Quality Gate", "AI-assisted pre-screening validates minimum fields before SME routing.", "Eliminates rework from late requirements"),
            ("Request Classification", "Automated: Defined Need, Forced Update, or Speculative. Right-sizes governance.", "Matches effort to complexity"),
        ],
        "roadmap": [
            ["Deploy unified intake form", "Implement deal-killer checklist", "Define minimum fields", "Retire redundant channels"],
            ["Progressive form logic", "AI-assisted completeness scoring", "Integrate feedback platform"],
            ["Self-service mini-RFP tools", "Auto-classification (DMN)", "Shift-left: questions upfront"],
            ["Pre-onboarding idea funnel", "Threshold-based escalation", "Full intake analytics"],
        ],
    },
    {
        "name": "Prioritization",
        "subtitle": "Backlog ranking, pathway selection, and resource allocation.",
        "bpmn_image": "sp2-planning.png",
        "bpmn_title": "SP2: Planning and Routing",
        "bpmn_desc": "Preliminary analysis, backlog prioritization, DMN-1 risk tier classification, DMN-2 pathway assignment. Unacceptable risk routes to rejection.",
        "findings": [
            ("No formal prioritization formula", "Teams 'horse trade' internally. No force-ranking mechanism across the enterprise.", False),
            ("'Whoever screams loudest' gets priority", "EVP support pushes other reviews down. No SLA enforcement possible.", False),
            ("Business Council quorum failures", "Monthly meetings draw 2-3 of 8-10 members. Evolved to email voting.", False),
            ("All requests follow same heavyweight process", "No auto-approval for low-risk. No de-prioritization guidelines.", False),
            ("Capacity management disconnected", "PI planning allocates to roadmap; exceptions create unanticipated workload.", False),
            ("Governance Specialist sees path forward", "Scoring intake responses can drive prioritization. Standards must be defined upfront.", True),
        ],
        "quote": "If you set up repercussion for SLA not achieved, what's the solution? There's no real effective prioritization.",
        "quote_source": "Security Architect",
        "raci": "RACI: R/A Governance | C Business, Finance, Technical Assessment",
        "recs": [
            ("Quantitative Scoring Formula", "DMN-driven: business impact, strategic alignment, urgency, risk tier, capacity.", "Replaces subjective horse trading"),
            ("3-Pathway Routing: Buy/Build/Enable", "Enable (Vendor Affinity) skips funding validation, reduces evaluation scope.", "30-day turnaround for Enable vs 6-9 months"),
            ("Tiered Governance Depth", "Risk tier determines intensity. Low-risk gets streamlined review.", "Frees 40%+ reviewer capacity"),
            ("Capacity-Aware Scheduling", "Integrate team capacity into routing. Display queue depth at submission.", "Sets realistic expectations"),
        ],
        "roadmap": [
            ["Define scoring formula", "Implement 3-pathway routing", "Publish criteria", "Queue transparency"],
            ["DMN decision table", "Capacity integration pilot", "Auto-routing for low-risk"],
            ["Enable pathway live", "AI fast-track pathway", "SLA enforcement"],
            ["Full portfolio dashboard", "Predictive wait times", "Continuous calibration"],
        ],
    },
    {
        "name": "Funding / Finance",
        "subtitle": "Financial analysis, budget authorization, and TCO assessment.",
        "bpmn_image": None,
        "bpmn_title": None,
        "bpmn_desc": None,
        "findings": [
            ("Formal denial restarts entire process", "Cannot reroute coding matrix issues to FP&A. Minor corrections require full restart.", False),
            ("Financial form outdated and locked", "Shows 2024 in 2026. Password unknown (owner left). Downloaded, completed offline.", False),
            ("Enable pathway forced through funding", "Vendor Affinity products require no org investment but process demands justification.", False),
            ("Business case required multiple times", "Redundancy across intake forms all describing the same business case.", False),
        ],
        "quote": "There's redundancy on what all these intake forms are... they're all describing what is the business case.",
        "quote_source": "Product Manager",
        "raci": "RACI: R/A Finance | C Governance, Procurement | I Oversight",
        "recs": [
            ("Finance Rework Loop", "Correction pathway for coding matrix issues. Minor fixes route to FP&A directly.", "Eliminates weeks of rework"),
            ("Enable Pathway Bypass", "DMN detects 'no org investment' and bypasses financial gates entirely.", "Removes friction for zero-cost tools"),
            ("Modernized Financial Form", "Dynamic digital version replacing locked form. Pre-populate from intake data.", "Single source of truth"),
            ("Consolidated Business Case", "One business case at intake, enriched progressively. No repeated justification.", "60%+ form burden reduction"),
        ],
        "roadmap": [
            ["Replace locked form", "Define Enable bypass rules", "Map financial data flow"],
            ["Finance rework loop", "Pre-populated fields", "Conditional fields per pathway"],
            ["Automated TCO/ROI calc", "Finance dashboard integration"],
            ["Full financial lifecycle", "Budget forecasting integration"],
        ],
    },
    {
        "name": "Sourcing",
        "subtitle": "Vendor landscape, risk evaluation, due diligence, and procurement.",
        "bpmn_image": "sp3-sourcing.png",
        "bpmn_title": "SP3: Evaluation and Due Diligence",
        "bpmn_desc": "5 parallel evaluation streams (Tech Architecture, Security, Risk/Compliance/Legal, Financial Analysis, Vendor Landscape) with 5-day SLA and vendor DD integration.",
        "findings": [
            ("RAE: 80 questions, 28-29 day actual vs 14 day target", "Internal questionnaire assigns inherent risk tier and determines DD level.", False),
            ("DD: 830 questions, 75 day internal review", "Down from 144 days. Vendor completion 30 days (ahead of 42-day target).", False),
            ("'Sourcing doesn't source'", "They manage contract lifecycle. True sourcing falls on requesters.", False),
            ("NDA timing disagreement", "Legal and Security want NDA first. Other teams disagree. No consistent policy.", False),
            ("TPRM performance improving", "Output doubled to 335 assessments/year. Team of 8 reducing timelines.", True),
            ("OneTrust deployed for risk tracking", "API connections between Ariba and Oracle. Risk assessment foundation exists.", True),
        ],
        "quote": "That's a fantastic idea that should be in the slides.",
        "quote_source": "Risk Management Lead, on shift-left strategy",
        "raci": "RACI: R Procurement | A Governance | C Compliance | I Oversight",
        "recs": [
            ("NDA-First Gate", "Enforce NDA execution before detailed vendor engagement.", "Protects IP, aligns Legal and Security"),
            ("Shift-Left: Self-Service Mini-RFP", "Structured RFP tools before formal onboarding. Maintains competitive leverage.", "Better vendor selection"),
            ("Tiered Due Diligence", "Risk tier determines DD depth. Low-risk: automated only. High: full 830-question.", "60% DD cycle time reduction for low-risk"),
            ("Vendor-Level Aggregation", "Single vendor spanning 10+ BUs shares compliance artifacts.", "Eliminates redundant assessments"),
        ],
        "roadmap": [
            ["NDA-first enforcement", "Risk-tiered DD matrix", "Vendor de-duplication"],
            ["Mini-RFP templates", "Abbreviated DD for low-risk", "OneTrust optimization"],
            ["Vendor artifact sharing", "Automated landscape scan", "AI-assisted routing"],
            ["Full self-service portal", "Competitive intelligence", "Vendor performance scoring"],
        ],
    },
    {
        "name": "Cybersecurity",
        "subtitle": "Security assessment, vulnerability management, and architecture review.",
        "bpmn_image": "sp3-sourcing.png",
        "bpmn_title": "SP3: Security Assessment Stream",
        "bpmn_desc": "Security assessment runs as one of 5 parallel streams in SP3. Tiered approach: automated baseline, elevated manual review, or full penetration testing based on risk tier.",
        "findings": [
            ("'Security is our biggest bottleneck'", "Primary reason ARB SLA takes 2 weeks. Understaffed security architecture.", False),
            ("'I need three of me right now'", "Request volume increasing, especially AI reviews. Cannot scale with current staffing.", False),
            ("No 'secure by design' standard", "Teams don't know minimum controls. Enforcement 'fairly loose.'", False),
            ("3x vendor contact for AI", "Tech risk, cybersecurity, and TPRM each contact vendor independently.", False),
            ("Zero automation", "All routing manual. No workflow integration. Dependency on individual knowledge.", False),
            ("Multiple security platforms", "Each addresses specific control requirements. Different security needs covered.", True),
        ],
        "quote": "Security is our biggest bottleneck... They are the reason our SLA takes two weeks.",
        "quote_source": "Architecture Lead",
        "raci": "RACI: R/A Technical Assessment | C Governance, Compliance",
        "recs": [
            ("Tiered Security Assessment", "DMN-driven. Baseline: automated. Elevated: auto+manual. Major: full+pen test.", "70% handled by automated baseline"),
            ("Consolidated Vendor Contact", "Single coordinated engagement replacing 3 independent contact points.", "Eliminates 2/3 vendor touchpoints"),
            ("Security Baseline Definition", "Minimum controls published. Baseline, elevated, and major hierarchies.", "Clear expectations for all"),
            ("Parallel Evaluation", "Security concurrent with Architecture, Compliance, Financial reviews.", "Eliminates 4-8 weeks sequential waiting"),
        ],
        "roadmap": [
            ["Define 3-tier framework", "Consolidate vendor contact", "Document baselines"],
            ["Automated baseline checks", "Parallel evaluation", "AI stream consolidation"],
            ["Full tiered assessment", "Automated pen test scheduling", "Security SLA enforcement"],
            ["Security posture dashboard", "Continuous monitoring", "Risk-based resource allocation"],
        ],
    },
    {
        "name": "Enterprise Architecture",
        "subtitle": "Architecture review, integration planning, and standards compliance.",
        "bpmn_image": "sp3-sp4-ea.png",
        "bpmn_title": "SP3-SP4: Architecture Review Lifecycle",
        "bpmn_desc": "Architecture review spans evaluation (SP3) through contracting (SP4). Governance Facilitator pre-screens artifacts, manages follow-ups, and runs ARB and SDRB.",
        "findings": [
            ("Governance Facilitator role (model for E2E)", "Pre-screens artifacts, manages follow-ups. Runs ARB and SDRB.", True),
            ("ARB: 2-week SLA. SDRB: same-day", "Structured workflow with JIRA integration. Most disciplined team.", True),
            ("AI tooling adoption", "Cursor AI for diagrams. Custom agent scans designs for pattern conformance.", True),
            ("HLD/PSS/SDD complexity", "Requesters confused about which artifact path applies to their situation.", False),
            ("Committee overlap with TBC", "Present to ARB, then TBC, then AI governance. Overlapping scope.", False),
            ("Domain-funded model constraints", "'If they want us to pursue all that work, they need to fund us appropriately.'", False),
        ],
        "quote": "I like how you're running your shop.",
        "quote_source": "Consulting Team, on Architecture governance",
        "raci": "RACI: R/A Technical Assessment | C Governance",
        "recs": [
            ("Scale Governance Facilitator Model", "Blueprint for E2E 'quarterback.' Expand concierge across all domains.", "Proven model reduces requester burden"),
            ("Define Clear Committee Scope", "Architecture questions by architects only. TBC: business case. AI Gov: AI risks.", "Eliminates redundant presentations"),
            ("Simultaneous Engagement Model", "All major voting players engaged simultaneously when request comes in.", "Architecture Lead's own proposal"),
            ("Domain-Based Auto-Assignment", "Automated routing by requesting domain. EA leader to architect by bandwidth.", "Reduces triage, balances workload"),
        ],
        "roadmap": [
            ["Document scope boundaries", "Define concierge role", "Map domain routing"],
            ["Simultaneous engagement pilot", "Automated domain routing", "Consolidated artifacts"],
            ["Concierge operational", "Cross-committee enforcement", "AI pattern conformance"],
            ["Full parallel model", "Architecture metrics", "Workload resource planning"],
        ],
    },
    {
        "name": "Compliance",
        "subtitle": "Regulatory compliance review across OCC 2023-17, GDPR/CCPA, SOX, DORA.",
        "bpmn_image": None,
        "bpmn_title": None,
        "bpmn_desc": None,
        "findings": [
            ("End-of-year compliance failures", "PII found in uncontrolled systems. Enforcement 'fairly loose.'", False),
            ("Multiple teams assess independently", "Risk, Legal, Privacy, Compliance each conduct own review without shared findings.", False),
            ("No standardized business case definition", "TBC overlaps financial review. Architecture assessing financials.", False),
            ("Contract deviations not reportable", "No structured format. Unknown compliance status for older contracts.", False),
            ("Regulatory awareness improving", "Teams recognize OCC 2023-17, DORA, EU AI Act requirements.", True),
            ("OneTrust for risk documentation", "Control gap documentation and risk acceptance tracking. Audit trail foundation.", True),
        ],
        "quote": None,
        "quote_source": None,
        "raci": "RACI: R/A Compliance | C Governance, Oversight",
        "recs": [
            ("Consolidated Compliance Stream", "Single review combining Risk, Legal, Privacy, Compliance. Shared findings.", "One review instead of four"),
            ("Regulatory Annotation Framework", "Every task mapped to applicable regulations. Auto compliance artifact generation.", "Audit-ready documentation"),
            ("Contract Deviation Tracking", "Structured format in OneTrust. Automated alerting for expiring risk acceptances.", "Full visibility into contract risk"),
            ("Compliance Quality Gates", "Phase-boundary checks. Automated validation of required artifacts.", "Prevents downstream failures"),
        ],
        "roadmap": [
            ["Map regulatory requirements", "Define consolidated stream", "Deviation template"],
            ["Shared findings database", "Auto compliance artifacts", "Phase quality gates"],
            ["Full regulatory annotation", "Deviation dashboard", "Expiry alerts"],
            ["Compliance analytics", "Regulatory change mgmt", "Audit readiness reporting"],
        ],
    },
    {
        "name": "AI Governance",
        "subtitle": "AI risk classification, model validation, and the most urgent governance challenge.",
        "bpmn_image": "sp3-sourcing.png",
        "bpmn_title": "SP3: AI Governance Stream",
        "bpmn_desc": "AI governance runs as a parallel evaluation stream in SP3. Consolidates 3 committees to 1, with fast-track pathway targeting 2-week turnaround for pre-approved risk postures.",
        "findings": [
            ("60+ items in queue", "Multiple tools for same function. No alignment to AI strategy.", False),
            ("3 separate AI committees", "Working Group, Cyber Review, Risk Review, plus Governance Committee. Sequential.", False),
            ("Overly restrictive AI addendum", "Extended vendor negotiations. External firms push back on AI terms.", False),
            ("3 additional AI questionnaires 'snuck up'", "Risk Management Lead 'really annoyed they exist.' Merging underway.", False),
            ("Vendors contacted 3x for AI", "Tech risk, cybersecurity, TPRM each contact independently.", False),
            ("No consistent escalation standard", "Unclear if adding models to packages = new use case.", False),
        ],
        "quote": "It's an ever-changing world that's different every other day... there needs to be some type of standardization.",
        "quote_source": "Risk Management Lead",
        "raci": "RACI: R/A AI Review | C Technical Assessment, Compliance | I Governance",
        "recs": [
            ("Consolidate 3 AI Committees to 1", "Single AI review stream. Cross-functional representation in one meeting.", "Eliminates months of sequential queuing"),
            ("AI Fast-Track Pathway", "Pre-defined risk posture. AI+Security in parallel only. 2-week target.", "2 weeks vs 6-9 months"),
            ("Unified AI Questionnaire", "Merge 3 questionnaires. One vendor engagement. Distributed findings.", "1 vendor contact instead of 3"),
            ("AI No-Go List", "Non-starter models/vendors communicated early. DMN red/green decisions at intake.", "Zero wasted effort on non-starters"),
        ],
        "roadmap": [
            ["Publish AI no-go list", "Define single AI stream", "Merge questionnaires"],
            ["AI fast-track pilot", "Consolidated committee", "Unified vendor engagement"],
            ["Pre-defined risk posture", "AI SLA targets (2 weeks)", "Auto classification"],
            ["Full AI governance lifecycle", "Model inventory", "EU AI Act framework"],
        ],
    },
    {
        "name": "Privacy",
        "subtitle": "Data protection, classification, and cross-border transfer compliance.",
        "bpmn_image": None,
        "bpmn_title": None,
        "bpmn_desc": None,
        "findings": [
            ("Privacy review embedded across teams", "Privacy SME participates in vendor risk, AI governance, compliance. No dedicated stream.", False),
            ("PII found in uncontrolled systems", "End-of-year compliance failures. Data classification not consistently enforced.", False),
            ("Different risk thresholds per team", "Varying risk acceptance for AI privacy concerns. No unified classification.", False),
            ("Data questions could be answered earlier", "Hosting, storage, transmission resolvable through proper sourcing events.", False),
            ("Privacy SME engaged across workstreams", "Active participation demonstrates organizational awareness.", True),
            ("GDPR/CCPA framework exists", "Challenge is consistent application during onboarding, not absence of policy.", True),
        ],
        "quote": None,
        "quote_source": None,
        "raci": "RACI: R/A Compliance (Privacy) | C Governance, Business",
        "recs": [
            ("Data Classification at Intake", "Mandatory fields in unified form. Determines DPIA vs standard vs no PII.", "Right-sizes privacy effort from day one"),
            ("Privacy as Parallel Stream", "Dedicated branch concurrent with security and compliance.", "Consistent, predictable timeline"),
            ("Unified Privacy Risk Classification", "Single framework replacing team-by-team thresholds. DMN-driven.", "Eliminates inconsistent decisions"),
            ("Early Data Residency Screening", "Validated at sourcing stage. Prevents late cross-border discoveries.", "Eliminates rework from late findings"),
        ],
        "roadmap": [
            ["Add classification to intake", "Define privacy triggers", "Map residency requirements"],
            ["Parallel privacy stream", "Unified classification", "DPIA automation template"],
            ["Cross-border validation", "Privacy-by-design checklist", "Automated PII detection"],
            ["Privacy impact dashboard", "Regulatory monitoring", "Annual privacy audit"],
        ],
    },
    {
        "name": "Commercial Counsel",
        "subtitle": "Contract negotiation, SOW/MSA redlining, and vendor agreement execution.",
        "bpmn_image": "sp4-commercial.png",
        "bpmn_title": "SP4: Contracting and Build",
        "bpmn_desc": "Buy path: contract refinement, PoC, tech risk review, negotiation with 7-day SLA. Build path: requirements definition and PDLC sub-process. Both converge to finalization.",
        "findings": [
            ("'Dumpster Fire #1': 2 people, 30+ contracts/month", "Manual review sustained 4 years unsustainably. CLM requested 5-6 years, no funding.", False),
            ("Contract negotiation: up to 1.5 years", "Security exhibits drive longest negotiations. No reportable format for deviations.", False),
            ("Sourcing makes quasi-legal decisions", "Legal bottleneck forces sourcing beyond expertise. Creates risk.", False),
            ("Only 2 legal partners for all contracts", "Capacity hasn't scaled with growth.", False),
            ("NDA timing inconsistency", "Legal and Security want NDA first. Others disagree. No enforcement.", False),
            ("Ariba for contract workspace", "Platform exists for NDAs, registration, contracts. Not fully used.", True),
        ],
        "quote": "That team desperately needs automation... The firm needs them to be automated.",
        "quote_source": "Risk Management Lead",
        "raci": "RACI: R/A Contracting | C Finance, Governance, Compliance | I Oversight",
        "recs": [
            ("Contract Review Automation", "AI-assisted for standard clauses, redlining, deviation detection.", "50%+ reduction in manual review"),
            ("Standardized Templates", "Pre-approved MSA/SOW by vendor tier and risk. Reduces negotiation scope.", "Days instead of months"),
            ("Parallel Contracting", "Begin negotiation concurrent with due diligence, not sequentially.", "Compresses 2-3 month phase"),
            ("Contract Deviation Register", "Structured tracking with risk classification, expiry dates, automated alerts.", "Full visibility into contract risk"),
        ],
        "roadmap": [
            ["Template library", "Parallel contracting design", "Deviation format"],
            ["Contract automation pilot", "AI redlining", "Deviation register"],
            ["Full parallel contracting", "Automated deviation alerts", "Contract analytics"],
            ["CLM system evaluation", "Historical remediation", "Vendor contract scoring"],
        ],
    },
    {
        "name": "Third-Party Risk Management",
        "subtitle": "Vendor lifecycle management, monitoring, and OCC 2023-17 compliance backbone.",
        "bpmn_image": "tprm-journey.png",
        "bpmn_title": "TPRM: End-to-End Vendor Lifecycle",
        "bpmn_desc": "Cross-cutting TPRM journey spanning all 5 phases: intake risk screening, due diligence with OneTrust, contract controls, ongoing monitoring, and vendor performance scoring.",
        "findings": [
            ("System landscape fragmented", "ServiceNow, Ariba, OneTrust, Oracle. Data in 4+ systems with manual PDF exports.", False),
            ("Process ownership vacuum", "'Somebody needs to be empowered to say I own this.' START: 'half a person.'", False),
            ("4 people contacted same vendor on breach", "No clear ownership for incident communication. Same questions repeated.", False),
            ("Tech owner identification weak", "Business Owner tracked. Vendor Owner by TPRM team. Tech Owner weakly tracked.", False),
            ("DD performance improving dramatically", "144 days (2019) to 75 days. 335 assessments/year. Output doubled.", True),
            ("Shift-left strategy endorsed", "'That's a fantastic idea that should be in the slides.' Risk-based self-service.", True),
        ],
        "quote": "If we want this to really click... I can't have the architect review group being a critical portion with only two people.",
        "quote_source": "Risk Management Lead",
        "raci": "RACI: R/A Governance | C Procurement, Contracting, Technical, Compliance | I Oversight",
        "recs": [
            ("Empowered Process Owner", "Dedicated strategic owner + 1 workflow manager + 2-3 PMs. Models Architecture Facilitator.", "Single accountability for E2E"),
            ("Mandatory Ownership Assignment", "Business, Technical, Support owners at onboarding. Annual validation. Feeds CMDB.", "Resolves #1 post-onboarding gap"),
            ("Integrated System of Record", "Connect ServiceNow, Ariba, OneTrust, Oracle. Eliminate manual PDF exports.", "Single source of truth across 4+ systems"),
            ("Incident Response Coordination", "Single vendor contact for incidents. Clear ownership matrix: who communicates, investigates, remediates.", "Prevents 4 people contacting same vendor"),
        ],
        "roadmap": [
            ["Define process owner role", "Ownership assignment form", "Incident communication matrix"],
            ["Process owner appointment", "System integration reqs", "Mandatory ownership"],
            ["OneTrust-Ariba-SN pilot", "Incident response playbook", "Annual validation design"],
            ["Full integrated SoR", "TPRM analytics dashboard", "Concentration risk monitoring"],
        ],
    },
]

import os

img_dir = "/Users/proth/repos/sla/docs/presentations/bpmn-images"

for domain in domains:
    # Section divider slide
    add_section_slide(domain["name"], domain["subtitle"])

    # BPMN context slide (woven into domain, v2-roadmap pattern)
    if domain.get("bpmn_image"):
        img_path = os.path.join(img_dir, domain["bpmn_image"])
        if os.path.exists(img_path):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(slide, KPMG_WHITE)
            add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5),
                         domain["bpmn_title"], font_size=22, bold=True, color=KPMG_BLUE,
                         font_name='Open Sans Condensed')
            add_text_box(slide, Inches(0.6), Inches(0.8), Inches(11.5), Inches(0.4),
                         domain["bpmn_desc"], font_size=9, color=RGBColor(0x66, 0x66, 0x66))
            pic = slide.shapes.add_picture(img_path, Inches(0.3), Inches(1.3), Inches(12.7))
            max_height = Inches(5.7)
            if pic.height > max_height:
                ratio = max_height / pic.height
                pic.width = int(pic.width * ratio)
                pic.height = max_height
                pic.left = int((prs.slide_width - pic.width) / 2)
            add_footer(slide)

    # Current state slide
    add_current_state_slide(
        f"{domain['name']}: Current State",
        domain["findings"],
        domain.get("quote"),
        domain.get("quote_source"),
        domain.get("raci", "")
    )

    # Recommendations + roadmap slide
    add_recommendations_slide(
        f"{domain['name']}: Recommendations and Roadmap",
        domain["recs"],
        domain["roadmap"]
    )

# ============================================================
# SUPPLEMENTARY ANALYTICS SECTION
# ============================================================

add_section_slide("Decision Framework and Analytics",
                  "DMN decision tables, RACI governance matrix, heatmap, and bottleneck analysis")

# --- Orchestrator Overview slide ---
orch_path = os.path.join(img_dir, "orchestrator-painpoints.png")
if os.path.exists(orch_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, KPMG_WHITE)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.5),
                 "BPMN: End-to-End Orchestrator with Pain Points", font_size=22, bold=True,
                 color=KPMG_BLUE, font_name='Open Sans Condensed')
    add_text_box(slide, Inches(0.6), Inches(0.8), Inches(11.5), Inches(0.4),
                 "Hierarchical 5-phase model: Request and Triage, Planning and Routing, Evaluation and DD, Contracting and Build, UAT and Go-Live. 3 decision gateways, vendor pool with message flows.",
                 font_size=9, color=RGBColor(0x66, 0x66, 0x66))
    pic = slide.shapes.add_picture(orch_path, Inches(0.3), Inches(1.3), Inches(12.7))
    max_height = Inches(5.7)
    if pic.height > max_height:
        ratio = max_height / pic.height
        pic.width = int(pic.width * ratio)
        pic.height = max_height
        pic.left = int((prs.slide_width - pic.width) / 2)
    add_footer(slide)

# --- E2E Process Flow slide ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, KPMG_WHITE)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
             "End-to-End Onboarding Process Flow", font_size=24, bold=True,
             color=KPMG_BLUE, font_name='Open Sans Condensed')
add_text_box(slide, Inches(0.6), Inches(1.0), Inches(11), Inches(0.4),
             "5-phase hierarchical BPMN 2.0 model with DMN-driven gateways, parallel evaluation, and SLA timers",
             font_size=10, color=RGBColor(0x66, 0x66, 0x66))

phases = [
    ("SP1: Request\nand Triage", "Review, Gather, Submit,\nTriage", "SLA: 2 Days"),
    ("SP2: Planning\nand Routing", "Prelim Analysis, Backlog,\nDMN-1, DMN-2", "SLA: 3 Days"),
    ("SP3: Evaluation\nand DD", "5 Parallel: Tech, Security,\nRisk, Finance, Vendor", "SLA: 5 Days"),
    ("SP4: Contracting\nand Build", "Buy: Negotiate, PoC\nBuild: Reqs, PDLC", "SLA: 7 Days"),
    ("SP5: UAT\nand Go-Live", "UAT, Approval, Onboard,\nNotify, Close", "SLA: 3 Days"),
]
for i, (name, tasks, sla) in enumerate(phases):
    left = Inches(0.5 + i * 2.5)
    top = Inches(1.6)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.2), Inches(1.4))
    box.fill.solid()
    box.fill.fore_color.rgb = KPMG_WHITE
    box.line.color.rgb = KPMG_BLUE
    box.line.width = Pt(2)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.05), Inches(2), Inches(0.45),
                 name, font_size=9, bold=True, color=KPMG_BLUE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.5), Inches(2), Inches(0.5),
                 tasks, font_size=7, color=RGBColor(0x66, 0x66, 0x66), alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), top + Inches(1.05), Inches(2), Inches(0.25),
                 sla, font_size=8, bold=True, color=KPMG_EMERALD, alignment=PP_ALIGN.CENTER)
    if i < 4:
        add_text_box(slide, left + Inches(2.2), top + Inches(0.5), Inches(0.3), Inches(0.3),
                     "\u25B6", font_size=14, color=KPMG_MEDIUM_BLUE, alignment=PP_ALIGN.CENTER)

cards_data = [
    ("4 DMN Tables", "Risk Tier, Pathway, Governance, SLA Breach"),
    ("3 Decision Gateways", "Build vs Buy?, Eval Approved?, Risk Acceptable?"),
    ("5 Parallel Streams", "Tech Arch, Security, Risk, Financial, Vendor"),
    ("SLA Timers", "Non-interrupting boundary timers with DMN-4 escalation"),
]
for i, (ct, cb) in enumerate(cards_data):
    add_card(slide, Inches(0.5 + i * 3.15), Inches(3.3), Inches(2.9), Inches(0.9), ct, cb,
             [KPMG_BLUE, KPMG_GOLD, KPMG_EMERALD, KPMG_LIGHT_BLUE][i])

# DMN flow description
add_text_box(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(2.5),
             "Decision Model Notation (DMN) Integration:\n\n"
             "\u2022 OB-DMN-1 (Risk Tier): 5 inputs (Data Sensitivity, Regulatory Exposure, Operational Criticality, "
             "AI Complexity, 3rd-Party Dependency) \u2192 4 tiers (Unacceptable, High, Limited, Minimal)\n"
             "\u2022 OB-DMN-2 (Pathway): 6 inputs including Existing Vendor Relationship \u2192 Fast-Track, Standard-Buy, "
             "Standard-Build, or Hybrid\n"
             "\u2022 OB-DMN-3 (Governance): Risk Tier + Pathway \u2192 Advisory Board (7d), Committee (5d), "
             "Fast Path (2d), or Auto-Approve (1d)\n"
             "\u2022 OB-DMN-4 (SLA Breach): Breached Phase + Days Beyond + Risk Tier \u2192 4-level escalation ladder",
             font_size=9, color=KPMG_DARK_GRAY)
add_footer(slide)

# --- RACI Matrix slide ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, KPMG_WHITE)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
             "RACI Matrix by Governance Topic", font_size=24, bold=True,
             color=KPMG_BLUE, font_name='Open Sans Condensed')

raci_headers = ["Topic", "Business", "Governance", "Finance", "Procurement", "Contracting",
                "Technical", "AI Review", "Compliance", "Oversight", "Automation", "Vendor"]
raci_data = [
    ["Intake", "R/A", "C", "C", "-", "-", "-", "-", "C", "-", "-", "-"],
    ["Prioritization", "C", "R/A", "C", "-", "-", "C", "-", "-", "-", "-", "-"],
    ["Funding", "-", "C", "R/A", "C", "-", "-", "-", "-", "I", "-", "-"],
    ["Sourcing", "-", "A", "-", "R", "-", "C", "-", "C", "I", "-", "R"],
    ["Cyber", "-", "C", "-", "-", "-", "R/A", "-", "C", "-", "-", "C"],
    ["EA", "-", "C", "-", "-", "-", "R/A", "-", "-", "-", "-", "C"],
    ["Compliance", "-", "C", "-", "-", "-", "-", "-", "R/A", "C", "-", "C"],
    ["AI Governance", "-", "I", "-", "-", "-", "C", "R/A", "C", "-", "-", "-"],
    ["Privacy", "C", "C", "-", "-", "-", "-", "-", "R/A", "-", "-", "-"],
    ["Comm. Counsel", "-", "C", "C", "-", "R/A", "-", "-", "C", "I", "-", "R"],
    ["TPRM", "-", "R/A", "-", "C", "C", "C", "-", "C", "C", "C", "C"],
]

rows_count = len(raci_data) + 1
cols_count = len(raci_headers)
table_shape = slide.shapes.add_table(rows_count, cols_count,
                                      Inches(0.3), Inches(1.1), Inches(12.7), Inches(4.5))
tbl = table_shape.table

for i, h in enumerate(raci_headers):
    cell = tbl.cell(0, i)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = KPMG_WHITE
        p.font.name = 'Open Sans'
    cell.fill.solid()
    cell.fill.fore_color.rgb = KPMG_GOLD if i == 3 else KPMG_BLUE

for r, row in enumerate(raci_data):
    for c, val in enumerate(row):
        cell = tbl.cell(r + 1, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.name = 'Open Sans'
            p.font.color.rgb = KPMG_DARK_GRAY
            if val in ("R/A", "R"):
                p.font.bold = True
                p.font.color.rgb = KPMG_BLUE
            if c == 0:
                p.font.bold = True

add_text_box(slide, Inches(0.3), Inches(5.8), Inches(12), Inches(0.3),
             "R = Responsible | A = Accountable | C = Consulted | I = Informed   |   Finance = new functional area (split from Governance)",
             font_size=8, color=RGBColor(0x66, 0x66, 0x66))
add_footer(slide)

# --- DMN-1 Risk Tier Classification slide ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, KPMG_WHITE)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.5),
             "OB-DMN-1: Risk Tier Classification", font_size=22, bold=True,
             color=KPMG_BLUE, font_name='Open Sans Condensed')

# Hit policy badge
badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.5), Inches(0.4), Inches(1.5), Inches(0.35))
badge.fill.solid()
badge.fill.fore_color.rgb = KPMG_BLUE
badge.line.fill.background()
add_text_box(slide, Inches(10.5), Inches(0.4), Inches(1.5), Inches(0.35),
             "UNIQUE | Phase 2", font_size=8, bold=True, color=KPMG_WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.6), Inches(0.9), Inches(11), Inches(0.4),
             "5 risk dimensions (1-10 scale each) determine risk tier. Unacceptable tiers terminate onboarding.",
             font_size=10, color=RGBColor(0x66, 0x66, 0x66))

dmn1_headers = ["Data Sensitivity", "Regulatory Exposure", "Operational Crit.", "AI Complexity", "3P Dependency", "Risk Tier", "Eligible?"]
dmn1_rows = [
    ["\u2265 9", "\u2265 9", "-", "-", "-", "Unacceptable", "No"],
    ["\u2265 9", "-", "\u2265 9", "-", "-", "Unacceptable", "No"],
    ["-", "\u2265 9", "\u2265 9", "-", "-", "Unacceptable", "No"],
    ["\u2265 9", "-", "-", "\u2265 8", "-", "Unacceptable", "No"],
    ["-", "-", "-", "\u2265 9", "\u2265 9", "Unacceptable", "No"],
    ["\u2265 7", "\u2265 7", "-", "-", "-", "High", "No"],
    ["\u2265 7", "-", "\u2265 7", "-", "-", "High", "No"],
    ["-", "\u2265 7", "-", "\u2265 6", "-", "High", "No"],
    ["-", "-", "\u2265 7", "-", "\u2265 7", "High", "No"],
    ["[4..7)", "[4..7)", "-", "-", "-", "Limited", "Yes"],
    ["[4..7)", "-", "[4..7)", "-", "-", "Limited", "Yes"],
    ["-", "[4..7)", "[4..7)", "-", "-", "Limited", "Yes"],
    ["-", "-", "[4..7)", "[3..6)", "[4..7)", "Limited", "Yes"],
    ["< 4", "< 4", "< 4", "< 3", "-", "Minimal", "Yes"],
    ["< 4", "< 4", "< 4", "-", "< 4", "Minimal", "Yes"],
]

t_shape = slide.shapes.add_table(len(dmn1_rows) + 1, len(dmn1_headers),
                                  Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.5))
t = t_shape.table
for i, h in enumerate(dmn1_headers):
    cell = t.cell(0, i)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(7)
        p.font.bold = True
        p.font.color.rgb = KPMG_WHITE
        p.font.name = 'Open Sans'
    cell.fill.solid()
    cell.fill.fore_color.rgb = KPMG_DARK_NAVY if i >= 5 else KPMG_MEDIUM_BLUE

tier_colors = {
    "Unacceptable": KPMG_ROSE, "High": KPMG_AMBER,
    "Limited": KPMG_MEDIUM_BLUE, "Minimal": KPMG_EMERALD
}
for r, row in enumerate(dmn1_rows):
    for c, val in enumerate(row):
        cell = t.cell(r + 1, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(7)
            p.font.name = 'Open Sans'
            p.font.color.rgb = KPMG_DARK_GRAY
            if c == 5 and val in tier_colors:
                p.font.bold = True
                p.font.color.rgb = tier_colors[val]
            if c == 6:
                p.font.color.rgb = KPMG_EMERALD if val == "Yes" else KPMG_ROSE
add_footer(slide)

# --- DMN-2 + DMN-3 side by side ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, KPMG_WHITE)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.5),
             "OB-DMN-2: Pathway Routing  |  OB-DMN-3: Governance Routing", font_size=22, bold=True,
             color=KPMG_BLUE, font_name='Open Sans Condensed')

# DMN-2 table (left half)
add_text_box(slide, Inches(0.5), Inches(1.0), Inches(6), Inches(0.3),
             "OB-DMN-2: Pathway Routing (UNIQUE, Phase 1)", font_size=10, bold=True, color=KPMG_DARK_NAVY)
dmn2_h = ["Reuse", "Build/Buy", "Capacity", "Strategic", "TtV", "Vendor?", "Pathway"]
dmn2_rows = [
    ["\u2265 8", "Buy", "-", "-", "\u2265 7", "true", "Fast-Track"],
    ["\u2265 7", "Buy", "-", "\u2265 7", "\u2265 7", "true", "Fast-Track"],
    ["-", "Buy", "-", "-", "-", "-", "Standard-Buy"],
    ["-", "Build", "\u2265 6", "-", "-", "-", "Standard-Build"],
    ["-", "Hybrid", "-", "-", "-", "-", "Hybrid"],
    ["-", "Build", "< 6", "-", "-", "-", "Hybrid"],
]
t2 = slide.shapes.add_table(len(dmn2_rows) + 1, len(dmn2_h), Inches(0.3), Inches(1.35), Inches(6.3), Inches(2.5)).table
for i, h in enumerate(dmn2_h):
    cell = t2.cell(0, i)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(7)
        p.font.bold = True
        p.font.color.rgb = KPMG_WHITE
        p.font.name = 'Open Sans'
    cell.fill.solid()
    cell.fill.fore_color.rgb = KPMG_DARK_NAVY if i == 6 else KPMG_MEDIUM_BLUE
pathway_colors = {"Fast-Track": KPMG_EMERALD, "Standard-Buy": KPMG_MEDIUM_BLUE,
                   "Standard-Build": KPMG_BLUE, "Hybrid": KPMG_GOLD}
for r, row in enumerate(dmn2_rows):
    for c, val in enumerate(row):
        cell = t2.cell(r + 1, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(7)
            p.font.name = 'Open Sans'
            p.font.color.rgb = KPMG_DARK_GRAY
            if c == 6 and val in pathway_colors:
                p.font.bold = True
                p.font.color.rgb = pathway_colors[val]

# DMN-3 table (right half)
add_text_box(slide, Inches(6.8), Inches(1.0), Inches(6), Inches(0.3),
             "OB-DMN-3: Governance Routing (UNIQUE, Phase 4)", font_size=10, bold=True, color=KPMG_DARK_NAVY)
dmn3_h = ["Risk Tier", "Pathway", "Review Body", "SLA"]
dmn3_rows = [
    ["Unacceptable", "-", "Rejected", "-"],
    ["High", "-", "Advisory Board", "7 days"],
    ["Limited", "Standard-*", "Committee", "5 days"],
    ["Limited", "Fast-Track", "Fast Path", "2 days"],
    ["Minimal", "Fast-Track", "Auto-Approve", "1 day"],
    ["Minimal", "Standard-*", "Fast Path", "2 days"],
]
t3 = slide.shapes.add_table(len(dmn3_rows) + 1, len(dmn3_h), Inches(6.8), Inches(1.35), Inches(6), Inches(2.5)).table
for i, h in enumerate(dmn3_h):
    cell = t3.cell(0, i)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(8)
        p.font.bold = True
        p.font.color.rgb = KPMG_WHITE
        p.font.name = 'Open Sans'
    cell.fill.solid()
    cell.fill.fore_color.rgb = KPMG_DARK_NAVY if i >= 2 else KPMG_MEDIUM_BLUE
for r, row in enumerate(dmn3_rows):
    for c, val in enumerate(row):
        cell = t3.cell(r + 1, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.name = 'Open Sans'
            p.font.color.rgb = KPMG_DARK_GRAY
            if c == 2:
                if val == "Rejected":
                    p.font.bold = True
                    p.font.color.rgb = KPMG_ROSE
                elif val in ("Fast Path", "Auto-Approve"):
                    p.font.bold = True
                    p.font.color.rgb = KPMG_EMERALD

# DMN-4 description below
add_text_box(slide, Inches(0.5), Inches(4.2), Inches(12), Inches(2.5),
             "OB-DMN-4: SLA Breach Escalation Routing (FIRST, Cross-Cutting)\n\n"
             "4-level escalation ladder triggered by SLA timer boundary events:\n"
             "  Level 1: Monitor \u2014 Dashboard alert to Process Owner (2+ days beyond SLA)\n"
             "  Level 2: Notify-Sponsor \u2014 Email to Business Sponsor (3-5 days beyond)\n"
             "  Level 3: Escalate-Governance \u2014 Priority reassignment by Governance Lead (5-10 days beyond)\n"
             "  Level 4: Board-Alert / Auto-Reject \u2014 Advisory Board escalation or process termination (10-15 days beyond, High risk)\n\n"
             "Inputs: Breached Phase, Days Beyond SLA, Risk Tier\n"
             "Outputs: Escalation Action, Notification Target, Auto-Action",
             font_size=9, color=KPMG_DARK_GRAY)
add_footer(slide)

# --- Governance Heatmap slide ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, KPMG_WHITE)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.5),
             "Governance Topic Coverage Across Onboarding Phases", font_size=22, bold=True,
             color=KPMG_BLUE, font_name='Open Sans Condensed')

hm_headers = ["Topic", "SP1: Request\n& Triage", "SP2: Planning\n& Routing", "SP3: Eval\n& DD",
              "SP4: Contract\n& Build", "SP5: UAT\n& Go-Live", "Vendor\nPool"]
hm_data = [
    ["Intake", 2, 0, 0, 0, 1, 0],
    ["Prioritization", 0, 2, 0, 0, 0, 0],
    ["Funding", 1, 1, 1, 1, 0, 0],
    ["Sourcing", 0, 0, 2, 2, 0, 2],
    ["Cyber", 0, 0, 2, 2, 1, 1],
    ["EA", 0, 0, 2, 2, 1, 1],
    ["Compliance", 1, 0, 2, 1, 2, 1],
    ["AI Governance", 0, 1, 2, 0, 0, 0],
    ["Privacy", 1, 1, 1, 0, 0, 0],
    ["Comm. Counsel", 0, 0, 0, 2, 0, 2],
    ["TPRM", 1, 0, 2, 2, 2, 2],
]

ht = slide.shapes.add_table(len(hm_data) + 1, len(hm_headers),
                             Inches(0.3), Inches(1.1), Inches(12.7), Inches(5)).table
for i, h in enumerate(hm_headers):
    cell = ht.cell(0, i)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(7)
        p.font.bold = True
        p.font.color.rgb = KPMG_WHITE
        p.font.name = 'Open Sans'
    cell.fill.solid()
    cell.fill.fore_color.rgb = KPMG_BLUE

for r, row in enumerate(hm_data):
    for c, val in enumerate(row):
        cell = ht.cell(r + 1, c)
        if c == 0:
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.bold = True
                p.font.name = 'Open Sans'
                p.font.color.rgb = KPMG_DARK_NAVY
            cell.fill.solid()
            cell.fill.fore_color.rgb = KPMG_LIGHT_GRAY
        else:
            cell.text = "\u2713" if val > 0 else ""
            cell.fill.solid()
            if val >= 2:
                cell.fill.fore_color.rgb = KPMG_BLUE
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = KPMG_WHITE
                    p.font.size = Pt(9)
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER
            elif val == 1:
                cell.fill.fore_color.rgb = RGBColor(0x93, 0xA5, 0xC0)
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = KPMG_WHITE
                    p.font.size = Pt(9)
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER
            else:
                cell.fill.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFC)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(8)
                    p.alignment = PP_ALIGN.CENTER

add_text_box(slide, Inches(0.3), Inches(6.3), Inches(12), Inches(0.3),
             "\u25A0 Primary (2+ tasks)    \u25A0 Secondary (1 task)    \u25A1 No coverage",
             font_size=8, color=RGBColor(0x66, 0x66, 0x66))
add_footer(slide)

# --- Bottleneck Analysis slide ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, KPMG_WHITE)
add_text_box(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.5),
             "Onboarding Bottleneck Analysis: Before vs. After", font_size=24, bold=True,
             color=KPMG_BLUE, font_name='Open Sans Condensed')

bottlenecks = [
    ("Vendor Evidence Collection", 15, 3, KPMG_ROSE),
    ("Committee Scheduling", 20, 5, KPMG_AMBER),
    ("Contract Review / Redlining", 18, 5, KPMG_MEDIUM_BLUE),
    ("Security Assessment Queue", 12, 4, KPMG_LIGHT_BLUE),
    ("AI Governance Backlog", 10, 3, KPMG_EMERALD),
]

for i, (label, before, after, color) in enumerate(bottlenecks):
    y = Inches(1.2 + i * 0.85)
    add_text_box(slide, Inches(0.5), y + Inches(0.05), Inches(3), Inches(0.35),
                 label, font_size=10, color=KPMG_DARK_GRAY, alignment=PP_ALIGN.RIGHT)
    max_width = Inches(8)
    scale = 20
    # Before bar
    bar_before = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(3.7), y, Emu(int(max_width.emu * before / scale)), Inches(0.3))
    bar_before.fill.solid()
    bar_before.fill.fore_color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    bar_before.line.fill.background()
    add_text_box(slide, Inches(3.7) + Emu(int(max_width.emu * before / scale)) + Inches(0.05), y,
                 Inches(0.5), Inches(0.3), f"{before}d", font_size=8, color=RGBColor(0x64, 0x74, 0x8B))
    # After bar
    bar_after = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(3.7), y + Inches(0.35), Emu(int(max_width.emu * after / scale)), Inches(0.3))
    bar_after.fill.solid()
    bar_after.fill.fore_color.rgb = color
    bar_after.line.fill.background()
    add_text_box(slide, Inches(3.7) + Emu(int(max_width.emu * after / scale)) + Inches(0.05), y + Inches(0.35),
                 Inches(0.5), Inches(0.3), f"{after}d", font_size=8, bold=True, color=color)

# Legend
add_text_box(slide, Inches(8), Inches(0.7), Inches(2), Inches(0.25),
             "\u25A0 Current State    \u25A0 Future State", font_size=8, color=RGBColor(0x66, 0x66, 0x66))

# Impact cards
impact = [
    ("75 Days \u2192 20 Days", "Total E2E cycle reduction"),
    ("18 \u2192 4 Committees", "Consolidated governance"),
    ("335 Assessments/Year", "Automated intake pipeline"),
    ("60+ AI Queue \u2192 0", "Dedicated AI governance lane"),
]
for i, (val, lbl) in enumerate(impact):
    add_card(slide, Inches(0.5 + i * 3.15), Inches(5.8), Inches(2.9), Inches(0.9), val, lbl,
             [KPMG_EMERALD, KPMG_GOLD, KPMG_LIGHT_BLUE, KPMG_BLUE][i])
add_footer(slide)

# Save
output_path = "/Users/proth/repos/sla/docs/presentations/Software-Onboarding-Transformation-KPMG.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
