#!/usr/bin/env python3
"""Generate industry-benchmarks-slide.pptx from data (deterministic output)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
BLUE = RGBColor(0x00, 0x33, 0x8D)
MEDIUM_BLUE = RGBColor(0x00, 0x5E, 0xB8)
LIGHT_BLUE = RGBColor(0x00, 0x91, 0xDA)
DARK_NAVY = RGBColor(0x00, 0x1D, 0x48)
GOLD = RGBColor(0xD4, 0xA8, 0x43)
EMERALD = RGBColor(0x00, 0xA8, 0x6B)
ROSE = RGBColor(0xC6, 0x28, 0x28)
AMBER = RGBColor(0xE6, 0x7E, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
LIGHT_ROSE_BG = RGBColor(0xFE, 0xF2, 0xF2)
LIGHT_GREEN_BG = RGBColor(0xEC, 0xFD, 0xF5)
BORDER_ROSE = RGBColor(0xFE, 0xCA, 0xCA)
BORDER_GREEN = RGBColor(0xA7, 0xF3, 0xD0)
TABLE_ALT = RGBColor(0xF9, 0xFA, 0xFB)
TABLE_BORDER = RGBColor(0xE5, 0xE7, 0xEB)

FONT = "Open Sans"
FONT_HEADING = "Open Sans"


def add_text_run(paragraph, text, size=9, bold=False, color=DARK_GRAY, font_name=FONT):
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return run


def add_hyperlink_run(paragraph, text, url, size=8, color=LIGHT_BLUE):
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = FONT
    run.font.underline = True
    hlink = run.hyperlink
    hlink.address = url
    return run


def add_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    # Adjust corner radius
    shape.adjustments[0] = 0.15
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    shape.adjustments[0] = 0.05
    return shape


def build_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # === Header gradient bar (3 colored rectangles to simulate gradient) ===
    bar_h = Inches(0.06)
    third = Inches(10) // 3
    for i, c in enumerate([BLUE, LIGHT_BLUE, GOLD]):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(third * i), Emu(0), Emu(third + 10), bar_h)
        s.fill.solid()
        s.fill.fore_color.rgb = c
        s.line.fill.background()

    # === Title ===
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    add_text_run(p, "Software Onboarding: Comparative Industry Benchmarks", size=20, bold=True, color=BLUE, font_name=FONT_HEADING)

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.65), Inches(9), Inches(0.3))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    add_text_run(p2, "End-to-end cycle times across financial services, TPRM, and enterprise procurement \u2014 current state vs. industry vs. optimized target", size=8, color=MEDIUM_GRAY)

    # ==============================
    # LEFT COLUMN
    # ==============================
    left_x = Inches(0.5)
    card_w = Inches(2.9)

    # --- Current State card ---
    cy = Inches(1.05)
    card_h = Inches(1.55)
    add_rounded_rect(slide, left_x, cy, card_w, card_h, LIGHT_ROSE_BG, BORDER_ROSE)

    # Current State label
    tb = slide.shapes.add_textbox(left_x + Inches(0.15), cy + Inches(0.08), card_w - Inches(0.3), Inches(0.22))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    add_text_run(p, "CURRENT STATE (CLIENT-MEASURED)", size=7, bold=True, color=ROSE, font_name=FONT_HEADING)

    # Big metric: 6-9 months
    tb = slide.shapes.add_textbox(left_x + Inches(0.15), cy + Inches(0.3), card_w - Inches(0.3), Inches(0.4))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    add_text_run(p, "6\u20139", size=28, bold=True, color=ROSE, font_name=FONT_HEADING)
    add_text_run(p, "  months end-to-end", size=10, color=MEDIUM_GRAY)

    # Stat pills as text
    stats_current = ["18 sequential committees", "5+ intake channels", "75d DD internal review",
                     "2\u20133mo contracting", "60+ AI queue items", "2 people / 30+ contracts/mo"]
    tb = slide.shapes.add_textbox(left_x + Inches(0.12), cy + Inches(0.78), card_w - Inches(0.2), Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = None
    for i, stat in enumerate(stats_current):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(1); p.space_after = Pt(1)
        parts = stat.split(" ", 1)
        add_text_run(p, parts[0] + " ", size=7, bold=True, color=DARK_GRAY)
        add_text_run(p, parts[1], size=7, color=DARK_GRAY)

    # --- Reduction banner ---
    ry = cy + card_h + Inches(0.08)
    banner_h = Inches(0.55)
    add_rounded_rect(slide, left_x, ry, card_w, banner_h, BLUE)
    tb = slide.shapes.add_textbox(left_x + Inches(0.15), ry + Inches(0.02), card_w - Inches(0.3), banner_h - Inches(0.04))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    add_text_run(p, "\u2193  ", size=16, color=WHITE)
    add_text_run(p, "68\u201375% Reduction", size=16, bold=True, color=GOLD, font_name=FONT_HEADING)
    p2 = tf.add_paragraph()
    add_text_run(p2, "Through parallel evaluation, workflow automation, front-loaded knowledge capture, and policy-as-code", size=7, color=WHITE)

    # --- Target State card ---
    ty = ry + banner_h + Inches(0.08)
    target_h = Inches(1.55)
    add_rounded_rect(slide, left_x, ty, card_w, target_h, LIGHT_GREEN_BG, BORDER_GREEN)

    tb = slide.shapes.add_textbox(left_x + Inches(0.15), ty + Inches(0.08), card_w - Inches(0.3), Inches(0.22))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    add_text_run(p, "OPTIMIZED TARGET STATE", size=7, bold=True, color=EMERALD, font_name=FONT_HEADING)

    tb = slide.shapes.add_textbox(left_x + Inches(0.15), ty + Inches(0.3), card_w - Inches(0.3), Inches(0.4))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    add_text_run(p, "29\u201345", size=28, bold=True, color=EMERALD, font_name=FONT_HEADING)
    add_text_run(p, "  days (standard risk)", size=10, color=MEDIUM_GRAY)

    stats_target = ["5 parallel streams", "60% automation", "20d phase SLA total",
                    "30%+ reuse/redirect", "2wk AI fast-track", "13.5d minimal risk"]
    tb = slide.shapes.add_textbox(left_x + Inches(0.12), ty + Inches(0.78), card_w - Inches(0.2), Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = None
    for i, stat in enumerate(stats_target):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(1); p.space_after = Pt(1)
        parts = stat.split(" ", 1)
        add_text_run(p, parts[0] + " ", size=7, bold=True, color=DARK_GRAY)
        add_text_run(p, parts[1], size=7, color=DARK_GRAY)

    # ==============================
    # RIGHT COLUMN — Benchmark Table
    # ==============================
    table_x = Inches(3.65)
    table_y = Inches(1.05)
    table_w = Inches(6.1)
    cols = 3
    rows = 6  # header + 5 data rows

    table_shape = slide.shapes.add_table(rows, cols, table_x, table_y, table_w, Inches(2.85))
    table = table_shape.table

    # Column widths
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.1)

    # Header row
    headers = ["Benchmark Metric", "Industry / Peer Data", "Comparison"]
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        add_text_run(p, h, size=7, bold=True, color=WHITE, font_name=FONT_HEADING)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data
    bench_data = [
        {
            "metric": "End-to-End Onboarding",
            "sources": [("APQC", "https://www.apqc.org/resources/benchmarking/open-standards-benchmarking/measures/average-cycle-time-days-set-supplier"),
                        ("VISO TRUST", "https://visotrust.com/resources/vendor-onboarding/")],
            "data": "90\u2013120 days industry avg (standard risk)\n60\u201390 days competitor benchmark\n4\u20136 weeks medium-complexity vendor (APQC)",
            "bars": "180-270d (ours) | 90-120d (industry) | 29-45d (target)"
        },
        {
            "metric": "Due Diligence Review",
            "sources": [("Whistic 2025", "https://www.whistic.com/resources/blog/tprm-impact-report-midyear-update")],
            "data": "75 days internal review (down from 144d)\n37.3 assessments/mo avg vendor load\n335/yr completed (8-person team)",
            "bars": "75d (current) | ~45d (industry est.) | 5d SLA (target)"
        },
        {
            "metric": "Vendor Assessment",
            "sources": [("Bitsight/Ponemon", "https://www.riskrecon.com/ponemon-report-data-risk-in-the-third-party-ecosystem-study"),
                        ("Whistic 2025", "https://www.whistic.com/resources/blog/tprm-impact-report-midyear-update")],
            "data": "28\u201329 days RAE (target: 14d)\n75% reduction with automation (Bitsight)\n94% of firms under-assess (Whistic 2025)",
            "bars": "29d (current) | 14d (internal tgt) | 7d (automated)"
        },
        {
            "metric": "Contract Cycle Time",
            "sources": [("Sirion CLM 2026", "https://www.sirion.ai/library/contract-insights/contract-management-roi-benchmarks/")],
            "data": "2\u20133 months (up to 1.5yr security)\n55% faster with AI-driven CLM\n39% faster overall contract lifecycle",
            "bars": "60-90d (current) | 30-45d (w/ CLM) | 7d SLA (target)"
        },
        {
            "metric": "TPRM Team Size",
            "sources": [("Whistic 2025", "https://www.whistic.com/resources/blog/tprm-impact-report-midyear-update"),
                        ("Venminder 2025", "https://www.venminder.com/blog/highlights-state-of-third-party-risk-management-2025-survey")],
            "data": "8.5 FTEs industry avg (up from 5.6)\n97% would do deeper assessments if resourced\n179 hrs/mo per vendor on assessments",
            "bars": "8 FTEs (DD team) | 8.5 avg (industry)"
        },
    ]

    bar_colors_map = {
        0: [ROSE, AMBER, EMERALD],
        1: [ROSE, AMBER, EMERALD],
        2: [ROSE, AMBER, EMERALD],
        3: [ROSE, AMBER, EMERALD],
        4: [LIGHT_BLUE, AMBER],
    }
    bar_widths_map = {
        0: [1.1, 0.65, 0.25],
        1: [0.8, 0.5, 0.18],
        2: [0.55, 0.28, 0.15],
        3: [0.95, 0.45, 0.2],
        4: [0.5, 0.53],
    }

    for ri, row_data in enumerate(bench_data):
        data_row = ri + 1

        # Alternate row fill
        if ri % 2 == 1:
            for ci in range(cols):
                table.cell(data_row, ci).fill.solid()
                table.cell(data_row, ci).fill.fore_color.rgb = TABLE_ALT

        # Col 0: Metric name + sources with hyperlinks
        cell = table.cell(data_row, 0)
        cell.vertical_anchor = MSO_ANCHOR.TOP
        p = cell.text_frame.paragraphs[0]
        add_text_run(p, row_data["metric"], size=8, bold=True, color=DARK_GRAY)
        p2 = cell.text_frame.add_paragraph()
        p2.space_before = Pt(2)
        for si, (name, url) in enumerate(row_data["sources"]):
            if si > 0:
                add_text_run(p2, "; ", size=6, color=MEDIUM_GRAY)
            add_hyperlink_run(p2, name, url, size=6, color=LIGHT_BLUE)

        # Col 1: Industry data
        cell = table.cell(data_row, 1)
        cell.vertical_anchor = MSO_ANCHOR.TOP
        lines = row_data["data"].split("\n")
        for li, line in enumerate(lines):
            p = cell.text_frame.paragraphs[0] if li == 0 else cell.text_frame.add_paragraph()
            p.space_before = Pt(1)
            # Bold the first part (before first space after number/%)
            # Simple: bold everything up to first letter-space pattern
            parts = line.split(" ", 2) if line[0].isdigit() else line.split(" ", 1)
            if len(parts) >= 2 and (parts[0][0].isdigit() or parts[0].endswith("%")):
                bold_part = parts[0] + " " + parts[1] if len(parts) > 2 else parts[0]
                rest = parts[2] if len(parts) > 2 else parts[1]
                add_text_run(p, bold_part + " ", size=7, bold=True, color=DARK_GRAY)
                add_text_run(p, rest, size=7, color=DARK_GRAY)
            else:
                add_text_run(p, line, size=7, color=DARK_GRAY)

        # Col 2: Comparison bars (text labels — bars drawn as shapes would overlap table)
        cell = table.cell(data_row, 2)
        cell.vertical_anchor = MSO_ANCHOR.TOP
        bar_labels = row_data["bars"].split(" | ")
        bar_colors = bar_colors_map[ri]
        for bi, label in enumerate(bar_labels):
            p = cell.text_frame.paragraphs[0] if bi == 0 else cell.text_frame.add_paragraph()
            p.space_before = Pt(1)
            # Unicode block char to simulate bar
            bar_len = int(bar_widths_map[ri][bi] * 12)
            add_text_run(p, "\u2588" * bar_len + " ", size=7, color=bar_colors[bi])
            add_text_run(p, label, size=6, bold=True, color=DARK_GRAY)

    # Set row heights
    for ri in range(rows):
        table.rows[ri].height = Inches(0.22) if ri == 0 else Inches(0.53)

    # ==============================
    # CYCLE TIME REDUCTION STACK
    # ==============================
    stack_y = Inches(4.05)
    stack_x = Inches(3.65)

    tb = slide.shapes.add_textbox(stack_x, stack_y, Inches(3), Inches(0.2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    add_text_run(p, "CYCLE TIME REDUCTION STACK", size=8, bold=True, color=BLUE, font_name=FONT_HEADING)

    stack_layers = [
        ("Baseline (Industry Avg)", "90\u2013120 days", Inches(3.1), ROSE),
        ("+ BPMN Workflow Automation", "65\u201385 days", Inches(2.25), AMBER),
        ("+ Front-Loaded Knowledge Capture", "50\u201365 days", Inches(1.65), RGBColor(0xF5, 0x9E, 0x0B)),
        ("+ Parallel Eval & DMN Routing", "35\u201350 days", Inches(1.15), RGBColor(0x34, 0xD3, 0x99)),
        ("+ Policy-as-Code / Compliance", "29\u201345 days", Inches(0.85), EMERALD),
    ]

    bar_x = stack_x + Inches(2.2)
    for i, (label, value, width, color) in enumerate(stack_layers):
        row_y = stack_y + Inches(0.25) + Inches(i * 0.22)

        # Label
        tb = slide.shapes.add_textbox(stack_x, row_y, Inches(2.15), Inches(0.2))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        add_text_run(p, label, size=7, color=MEDIUM_GRAY)

        # Bar
        bar = add_bar(slide, bar_x, row_y + Inches(0.02), width, Inches(0.16), color)

        # Value on bar
        tb = slide.shapes.add_textbox(bar_x + Inches(0.05), row_y + Inches(0.01), width - Inches(0.1), Inches(0.18))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        add_text_run(p, value, size=7, bold=True, color=WHITE)

    # ==============================
    # FOOTER
    # ==============================
    footer_y = Inches(5.15)
    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), footer_y, Inches(9), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    line.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.5), footer_y + Inches(0.05), Inches(7.5), Inches(0.3))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    add_text_run(p, "Sources: 14 stakeholder sessions (Feb\u2013Mar 2026) \u2022 PRD \u00a713 \u2022 ", size=6, color=MEDIUM_GRAY)
    add_hyperlink_run(p, "APQC Open Standards", "https://www.apqc.org/resources/benchmarking/open-standards-benchmarking/measures/average-cycle-time-days-set-supplier", size=6)
    add_text_run(p, " \u2022 ", size=6, color=MEDIUM_GRAY)
    add_hyperlink_run(p, "Whistic TPRM 2025", "https://www.whistic.com/resources/blog/tprm-impact-report-midyear-update", size=6)
    add_text_run(p, " / ", size=6, color=MEDIUM_GRAY)
    add_hyperlink_run(p, "Venminder 2025", "https://www.venminder.com/blog/highlights-state-of-third-party-risk-management-2025-survey", size=6)
    add_text_run(p, " \u2022 ", size=6, color=MEDIUM_GRAY)
    add_hyperlink_run(p, "Bitsight/Ponemon", "https://www.riskrecon.com/ponemon-report-data-risk-in-the-third-party-ecosystem-study", size=6)
    add_text_run(p, " \u2022 ", size=6, color=MEDIUM_GRAY)
    add_hyperlink_run(p, "Sirion CLM 2026", "https://www.sirion.ai/library/contract-insights/contract-management-roi-benchmarks/", size=6)
    add_text_run(p, " \u2022 ", size=6, color=MEDIUM_GRAY)
    add_hyperlink_run(p, "VISO TRUST", "https://visotrust.com/resources/vendor-onboarding/", size=6)

    # Right-side branding
    tb = slide.shapes.add_textbox(Inches(8.2), footer_y + Inches(0.05), Inches(1.8), Inches(0.2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    add_text_run(p, "Software Onboarding Transformation", size=6, color=MEDIUM_GRAY)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    build_slide(prs)

    out_path = os.path.join(os.path.dirname(__file__), "industry-benchmarks-slide.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
