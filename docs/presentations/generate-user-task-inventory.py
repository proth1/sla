#!/usr/bin/env python3
"""Generate User Task Inventory PPTX from v17 onboarding model data.

v2: Adds BPMN process model diagram slides before each phase table.
Requires PNG images pre-rendered from SVGs in bpmn-images/png/.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
import os
import subprocess
import sys

# Brand colors (from Software-Onboarding-Transformation.pptx)
KPMG_BLUE = RGBColor(0x00, 0x33, 0x8D)
KPMG_GOLD = RGBColor(0xD4, 0xA8, 0x43)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
LIGHT_ROW = RGBColor(0xE8, 0xED, 0xF4)  # Light blue-gray alternating
MED_BORDER = RGBColor(0xBB, 0xBB, 0xBB)

# Slide dimensions (widescreen 13.33" x 7.5")
SLIDE_WIDTH = Emu(12191695)
SLIDE_HEIGHT = Emu(6858000)

# Table layout
LEFT_MARGIN = Inches(0.5)
TOP_MARGIN = Inches(1.3)
TABLE_WIDTH = Inches(12.33)
COL_WIDTHS = [Inches(2.5), Inches(2.0), Inches(7.83)]

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
IMG_DIR = os.path.join(SCRIPT_DIR, "bpmn-images", "png")

PHASES = [
    {
        "title": "SP0: Mini RFP Pre-Screening",
        "subtitle": "Nested within SP1 \u2014 Self-service vendor pre-screening before formal intake",
        "images": ["v17-sp0-mini-rfp.png"],
        "tasks": [
            ("Understand the Need", "quarterback-lane",
             "Requester describes business problem, urgency, and what capability they need. Collects problem statement, desired outcomes, and timeline."),
            ("Understand Vendor Context", "quarterback-lane",
             "Requester provides vendor/solution context \u2014 whether they have a vendor in mind, market awareness, and selection criteria."),
            ("Classification Check", "governance-lane",
             "Governance validates the auto-classification from intake data. Reviews risk tier, solution category, and confirms or overrides the system recommendation."),
            ("Deal-Killer Results", "quarterback-lane",
             "Displays OB-DMN-7 deal-killer pre-screen results (green/red flags). Requester reviews disqualifying criteria and decides whether to proceed or abandon."),
            ("Question Preview", "quarterback-lane",
             "Previews the dynamically selected RFP questions (from OB-DMN-8) before sending to vendor. Requester can review and confirm the question set."),
            ("Send to Vendor(s)", "quarterback-lane",
             "Requester confirms vendor contact details and triggers RFP transmission. Captures delivery method and vendor acknowledgment."),
            ("Identify Additional Solutions Vendors", "quarterback-lane",
             "When deal-killer flags are raised, requester identifies alternative vendors or solutions to consider instead."),
        ],
    },
    {
        "title": "SP0: Vendor Response",
        "subtitle": "Vendor Pool \u2014 External vendor engagement tasks",
        "images": ["v17-sp0-mini-rfp.png"],
        "tasks": [
            ("Receive from Vendor(s)", "Vendor (external)",
             "Vendor submits their RFP response, pricing, and supporting documentation."),
            ("Initial Evaluation of Redlines", "contracting-lane",
             "Legal/Contracting performs initial review of vendor contract redlines and terms that deviate from standard."),
            ("Flagged for Early Legal Review", "Legal (unassigned)",
             "Escalation task when vendor response contains significant legal concerns requiring early counsel review."),
        ],
    },
    {
        "title": "SP1: Request & Triage",
        "subtitle": "Phase 1 \u2014 Intake, screening, and initial routing",
        "images": ["v17-sp1-request-triage.png", "v17-sp1-execute-nda.png"],
        "tasks": [
            ("Review Problem and Existing Solutions", "quarterback-lane",
             "Requester reviews whether existing enterprise software already addresses the need. System shows current solution inventory; requester assesses fit (FullMatch/Partial/NeedsGuidance/NoMatch)."),
            ("Quarterback Assistance", "quarterback-lane",
             "Facilitator provides coaching when requester is unsure about existing solutions. Advisory guidance to help navigate the portfolio \u2014 followed by a \u201cPursue it?\u201d decision gateway."),
            ("Pursue Exception with Quarterback", "quarterback-lane",
             "When a partial match exists, the Quarterback helps the requester build a case for why a new solution is still needed despite existing options."),
            ("Leverage Existing Software", "quarterback-lane",
             "Requester documents decision to use an existing solution. Captures which solution, implementation plan, and expected timeline. Process ends here (\u201cSoftware Leveraged\u201d)."),
            ("Gather Documentation", "quarterback-lane",
             "Requester assembles supporting documents \u2014 business case, requirements, vendor materials, budget justification, data classification details."),
            ("Submit Software Request", "quarterback-lane",
             "Formal submission of the software request. Collects business problem, urgency level, pathway preference, business criticality, estimated budget, data types involved, AI involvement flag, and submitter certification."),
            ("Initial Triage and Evaluation", "quarterback-lane",
             "Quarterback/Facilitator evaluates the request. Assesses strategic alignment, checks for duplicates, reviews concentration risk, resource availability, and preliminary risk indicators. Produces triage decision and priority assignment."),
            ("Classify Request", "governance-lane",
             "Governance classifies the request type (Defined Need / Forced Update / Speculative) with rationale. Determines which downstream pathway and evaluation depth is required."),
            ("Completeness Gate", "automation-lane",
             "Automated quality gate verifying all required intake fields are populated and documentation is sufficient before advancing to planning."),
        ],
    },
    {
        "title": "SP2: Planning & Risk Scoping",
        "subtitle": "Phase 2 \u2014 Business case, risk appetite, and pathway routing",
        "images": ["v17-sp2-planning-routing.png"],
        "tasks": [
            ("Preliminary Analysis", "governance-lane",
             "Business case and privacy screening. Collects market research summary, strategic alignment assessment, risk appetite alignment (Within/Borderline/Exceeds), preliminary cost/ROI estimates, budget source, DPIA requirement, capacity impact score, and whether full evaluation is needed."),
            ("Backlog Prioritization", "governance-lane",
             "Reviews request against the governance backlog. Assesses capacity, competing priorities, and assigns queue position. Determines if the request can proceed immediately or must wait."),
            ("Buy vs Build", "unassigned",
             "Decision-support task where the governance team reviews pathway routing results and confirms whether the request follows a Buy, Build, or Hybrid pathway."),
        ],
    },
    {
        "title": "SP3: Evaluation & Due Diligence",
        "subtitle": "Phase 3 \u2014 Parallel swarm evaluation across 9 domains",
        "images": ["v17-sp3-risk-sme-assessments.png", "v17-sp3-vendor-sourcing.png"],
        "tasks": [
            ("Enterprise Architecture", "technical-assessment",
             "[Parallel] Reviews technology stack compatibility, integration architecture, scalability, data flow patterns, and alignment with enterprise reference architecture. Runs concurrently with Security, AI Governance, TPRM, Compliance, Privacy, and Strategic Sourcing assessments."),
            ("Security Assessment", "technical-assessment",
             "[Parallel] Cybersecurity posture evaluation. Collects security tier (1-6), encryption status, key management, MFA, privileged access, pen test dates, MTTP metrics, incident response plan, SOC 2/ISO 27001/PCI/FedRAMP certifications. Produces security risk rating and approval. Runs concurrently with other SME assessments."),
            ("AI Governance", "ai-review",
             "[Parallel] AI model risk classification review. Evaluates AI involvement, bias/fairness controls, model transparency, explainability, and alignment with SR 11-7 and EU AI Act requirements. Runs concurrently with other SME assessments."),
            ("TPRM", "governance-lane",
             "[Parallel] Third-Party Risk Management assessment. Identifies risks, scores probability/impact, defines mitigation strategies, and produces overall vendor risk rating. Runs concurrently with other SME assessments."),
            ("Compliance", "compliance-lane",
             "[Parallel] Maps applicable regulations to the solution. Reviews regulatory gaps, compliance status, and determines remediation requirements. Runs concurrently with other SME assessments."),
            ("Legal - Privacy", "compliance-lane",
             "[Parallel] Data Protection Impact Assessment (DPIA) and privacy review. Evaluates data handling, consent mechanisms, cross-border transfers, and retention policies. Runs concurrently with other SME assessments."),
            ("Strategic Sourcing and Procurement", "ai-review",
             "[Parallel] Market analysis and vendor competitive landscape assessment. Reviews vendor options, market maturity, switching costs, and procurement strategy. Runs concurrently with other SME assessments."),
            ("Vendor Due Diligence", "contracting-lane",
             "Comprehensive vendor assessment \u2014 financial health, references, customer satisfaction, operational resilience, sub-contractor dependencies, and contractual track record. Runs sequentially after parallel assessments complete."),
            ("Evaluate Vendor Response", "contracting-lane",
             "Scores and analyzes vendor RFP/proposal response. Gap analysis against requirements, pricing evaluation, and comparative vendor ranking. Runs sequentially after parallel assessments complete."),
        ],
    },
    {
        "title": "SP4: Governance Review & Contracting \u2014 Buy Pathway",
        "subtitle": "Phase 4 (Buy) \u2014 Requirements refinement, POC, negotiation, and contract execution",
        "images": ["v17-sp4-contracting.png"],
        "tasks": [
            ("Refine Requirements", "quarterback-lane",
             "Detailed requirements refinement based on due diligence findings. Updates acceptance criteria, integration requirements, and performance expectations."),
            ("Perform Proof of Concept", "technical-assessment",
             "POC execution and evaluation. Documents scope, test scenarios, results, and produces a go/no-go recommendation for full procurement."),
            ("Technology and Risk Evaluation", "technical-assessment",
             "Technical risk assessment \u2014 evaluates tech debt, scalability risks, complexity, vendor lock-in, and architectural fitness."),
            ("Negotiate Contract", "contracting-lane",
             "Contract negotiation. Captures contract type (MSA/SOW/SaaS/License), term, total value, payment terms, SLA commitments, audit rights, data protection clauses, termination rights, liability caps, and regulatory provisions."),
            ("Finalize Contract", "contracting-lane",
             "Final contract review and signature preparation. Legal sign-off on all terms, execution readiness confirmation."),
            ("Execute Contract", "contracting-lane",
             "Formal contract execution \u2014 captures signatures, effective dates, and archives executed agreement."),
            ("Contract Deviation Check", "contracting-lane",
             "Documents and reviews any deviations from standard contract terms or organizational policy. Requires risk acceptance documentation."),
            ("Compliance Review", "compliance-lane",
             "Validates that governance controls and monitoring requirements are defined and implementable for the procured solution."),
            ("Coding Matrix Correction", "quarterback-lane",
             "Corrects financial coding and categorization for the procurement (GL codes, cost centers, asset classification)."),
        ],
    },
    {
        "title": "SP4: Build Pathway & PDLC",
        "subtitle": "Phase 4 (Build) \u2014 Custom development lifecycle",
        "images": ["v17-pdlc.png"],
        "tasks": [
            ("Define Build Requirements", "quarterback-lane",
             "Defines build-specific requirements \u2014 architecture decisions, technology stack, timeline, resource needs, and security requirements for custom development."),
            ("Architecture Review", "technical-assessment",
             "Design review of proposed architecture. Evaluates patterns, security posture, performance characteristics, and alignment with enterprise standards."),
            ("Development and Build", "technical-assessment",
             "Tracks development progress. Collects sprint number, features completed, code review status, static analysis results, unit test coverage, SAST/secrets/dependency scan status."),
            ("Testing and Quality Assurance", "technical-assessment",
             "QA gate \u2014 test coverage, defect tracking, regression testing, and quality metric thresholds."),
            ("Integration and Deployment", "technical-assessment",
             "System integration testing, API validation, end-to-end testing, and deployment readiness assessment."),
        ],
    },
    {
        "title": "SP5: UAT & Go-Live",
        "subtitle": "Phase 5 \u2014 Acceptance testing, audit, onboarding, and closure",
        "images": ["v17-sp5-uat-golive.png"],
        "tasks": [
            ("Verify Conditions", "governance-lane",
             "Reviews and verifies that all governance conditions from prior approvals (conditional approvals, remediation items) have been satisfactorily addressed."),
            ("Perform Pilot / UAT", "quarterback-lane",
             "User acceptance testing execution. Collects test scenario counts (total/passed/failed), critical defects, pilot user count, pilot duration, user satisfaction score, training completion rate, and UAT sign-off from business owner."),
            ("3rd Line Audit Review", "oversight-lane",
             "Independent audit review by Internal Audit (3rd Line of Defense). Validates that governance controls are operating effectively and regulatory requirements are met. Approval gate."),
            ("Onboard Software", "automation-lane",
             "Technical onboarding \u2014 installation, configuration, data migration, access provisioning, and production readiness verification."),
            ("Assign Ownership", "governance-lane",
             "Assigns long-term operational ownership \u2014 names the operational owner, support team, escalation contacts, and ongoing monitoring responsibilities."),
            ("Close Request", "governance-lane",
             "Formal closure of the governance request. Captures lessons learned, archives documentation, and triggers operational handoff."),
        ],
    },
]

LANE_LABELS = {
    "quarterback-lane": "Quarterback",
    "governance-lane": "Governance",
    "contracting-lane": "Contracting",
    "technical-assessment": "Technical Assessment",
    "ai-review": "AI Review",
    "compliance-lane": "Compliance",
    "oversight-lane": "Oversight",
    "automation-lane": "Automation",
    "vendor-response": "Vendor Response",
    "Vendor (external)": "Vendor (External)",
    "Legal (unassigned)": "Legal (Unassigned)",
    "unassigned": "Governance Team",
}


def set_cell_border(cell, color=MED_BORDER):
    """Set thin borders on a cell."""
    from pptx.oxml.ns import qn
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge in ("lnL", "lnR", "lnT", "lnB"):
        ln = tcPr.find(qn(f"a:{edge}"))
        if ln is not None:
            tcPr.remove(ln)
        from lxml import etree
        ln = etree.SubElement(tcPr, qn(f"a:{edge}"), w="6350", cap="flat", cmpd="sng")
        sf = etree.SubElement(ln, qn("a:solidFill"))
        srgb = etree.SubElement(sf, qn("a:srgbClr"), val=f"{color}")


def format_cell(cell, text, font_size=9, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT):
    """Format a table cell with consistent styling."""
    cell.text = ""
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    # Tight margins
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)


def add_diagram_slide(prs, title, image_path):
    """Add a slide with a BPMN process model diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Phase title bar
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.6)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = KPMG_BLUE
    bar.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.08), Inches(11), Inches(0.45))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

    # Load image to get aspect ratio
    img = Image.open(image_path)
    img_w, img_h = img.size
    aspect = img_w / img_h

    # Available area: full width minus margins, below title bar
    avail_w = Inches(12.33)
    avail_h = Inches(6.5)  # 7.5 - 0.6 title - 0.4 padding
    avail_aspect = avail_w / avail_h

    if aspect > avail_aspect:
        # Image is wider — fit to width
        pic_w = avail_w
        pic_h = int(avail_w / aspect)
    else:
        # Image is taller — fit to height
        pic_h = avail_h
        pic_w = int(avail_h * aspect)

    # Center horizontally, position below title
    pic_x = (SLIDE_WIDTH - pic_w) // 2
    pic_y = Inches(0.7)

    slide.shapes.add_picture(image_path, pic_x, pic_y, pic_w, pic_h)


def add_title_slide(prs):
    """Add title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Blue bar at top
    from pptx.util import Emu as E
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = KPMG_BLUE
    bar.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(10), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "User Task Inventory"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = KPMG_BLUE
    run.font.name = "Calibri"

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(10), Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = "Software Onboarding Process v17"
    run2.font.size = Pt(20)
    run2.font.color.rgb = DARK_TEXT
    run2.font.name = "Calibri"

    # Summary line
    total = sum(len(ph["tasks"]) for ph in PHASES)
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(10), Inches(0.6))
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = f"{total} user tasks across {len(PHASES)} phases"
    run3.font.size = Pt(14)
    run3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    run3.font.name = "Calibri"

    # Gold accent bar
    bar2 = slide.shapes.add_shape(
        1, Inches(1), Inches(5.2), Inches(3), Inches(0.04)
    )
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = KPMG_GOLD
    bar2.line.fill.background()


def add_phase_slide(prs, phase):
    """Add a slide with a phase table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Phase title
    txBox = slide.shapes.add_textbox(LEFT_MARGIN, Inches(0.3), Inches(11), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = phase["title"]
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = KPMG_BLUE
    run.font.name = "Calibri"

    # Subtitle
    txBox2 = slide.shapes.add_textbox(LEFT_MARGIN, Inches(0.8), Inches(11), Inches(0.4))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = phase["subtitle"]
    run2.font.size = Pt(11)
    run2.font.italic = True
    run2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    run2.font.name = "Calibri"

    # Table
    rows = len(phase["tasks"]) + 1  # +1 for header
    table_shape = slide.shapes.add_table(rows, 3, LEFT_MARGIN, TOP_MARGIN, TABLE_WIDTH, Inches(0.4 * rows))
    table = table_shape.table

    # Column widths
    for i, w in enumerate(COL_WIDTHS):
        table.columns[i].width = w

    # Header row
    headers = ["Process Step", "Responsible Party", "Description"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = KPMG_BLUE
        format_cell(cell, h, font_size=10, bold=True, color=WHITE)
        set_cell_border(cell, color=KPMG_BLUE)

    # Data rows
    for r, (step, lane, desc) in enumerate(phase["tasks"], start=1):
        # Alternating row fill
        for c in range(3):
            cell = table.cell(r, c)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_ROW
            else:
                cell.fill.background()
            set_cell_border(cell)

        format_cell(table.cell(r, 0), step, font_size=9, bold=True)
        format_cell(table.cell(r, 1), LANE_LABELS.get(lane, lane), font_size=9)
        format_cell(table.cell(r, 2), desc, font_size=8)

    # Footer with task count
    txFoot = slide.shapes.add_textbox(LEFT_MARGIN, Inches(7.0), Inches(4), Inches(0.3))
    tf_foot = txFoot.text_frame
    pf = tf_foot.paragraphs[0]
    rf = pf.add_run()
    rf.text = f"{len(phase['tasks'])} tasks"
    rf.font.size = Pt(9)
    rf.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    rf.font.name = "Calibri"


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    add_title_slide(prs)

    # Orchestrator overview diagram (slide 2)
    orch_path = os.path.join(IMG_DIR, "v17-orchestrator.png")
    if os.path.exists(orch_path):
        add_diagram_slide(prs, "End-to-End Process Overview", orch_path)

    slide_count = 2  # title + orchestrator
    for phase in PHASES:
        # Diagram slide(s) for this phase
        seen = set()
        for img_name in phase.get("images", []):
            if img_name in seen:
                continue
            seen.add(img_name)
            img_path = os.path.join(IMG_DIR, img_name)
            if os.path.exists(img_path):
                add_diagram_slide(prs, f"Process Model: {phase['title']}", img_path)
                slide_count += 1
            else:
                print(f"  WARNING: Missing image {img_path}")

        # Task table slide
        add_phase_slide(prs, phase)
        slide_count += 1

    out_path = os.path.join(SCRIPT_DIR, "user-task-inventory.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")
    total = sum(len(ph["tasks"]) for ph in PHASES)
    print(f"  {slide_count} slides ({len(PHASES)} tables + diagram slides + title)")
    print(f"  {total} user tasks total")


if __name__ == "__main__":
    main()
